#!/usr/bin/env python3
"""
axi_bot.py
──────────
Axi — AIMS Orchestrator & Quality Monitor
Telegram Bot: оркестрация, мониторинг Task Registry, генерация документов.

Роль (docs/project-owner-answers.md §9):
  - Оркестратор: конфиги, связи между ботами, расписания, ключи
  - Монитор качества: периодически сканирует зависшие задачи, уведомляет владельца
  - DB задачи → роутинг через Omi

Переменные окружения:
    AXI_BOT_TOKEN / AXI_TELEGRAM_TOKEN — токен @BotFather
    AXI_ALLOWED_USER_IDS              — через запятую user id (личка: id чата = user id;
                                        группа: либо id группы, либо id отправителя должен быть в списке)
    AXI_OWNER_CHAT_IDS                — чаты владельца (для уведомлений о застрявших задачах)
    AXI_NAME                          — имя бота в ответах (default: Axi)
    ANTHROPIC_API_KEY                 — Claude API key для локального fallback
    AXI_ANTHROPIC_MODEL               — модель (default: claude-sonnet-4-6)
    AXI_AIMS_KNOWLEDGE_MODE           — off (default) | on_topic | always
    AXI_RESULTS_DIR                   — куда писать сгенерированные docx (default: /data/result)
    TASK_REGISTRY_URL                 — http://localhost:8765
    TASK_REGISTRY_STUCK_MINUTES       — порог зависших задач в минутах (default: 15)
    TASK_REGISTRY_MONITOR_INTERVAL    — интервал мониторинга в секундах (default: 900)
    QWEN_PC_ASSIST_STACK              — 1: прогрев 70B + PC Qwen при каждом сообщении
    AXI_LLM_LOCAL_FIRST               — 1: сначала Ollama, fallback → Anthropic. 0 (default): облако.
"""

from __future__ import annotations

import asyncio
import hashlib
import json
import logging
import os
import re
import sys
import tempfile
import time
from datetime import datetime, timezone
from pathlib import Path
from threading import Lock

from telegram import BotCommand, InlineKeyboardButton, InlineKeyboardMarkup, Update
from telegram.constants import ChatAction
from telegram.ext import (
    Application,
    CallbackQueryHandler,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    filters,
)

import sys as _sys
_CORE_PATH = str(Path(__file__).resolve().parent)
if _CORE_PATH not in _sys.path:
    _sys.path.insert(0, _CORE_PATH)
from core.metrics import record_llm_call  # noqa: E402

# ── Logging ────────────────────────────────────────────────────────────────────

logging.basicConfig(
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    level=logging.INFO,
)
log = logging.getLogger("axi")

# ── Domain module imports (Phase C Task 2 refactor) ───────────────────────────
# Each domain module is self-contained (reads env at import time).
# We re-export everything here so call sites in this file remain unchanged.

try:
    from axi_llm_client import (  # noqa: F401, F403
        AXI_SYSTEM_PROMPT as _llm_AXI_SYSTEM_PROMPT,
        AXI_SKILL_CONTEXT as _llm_AXI_SKILL_CONTEXT,
        _load_axi_skill_context,
        _axi_system_prompt,
        _get_llm_semaphore,
        _axi_llm_local_first_enabled,
        _animate_progress,
        _ft_log_example,
        _gemini_reply,
        _anthropic_reply,
        _anthropic_classify_intent,
        _local_ollama_reply_sync,
        _cloud_llm_reply,
        _llm_reply,
        _llm_reply_inner,
        _should_web_search,
        _PULSE_FRAMES,
    )
    # Override module-level globals with values from llm_client module
    # so call sites that reference AXI_SYSTEM_PROMPT / AXI_SKILL_CONTEXT
    # get consistent values.
    AXI_SYSTEM_PROMPT = _llm_AXI_SYSTEM_PROMPT
    AXI_SKILL_CONTEXT = _llm_AXI_SKILL_CONTEXT
    log.debug("axi_llm_client: imported successfully")
except ImportError as _e:
    log.warning("axi_llm_client import failed (using inline fallback): %s", _e)

try:
    from axi_omniroute_session import (  # noqa: F401
        _omniroute_base_urls,
        _omniroute_get_active_sessions_sync,
        _omniroute_close_session_sync,
        _omniroute_debug_session_tables_sync,
        _omniroute_preflight_free_one_slot_sync,
        _track_omniroute_session,
        _cleanup_omniroute_sessions_from_state,
        _extract_omniroute_content_from_body,
        _call_doctuning_openai_omniroute_sync,
        _extract_doc_context_sync,
        _DOCFILL_CONTEXT_SYSTEM,
    )
    log.debug("axi_omniroute_session: imported successfully")
except ImportError as _e:
    log.warning("axi_omniroute_session import failed (using inline fallback): %s", _e)

try:
    from axi_standards_check import (  # noqa: F401
        _build_doc_analysis_system as _sc_build_doc_analysis_system,
        _wants_standards_docx_result,
        _build_applied_corrections_section,
        _sources_note,
        _wants_standards_check,
        AIMS_ENABLE_CONTEXTUAL_STANDARD_DISCOVERY,
        _wants_contextual_discovery,
        _wants_registry_list,
        _wants_db_strategy,
        _extract_search_keywords,
        _register_corrected_standards_knowledge,
        _resolve_doc_path,
        _fetch_db_docs,
    )
    log.debug("axi_standards_check: imported successfully")
except ImportError as _e:
    log.warning("axi_standards_check import failed (using inline fallback): %s", _e)

try:
    from axi_doctuning import (  # noqa: F401
        _build_validate_system,
        _detect_template_reference_roles,
        _doctuning_ws,
        _save_doctuning_memo,
        _find_existing_training_pair,
        _save_doctuning_failure_case,
        _queue_doctuning_repair_case,
        _request_repairman_investigation_for_doctuning_failure,
        _read_document_text,
        _validate_and_extract_context_sync,
        _nim_validate_document_sync,
        _local_generate_fill_sync,
        _check_ollama_alive_sync,
        _check_vram_free_gb_sync,
        _repairman_fix_local_model_sync,
        _repairman_final_report_sync,
        _local_fill_with_repair,
        _xlsx_safe_text,
        _create_doctuning_xlsx_candidate,
        _xlsx_write_validation_failure,
        _docfill_save_dpo_pair,
        _docfill_save_training_pair,
        _doc_creator_120b_save,
        _doc_creator_120b_save_failed,
        _docfill_save_to_omi,
        AIMS_DOCUMENT_CREATOR_MODEL as _dt_AIMS_DOCUMENT_CREATOR_MODEL,
        AIMS_DOCUMENT_TEACHER_MODEL as _dt_AIMS_DOCUMENT_TEACHER_MODEL,
        AIMS_DOCUMENT_AUDIT_MODEL as _dt_AIMS_DOCUMENT_AUDIT_MODEL,
        AIMS_DOC_CREATOR_TRAINING_DIR as _dt_AIMS_DOC_CREATOR_TRAINING_DIR,
    )
    log.debug("axi_doctuning: imported successfully")
except ImportError as _e:
    log.warning("axi_doctuning import failed (using inline fallback): %s", _e)

try:
    from axi_job_pipeline import _queue_to_omi_batch_inbox as _jb_queue_to_omi_batch_inbox  # noqa: F401
    log.debug("axi_job_pipeline: imported successfully")
except ImportError as _e:
    log.warning("axi_job_pipeline import failed (using inline fallback): %s", _e)

# ─────────────────────────────────────────────────────────────────────────────

# ── Config ─────────────────────────────────────────────────────────────────────

AXI_BOT_TOKEN = (
    os.environ.get("AXI_BOT_TOKEN", "")
    or os.environ.get("AXI_TELEGRAM_TOKEN", "")
).strip()

AXI_NAME = os.environ.get("AXI_NAME", "Axi").strip()

def _parse_ids(env_key: str) -> frozenset[int]:
    raw = os.environ.get(env_key, "").strip()
    out: set[int] = set()
    for part in raw.split(","):
        part = part.strip()
        try:
            if part:
                out.add(int(part))
        except ValueError:
            pass
    return frozenset(out)

ALLOWED_CHATS   = _parse_ids("AXI_ALLOWED_USER_IDS")
OWNER_CHATS     = _parse_ids("AXI_OWNER_CHAT_IDS") or _parse_ids("AXI_ALLOWED_USER_IDS")

AXI_RESULTS_DIR = Path(os.environ.get("AXI_RESULTS_DIR", "/data/result"))
AXI_AIMS_KNOWLEDGE_MODE = os.environ.get("AXI_AIMS_KNOWLEDGE_MODE", "off").strip().lower()
AXI_WEB_SEARCH_ENABLED = False  # Disabled: NVIDIA NIM does not support Google Search grounding
AXI_CHAT_SHOW_SOURCES = os.environ.get("AXI_CHAT_SHOW_SOURCES", "0").strip().lower() in ("1", "true", "yes", "on")
AXI_CHAT_SHOW_TASK_REGISTRY_WARNINGS = os.environ.get(
    "AXI_CHAT_SHOW_TASK_REGISTRY_WARNINGS", "0"
).strip().lower() in ("1", "true", "yes", "on")

TASK_REGISTRY_STUCK_MINUTES  = float(os.environ.get("TASK_REGISTRY_STUCK_MINUTES", "15"))
TASK_REGISTRY_MONITOR_INTERVAL = int(os.environ.get("TASK_REGISTRY_MONITOR_INTERVAL", "900"))
TASK_REGISTRY_AUTO_CLEANUP_MINUTES = float(os.environ.get("TASK_REGISTRY_AUTO_CLEANUP_MINUTES", "60"))

QWEN_PC_ASSIST_STACK = os.environ.get("QWEN_PC_ASSIST_STACK", "0").strip().lower() in ("1", "true", "yes", "on")

# Pending intent clarification: {chat_id: {prompt, files, ts}}
_pending_intent: dict[int, dict] = {}
# Self-learning log for clarification confirmations → future fine-tuning data
_SELF_LEARN_LOG = AXI_RESULTS_DIR.parent / "axi_intent_learn.jsonl"

AXI_LLM_LOCAL_FIRST = os.environ.get("AXI_LLM_LOCAL_FIRST", "1").strip().lower() in ("1", "true", "yes", "on")
AXI_FT_LOG_ENABLED = os.environ.get("AXI_FT_LOG_ENABLED", "0").strip().lower() in ("1", "true", "yes", "on")
AXI_FT_LOG_DIR = Path(os.environ.get("AXI_FT_LOG_DIR", "/data/axi_ft_log"))
# Limit concurrent heavy LLM calls to prevent GPU saturation (96%+ util / thermal shutdown)
AXI_LLM_MAX_CONCURRENT = max(1, int(os.environ.get("AXI_LLM_MAX_CONCURRENT", "1")))
AXI_ANALYZE_WAIT_SEC = int(os.environ.get("AXI_ANALYZE_WAIT_SEC", "1800"))
AXI_ANALYZE_BATCH_WAIT_SEC = int(os.environ.get("AXI_ANALYZE_BATCH_WAIT_SEC", "20"))
AXI_DIALOG_LOG_MAX = max(2, int(os.environ.get("AXI_DIALOG_LOG_MAX", "10")))
AXI_FORCE_REPLY_LANG = os.environ.get("AXI_FORCE_REPLY_LANG", "en").strip().lower()
# In groups, require explicit @Axi / "Axi ..." mention only when enabled.
AXI_GROUP_REQUIRE_MENTION = os.environ.get("AXI_GROUP_REQUIRE_MENTION", "0").strip().lower() in (
    "1", "true", "yes", "on"
)
AXI_GROUP_ALLOW_ALL_MEMBERS = os.environ.get("AXI_GROUP_ALLOW_ALL_MEMBERS", "1").strip().lower() in (
    "1", "true", "yes", "on"
)

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "").strip()
ANTHROPIC_MODEL   = os.environ.get("AXI_ANTHROPIC_MODEL", "claude-sonnet-4-6").strip()

# Document-generation model roles.
# Creator: orchestrator-layer draft model (resolved via registry/slot abstraction).
# Teacher / Audit: external OmniRoute models — standards discovery, quality judging, audit.
# NOTE: nemotron-120B should NOT be assumed active per project spec.
# Use slot-based resolution (slot120/slot32) or dynamic registry lookup instead.
AIMS_DOCUMENT_CREATOR_MODEL: str = (
    os.environ.get("AIMS_DOCUMENT_CREATOR_MODEL")
    or os.environ.get("AIMS_LOCAL_DRAFT_MODEL")
    or "axi_omi_sphere"  # Resolved via registry abstraction
)
AIMS_DOCUMENT_TEACHER_MODEL: str = (
    os.environ.get("AIMS_DOCUMENT_TEACHER_MODEL")
    or os.environ.get("OMNIROUTE_TEACHER_MODEL")
    or "doc-training-standards-best"
)
AIMS_DOCUMENT_AUDIT_MODEL: str = (
    os.environ.get("AIMS_DOCUMENT_AUDIT_MODEL")
    or os.environ.get("OMNIROUTE_AUDIT_MODEL")
    or "doc-training-pair-audit-combo"
)

# ChatLLM conversational layer — slot14 lightweight chat model.
# Roles: user-facing chat, clarification, summaries, lightweight workflow guidance,
# fallback when heavier models are busy.
# Resolved via registry/slot abstraction (do not hardcode model names).
AIMS_CHATLLM_MODEL: str = (
    os.environ.get("AIMS_CHATLLM_MODEL")
    or os.environ.get("AIMS_CHAT_MODEL")
    or "omi-ft-14b-v16"  # Legacy: will be replaced by slot14 registry resolution
)
AIMS_CHATLLM_PC_URL: str = (
    os.environ.get("AIMS_CHATLLM_PC_URL")
    or os.environ.get("PC_OLLAMA_URL")
    or f"http://{os.environ.get('ANDREI_HOST', '10.77.77.2')}:11434"
)
AIMS_CHATLLM_DGX_URL: str = (
    os.environ.get("AIMS_CHATLLM_DGX_URL")
    or os.environ.get("OLLAMA_LOCAL_URL", "http://localhost:11434")
)

# Workflow / repair / coding / primary action classifier — 32B on DGX.
# Never route document creation here.
AIMS_WORKFLOW_MODEL: str = (
    os.environ.get("AIMS_WORKFLOW_MODEL")
    or os.environ.get("AIMS_REPAIR_MODEL")
    or "axi_omi_sphere"
)

# 120B training data pool root — organized by pair type.
AIMS_DOC_CREATOR_TRAINING_DIR: Path = Path(
    os.environ.get("AIMS_DOC_CREATOR_TRAINING_DIR", "aims_workspace/training/document_creator_120b")
)

_PENDING_ANALYZE: dict[int, dict[str, object]] = {}
_PENDING_VIDEO: dict[int, dict[str, object]] = {}
_PENDING_DOCFILL: dict[int, dict[str, object]] = {}  # step: await_blank | await_example | validating | await_approval
_DOCFILL_VRAM_QUEUE: dict[int, str] = {}            # chat_id → blank_text, waiting for VRAM headroom
_PENDING_DOCREVIEW: dict[int, dict[str, object]] = {}   # step: await_doc | ready
_PENDING_DOC_INTAKE: dict[int, dict[str, object]] = {}  # step: await_more | await_task | await_mode
_PENDING_AI_LABEL_REMOVAL: dict[int, dict[str, object]] = {}  # prompt + ts for document-followup routing
_DOC_INTAKE_HANDLED_CALLBACKS: set[str] = set()  # Telegram callback_query.id de-duplication
DOCFILL_VRAM_MIN_FREE_GB: float = float(os.environ.get("DOCFILL_VRAM_MIN_FREE_GB", "36"))
_VEO_GEN_BUSY: set[int] = set()
_VEO_GEN_BUSY_LOCK = asyncio.Lock()
_AXI_DIALOG: dict[int, list[dict[str, str]]] = {}
_SEEN_UPDATES: dict[str, float] = {}
_SEEN_UPDATES_LOCK = Lock()

# Путь к БД Omi и рабочему пространству (для чтения документов)
OMI_DB_PATH   = Path(os.environ.get("OMI_DB_PATH",   "/data/aims_registry.db"))
AIMS_WORKSPACE = Path(os.environ.get("AIMS_WORKSPACE", "/data"))
_AXI_DIALOG_PATH = AIMS_WORKSPACE / "axi_dialog_state.json"
_AXI_DIALOG_LOCK = Lock()
_AXI_PENDING_ANALYZE_PATH = AIMS_WORKSPACE / "axi_pending_analyze_state.json"
_AXI_PENDING_LOCK = Lock()

# ── Document intake router ─────────────────────────────────────────────────────
_DI_INTAKE_DIR   = Path(os.environ.get("TELEGRAM_DOC_INTAKE_DIR", "/data/telegram_intake/cases"))
_DI_INBOX_DIR    = Path(os.environ.get("DOCS_AGENT_INBOX_DIR",    "/data/docs_agent/inbox"))
_DI_SYNC_ENABLED = os.environ.get("DOCS_AGENT_TELEGRAM_SYNC_ENABLED", "0") == "1"

# ── System prompt (axi_bot_refocus_patch.md) ───────────────────────────────────

AXI_SYSTEM_PROMPT = os.environ.get(
    "AXI_SYSTEM_PROMPT",
    "You are Axi, a practical AI assistant embedded in a Telegram chat. "
    "Your primary job: answer questions, analyse documents, draft text, and generate files — directly, without ceremony.\n\n"
    "BEHAVIOUR RULES (follow strictly):\n"
    "1. Be concise. Do NOT open with a self-introduction, capability statement, or summary of what you 'see in the system' "
    "unless the user asked for it. Jump straight to the answer or the action.\n"
    "2. Do NOT volunteer ISO 55001 / AIMS / P-code / E-code mappings unless the user explicitly asks for them. "
    "If the user asks a general question (safety, engineering, drafting, web search, coding, etc.) — "
    "answer it as a general assistant.\n"
    "3. You CAN generate Word (.docx) files yourself. When a user asks for a report, updated document, or analysis "
    "in Word format: generate the content, save it as .docx using python-docx, and send it back in the same chat. "
    "Do NOT tell the user to ask Omi for docx generation — you handle it directly.\n"
    "4. Omi is a separate bot that manages the AIMS database and filesystem (move files, backup DB). "
    "Only redirect to Omi for: file moves/archives, DB backups, /docgen of multi-document bundles. "
    "For registry queries (what documents are registered, search by keyword/type/date, show full registry) — "
    "query aims_registry.db directly and return the result to the user. Do NOT tell the user to ask Omi.\n"
    "5. When a task file is in context (OCR text available): use it for the answer, but do NOT repeat the task metadata "
    "('Task #N, inbox file, ISO mapping…') every time. Mention it only if directly relevant.\n"
    "6. Web search: not available in local-only mode.\n"
    "7. Language: reply in English by default. Switch to the user's language only if they explicitly write in another language.\n"
    "8. For document review/improvement tasks: do the work (read OCR text, apply requested standards/rules, produce "
    "the updated content or file). Do not describe your plan at length — just do it and send the result.\n\n"
    "TERMINOLOGY (do not confuse):\n"
    "- **register a file / register file** = verb phrase: the *action* of adding a file into the AIMS registry database "
    "(ingest, OCR pipeline, row in aims_registry).\n"
    "- **the register / registry** = noun: the *document registry* as a system (the database of records), not a single file.\n"
    "- In user messages, «register file» usually means *perform registration*, not «the register object».\n\n"
    "ARTIFACT INTEGRITY (critical):\n"
    "- Do not claim a pipeline task is complete or that a file was delivered to the chat unless the runtime actually "
    "wrote the artifact and sent it via Telegram in the same handler turn.\n"
    "- Do not invent progress percentages unless they come from task DB state.\n"
    "- Do not promise that the file will be attached in the next message; either attach now or say generation failed.\n"
    "- If the user says they did not receive a file: acknowledge, check task/tools status — do not fabricate errors.",
).strip()


def _load_axi_skill_context() -> str:
    """Load Axi runtime skill pack for internal prompt injection."""
    try:
        from agents.agent_skill_loader import load_agent_skill_text

        text = load_agent_skill_text("axi")
        if text:
            digest = hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]
            log.info("Axi skill pack loaded: agent=axi chars=%d sha256=%s", len(text), digest)
        else:
            log.warning("Axi skill pack empty: agent=axi")
        return text or ""
    except Exception as exc:
        log.warning("Axi skill pack load failed: %s", exc)
        return ""


AXI_SKILL_CONTEXT = _load_axi_skill_context()


def _axi_system_prompt(base_prompt: str | None = None) -> str:
    """Build Axi system prompt with optional runtime skill context."""
    base = base_prompt if base_prompt is not None else AXI_SYSTEM_PROMPT
    if AXI_SKILL_CONTEXT:
        return base + "\n\n# Axi runtime skill pack\n" + AXI_SKILL_CONTEXT
    return base

# ── Task Registry client ───────────────────────────────────────────────────────

sys.path.insert(0, str(Path(__file__).parent))
sys.path.insert(0, str(Path(__file__).parent / "omi_telegram"))
try:
    from task_registry_api import TaskRegistryClient as _TRClient
    _tr_client: "_TRClient | None" = _TRClient()
    log.info("Task Registry client: %s", _tr_client.base_url)
except Exception as _tre:
    _tr_client = None
    log.warning("Task Registry unavailable: %s", _tre)


def _tr_register(desc: str, chat_id: str, source: str = "axi") -> str:
    if _tr_client is None:
        return ""
    try:
        return _tr_client.register(desc, source=source, chat_id=chat_id) or ""
    except Exception as e:
        log.debug("_tr_register: %s", e)
        return ""

def _tr_start(tid: str, assigned_to: str = "axi") -> None:
    if not tid or _tr_client is None:
        return
    try:
        _tr_client.start(tid, assigned_to=assigned_to)
    except Exception as e:
        log.debug("_tr_start: %s", e)

def _tr_done(tid: str, summary: str = "") -> None:
    if not tid or _tr_client is None:
        return
    try:
        _tr_client.done(tid, result_summary=summary[:200])
    except Exception as e:
        log.debug("_tr_done: %s", e)

def _tr_stuck(tid: str, error: str = "") -> None:
    if not tid or _tr_client is None:
        return
    try:
        _tr_client.stuck(tid, error=error[:200])
    except Exception as e:
        log.debug("_tr_stuck: %s", e)


_TASK_CLOSE_PHRASES = (
    "выполнено",
    "готово",
    "закрыть задачу",
    "закрой задачу",
    "close task",
    "task done",
    "completed",
    "resolved",
)
_TASK_ID_RE = re.compile(r"\b(task_[0-9a-f]{12,})\b", re.IGNORECASE)


def _extract_task_id(text: str) -> str:
    m = _TASK_ID_RE.search(text or "")
    return (m.group(1) if m else "").lower()


def _looks_like_task_close_intent(text: str) -> bool:
    t = (text or "").strip().lower()
    if not t:
        return False
    if _extract_task_id(t):
        return True
    return any(p in t for p in _TASK_CLOSE_PHRASES)


def _close_omi_task_from_chat(chat_id: str, text: str) -> tuple[bool, str]:
    """
    Best-effort close flow for user's "done/close task" messages.
    Prevents creating a new task when user actually wants to close a hanging Omi task.
    """
    if _tr_client is None:
        return False, "Task Registry недоступен."
    wanted_id = _extract_task_id(text)
    try:
        if wanted_id:
            task = _tr_client.get(wanted_id) or {}
            if str(task.get("chat_id", "")) != str(chat_id):
                return True, f"Задача `{wanted_id}` не найдена в этом чате."
            status = str(task.get("status", "")).lower()
            if status in {"done", "failed"}:
                return True, f"Задача `{wanted_id}` уже закрыта (`{status}`)."
            _tr_done(wanted_id, summary="manual_close_by_user")
            return True, f"✅ Закрыл задачу `{wanted_id}` по вашему подтверждению."

        stuck = _tr_client.find_stuck(older_than_minutes=0) or []
        candidates = [
            t for t in stuck
            if str(t.get("chat_id", "")) == str(chat_id)
            and str(t.get("assigned_to", "")).lower() == "omi"
            and str(t.get("status", "")).lower() in {"in_progress", "stuck", "pending"}
        ]
        if not candidates:
            return True, "В этом чате нет зависших задач Omi для закрытия."

        # Prefer the oldest hanging one to drain backlog deterministically.
        task = sorted(candidates, key=lambda x: (x.get("created_at") or "", x.get("id") or 0))[0]
        tid = str(task.get("task_id", "")).strip()
        if not tid:
            return True, "Нашёл зависшую задачу, но без task_id — не могу закрыть автоматически."
        _tr_done(tid, summary="manual_close_by_user")
        return True, f"✅ Закрыл зависшую задачу `{tid}` по вашему подтверждению."
    except Exception as e:
        log.warning("manual task close failed chat=%s: %s", chat_id, e)
        return True, f"Не смог закрыть задачу автоматически: {type(e).__name__}."

# ── Owner notification helper ──────────────────────────────────────────────────

import httpx as _httpx

async def _notify_owners_async(text: str) -> None:
    """Send a plain-text message to all OWNER_CHATS via Bot API (no context needed)."""
    if not AXI_BOT_TOKEN or not OWNER_CHATS:
        return
    url = f"https://api.telegram.org/bot{AXI_BOT_TOKEN}/sendMessage"
    async with _httpx.AsyncClient(timeout=10.0) as client:
        for chat_id in OWNER_CHATS:
            try:
                await client.post(url, json={"chat_id": chat_id, "text": text})
            except Exception as e:
                log.warning("notify_owners: failed chat=%s: %s", chat_id, e)


# ── NVIDIA NIM API helper ─────────────────────────────────────────────────────

async def _gemini_reply(
    text: str,
    *,
    extra_context: str = "",
    use_search: bool = False,
    system_override: str | None = None,
) -> str:
    """Call NVIDIA NIM API (llama-3.1-405b-instruct). Returns the text response or an error string."""
    nim_url = os.environ.get("NVIDIA_NIM_URL", "http://127.0.0.1:8082").rstrip("/")
    nim_key = os.environ.get("NVIDIA_NIM_API_KEY", "")

    sys_prompt = _axi_system_prompt(system_override)
    user_text = f"{extra_context}\n\n{text}".strip() if extra_context else text

    timeout = _httpx.Timeout(60.0, connect=10.0)
    model = "meta/llama-3.1-405b-instruct"

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": sys_prompt},
            {"role": "user", "content": user_text},
        ],
        "temperature": 0.7,
        "max_tokens": 4096,
    }

    headers = {"Content-Type": "application/json"}
    if nim_key:
        headers["Authorization"] = f"Bearer {nim_key}"

    try:
        async with _httpx.AsyncClient(timeout=timeout) as client:
            url = f"{nim_url}/v1/chat/completions"
            resp = await client.post(url, json=payload, headers=headers)
            if resp.status_code == 200:
                body = resp.json()
                answer = (
                    body.get("choices", [{}])[0]
                    .get("message", {})
                    .get("content", "")
                ).strip()
                if answer:
                    return answer
            else:
                log.warning("NIM %s status=%s", model, resp.status_code)
    except (_httpx.TimeoutException, _httpx.NetworkError):
        pass
    except Exception as e:
        log.warning("NIM error: %s", e)

    return f"{AXI_NAME}: не удалось получить ответ от NIM."


# ── Anthropic fallback ─────────────────────────────────────────────────────────

async def _anthropic_reply(
    text: str,
    *,
    extra_context: str = "",
    system_override: str | None = None,
) -> str:
    """Call Claude via OmniRouter. Used as fallback when NIM is unavailable."""
    omnirouter_url = os.environ.get("AIMS_OMNIROUTER_URL", "http://127.0.0.1:20129").rstrip("/")
    auth_token = os.environ.get("AIMS_CLAUDE_PROXY_TOKEN", "aims-local-repair-token")
    model = os.environ.get("AIMS_ANTHROPIC_MODEL", ANTHROPIC_MODEL)

    sys_prompt = _axi_system_prompt(system_override)
    user_text = f"{extra_context}\n\n{text}".strip() if extra_context else text
    # Cloud fallback: truncate very large payloads to avoid 90s timeout
    _CLOUD_MAX_CHARS = 40_000
    if len(user_text) > _CLOUD_MAX_CHARS:
        log.info("omnirouter: truncating payload %d→%d chars", len(user_text), _CLOUD_MAX_CHARS)
        user_text = user_text[:_CLOUD_MAX_CHARS] + "\n\n[...truncated for cloud processing...]"

    payload = {
        "model": model,
        "max_tokens": 8096,
        "system": sys_prompt,
        "messages": [{"role": "user", "content": user_text}],
    }
    _OMNIROUTER_HDRS = {
        "Authorization": f"Bearer {auth_token}",
        "Content-Type": "application/json",
    }
    for _attempt in range(2):
        timeout = _httpx.Timeout(90.0 if _attempt == 0 else 60.0, connect=10.0)
        try:
            async with _httpx.AsyncClient(timeout=timeout) as client:
                resp = await client.post(f"{omnirouter_url}/v1/messages", json=payload, headers=_OMNIROUTER_HDRS)
            if resp.status_code == 200:
                body = resp.json()
                text_out = (body.get("content") or [{}])[0].get("text", "").strip()
                if text_out:
                    _ft_log_example(
                        [{"role": "system", "content": sys_prompt}, {"role": "user", "content": user_text}],
                        text_out,
                        source=f"omnirouter:{model}",
                    )
                    return text_out
            else:
                log.warning("omnirouter status=%s body=%s", resp.status_code, resp.text[:200])
                if resp.status_code in (402, 529):
                    asyncio.ensure_future(_notify_owners_async(
                        f"🚨 Axi: OmniRouter ошибка {resp.status_code} — "
                        f"{'средства на счёте закончились' if resp.status_code == 402 else 'сервис перегружен (529)'}."
                    ))
                break
        except Exception as e:
            log.warning("omnirouter error [attempt %d]: %s: %s", _attempt + 1, type(e).__name__, e)
            if _attempt < 1:
                await asyncio.sleep(3)
    return f"{AXI_NAME}: не удалось получить ответ от OmniRouter."


async def _anthropic_classify_intent(prompt: str, file_names: str) -> str:
    """Classify routing intent via Omnirouter (NVIDIA NIM) tool_use — returns guaranteed enum value."""
    omnirouter_url = os.environ.get("AIMS_OMNIROUTER_URL", "http://127.0.0.1:20129").rstrip("/")
    auth_token = os.environ.get("AIMS_CLAUDE_PROXY_TOKEN", "aims-local-repair-token")
    model = os.environ.get("AIMS_INTENT_MODEL", "llama405b")

    clf_tools = [{
        "name": "set_intent",
        "description": "Set the classified routing intent for this file-processing request",
        "input_schema": {
            "type": "object",
            "properties": {
                "intent": {
                    "type": "string",
                    "enum": ["docx", "edit", "chat", "ask", "training_pair"],
                    "description": (
                        "docx=generate analysis/report as Word document; "
                        "edit=modify/update data inside the spreadsheet; "
                        "chat=text answer only, no file output; "
                        "ask=ambiguous, need clarification from user; "
                        "training_pair=prepare doctuning training pair from uploaded documents"
                    ),
                }
            },
            "required": ["intent"],
        },
    }]
    payload = {
        "model": model,
        "max_tokens": 64,
        "tools": clf_tools,
        "tool_choice": {"type": "tool", "name": "set_intent"},
        "messages": [{
            "role": "user",
            "content": (
                f"Attached files: {file_names}\n"
                f"User request: {prompt}\n\n"
                "Classify the routing intent."
            ),
        }],
    }
    _OMNIROUTER_HDRS = {
        "Authorization": f"Bearer {auth_token}",
        "Content-Type": "application/json",
    }
    try:
        async with _httpx.AsyncClient(timeout=_httpx.Timeout(20.0, connect=10.0)) as client:
            resp = await client.post(f"{omnirouter_url}/v1/messages", json=payload, headers=_OMNIROUTER_HDRS)
        if resp.status_code == 200:
            for block in resp.json().get("content", []):
                if block.get("type") == "tool_use" and block.get("name") == "set_intent":
                    label = block.get("input", {}).get("intent", "ask")
                    return label if label in ("docx", "edit", "chat", "ask", "training_pair") else "ask"
        else:
            log.warning("omnirouter classify status=%s", resp.status_code)
    except Exception as e:
        log.warning("omnirouter classify error: %s", e)
    return "ask"


def _axi_llm_local_first_enabled() -> bool:
    return AXI_LLM_LOCAL_FIRST


# Semaphore: at most AXI_LLM_MAX_CONCURRENT heavy LLM calls at once.
# Prevents stacking two 32B–72B inferences on the DGX GPU simultaneously.
_llm_semaphore: asyncio.Semaphore | None = None


def _get_llm_semaphore() -> asyncio.Semaphore:
    global _llm_semaphore
    if _llm_semaphore is None:
        _llm_semaphore = asyncio.Semaphore(AXI_LLM_MAX_CONCURRENT)
    return _llm_semaphore


_PULSE_FRAMES = ("🧠", "💭", "⚙️", "🔄")


async def _animate_progress(
    msg,
    base_text: str,
    stop_event: asyncio.Event,
    *,
    interval: float = 12.0,
) -> None:
    """Edit progress_msg every interval seconds with elapsed time + rotating emoji."""
    t0 = time.perf_counter()
    idx = 0
    while not stop_event.is_set():
        await asyncio.sleep(interval)
        if stop_event.is_set():
            break
        elapsed = int(time.perf_counter() - t0)
        emoji = _PULSE_FRAMES[idx % len(_PULSE_FRAMES)]
        try:
            await msg.edit_text(f"⏱ {base_text} {elapsed}с {emoji}")
        except Exception:
            pass
        idx += 1


def _ft_log_example(
    messages: list[dict],
    response: str,
    source: str,
) -> None:
    """Append a chat example to daily JSONL for future fine-tuning. Silent on errors."""
    if not AXI_FT_LOG_ENABLED:
        return
    try:
        import json as _json
        from datetime import datetime, timezone
        AXI_FT_LOG_DIR.mkdir(parents=True, exist_ok=True)
        date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")
        record = _json.dumps(
            {
                "messages": messages + [{"role": "assistant", "content": response}],
                "source": source,
                "ts": datetime.now(timezone.utc).isoformat(),
            },
            ensure_ascii=False,
        )
        with (AXI_FT_LOG_DIR / f"axi_ft_{date_str}.jsonl").open("a", encoding="utf-8") as fh:
            fh.write(record + "\n")
    except Exception as e:
        log.debug("_ft_log_example write failed: %s", e)


def _local_ollama_reply_sync(
    text: str,
    *,
    extra_context: str = "",
    system_override: str | None = None,
) -> str:
    """Primary local registry-resolved path for Axi responses."""
    from ollama_resolve import effective_ollama_base_url, heavy_ollama_model_name

    from omi_ollama import ollama_chat

    sys_prompt = _axi_system_prompt(system_override)
    user_text = f"{extra_context}\n\n{text}".strip() if extra_context else text
    messages = [
        {"role": "system", "content": sys_prompt},
        {"role": "user", "content": user_text},
    ]
    base = effective_ollama_base_url().rstrip("/")
    if not base:
        raise RuntimeError("Ollama base URL is empty (OLLAMA_BASE_URL / DGX_OLLAMA_URL / …)")
    model = heavy_ollama_model_name()
    with record_llm_call(provider="local", model=model):
        result = ollama_chat(messages, model=model, base_url=base, timeout=300.0)
    _ft_log_example(messages, result, source=f"local:{model}")
    return result


async def _cloud_llm_reply(
    text: str,
    *,
    extra_context: str = "",
    use_search: bool = False,
    system_override: str | None = None,
) -> str:
    """Call NVIDIA NIM (llama-3.1-405b) or fallback to Anthropic."""
    with record_llm_call(provider="nim", model="meta/llama-3.1-405b-instruct"):
        result = await _gemini_reply(
            text,
            extra_context=extra_context,
            use_search=use_search,
            system_override=system_override,
        )
    if "не удалось получить ответ от NIM" in result:
        log.info("NIM unavailable — trying Anthropic fallback")
        with record_llm_call(provider="anthropic", model=ANTHROPIC_MODEL):
            result = await _anthropic_reply(
                text, extra_context=extra_context, system_override=system_override,
            )
    return result


async def _llm_reply(
    text: str,
    *,
    extra_context: str = "",
    use_search: bool = False,
    system_override: str | None = None,
) -> str:
    """Route to local Ollama or cloud, serialised by semaphore to prevent GPU saturation."""
    async with _get_llm_semaphore():
        return await _llm_reply_inner(
            text,
            extra_context=extra_context,
            use_search=use_search,
            system_override=system_override,
        )


async def _llm_reply_inner(
    text: str,
    *,
    extra_context: str = "",
    use_search: bool = False,
    system_override: str | None = None,
) -> str:
    """
    При AXI_LLM_LOCAL_FIRST=1: сначала Ollama, затем NIM.
    Fallback: Anthropic Claude.
    """
    if use_search or not _axi_llm_local_first_enabled():
        return await _cloud_llm_reply(
            text,
            extra_context=extra_context,
            use_search=use_search,
            system_override=system_override,
        )
    try:
        out = await asyncio.to_thread(
            _local_ollama_reply_sync,
            text,
            extra_context=extra_context,
            system_override=system_override,
        )
        if (out or "").strip():
            return out.strip()
    except Exception as e:
        log.warning("Axi: local Ollama failed, using NIM: %s", e)
        try:
            from failure_detector import get_detector
            get_detector().from_axi({"ok": False, "error": str(e), "source": "local_ollama"})
        except Exception:
            pass
    return await _cloud_llm_reply(
        text,
        extra_context=extra_context,
        use_search=False,
        system_override=system_override,
    )


# ── Web search ─────────────────────────────────────────────────────────────────

def _should_web_search(text: str) -> bool:
    if not AXI_WEB_SEARCH_ENABLED:
        return False
    low = text.lower()
    return any(kw in low for kw in (
        "search", "find", "look up", "google", "latest", "current", "today",
        "news", "price", "цена", "найди", "поищи", "поиск", "погода", "курс",
        "сейчас", "сегодня", "последн", "актуальн",
    ))


# ── XLSX editing ───────────────────────────────────────────────────────────────


def _parse_xlsx_table_response(
    answer: str,
    original_sheets: dict[str, list[list]],
) -> dict[str, list[list[str]]]:
    """Parse LLM table response: [SHEET:] format first, markdown table fallback."""
    result_sheets: dict[str, list[list[str]]] = {}
    current_sheet: str | None = None

    for line in answer.splitlines():
        stripped = line.strip()
        sheet_match = re.match(r"\[SHEET:\s*(.+?)\]", stripped)
        if sheet_match:
            current_sheet = sheet_match.group(1).strip()
            result_sheets[current_sheet] = []
            continue
        # Auto-assign to first original sheet if LLM skips [SHEET:] header
        if current_sheet is None and "|" in stripped and not re.match(r"^[\|\s\-:]+$", stripped):
            current_sheet = next(iter(original_sheets), "Sheet1")
            result_sheets[current_sheet] = []
        if current_sheet is None:
            continue
        # Skip empty lines and markdown separator rows (|---|---|)
        if not stripped or re.match(r"^[\|\s\-:]+$", stripped):
            continue
        if "|" in stripped:
            cells = [c.strip() for c in stripped.strip("|").split("|")]
            result_sheets[current_sheet].append(cells)
        # Skip non-table lines (explanatory text without pipes)

    return result_sheets


async def _edit_xlsx_with_llm(
    xlsx_path: Path,
    instruction: str,
) -> Path:
    """Read xlsx, ask LLM to apply instruction, write modified xlsx."""
    from openpyxl import load_workbook, Workbook  # type: ignore

    wb_in = load_workbook(str(xlsx_path), data_only=True)
    sheets_raw: dict[str, list[list]] = {}
    for ws in wb_in.worksheets:
        sheets_raw[ws.title] = [list(row) for row in ws.iter_rows(values_only=True)]
    wb_in.close()

    table_parts: list[str] = []
    for sheet_name, rows in sheets_raw.items():
        table_parts.append(f"[SHEET: {sheet_name}]")
        for row in rows:
            table_parts.append(" | ".join("" if v is None else str(v) for v in row))
    table_text = "\n".join(table_parts)
    first_sheet = next(iter(sheets_raw), "Sheet1")

    prompt = (
        f"You are editing an Excel spreadsheet. Output ONLY the table data.\n\n"
        f"Current table:\n{table_text[:60000]}\n\n"
        f"Instruction: {instruction}\n\n"
        f"Required output format:\n"
        f"[SHEET: {first_sheet}]\n"
        f"cell1 | cell2 | cell3\n"
        f"cell1 | cell2 | cell3\n\n"
        f"Rules: include ALL rows and ALL columns unchanged except where instructed. "
        f"Empty cell = empty string. Start IMMEDIATELY with [SHEET: {first_sheet}]. "
        f"No markdown, no code fences, no explanation before or after the table."
    )
    answer = await _llm_reply(prompt)
    if not answer.strip() or answer.startswith(f"{AXI_NAME}:"):
        raise RuntimeError(answer.strip() or "LLM returned empty response")

    result_sheets = _parse_xlsx_table_response(answer, sheets_raw)
    if not result_sheets:
        raise ValueError(f"LLM returned no parseable table data (response length: {len(answer)})")

    wb_out = Workbook()
    wb_out.remove(wb_out.active)
    for sheet_name, rows in result_sheets.items():
        ws_out = wb_out.create_sheet(title=sheet_name[:31])
        for row in rows:
            ws_out.append(row)

    AXI_RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    out_path = AXI_RESULTS_DIR / f"{xlsx_path.stem}_edited_{int(time.time())}.xlsx"
    wb_out.save(str(out_path))
    return out_path


async def _transform_xlsx_with_code(
    xlsx_path: Path,
    code_instruction: str,
) -> tuple[Path, str]:
    """Ask LLM to generate transformation code, run it in sandbox on XLSX.

    The LLM generates Python code that:
      - reads from input.xlsx via openpyxl
      - modifies the data
      - saves to input.xlsx (or output.xlsx)

    Returns (output_xlsx_path, stdout_from_execution)
    """
    from ops.workers.engineer_worker import execute_code_on_xlsx

    # Ask LLM to generate the transformation code
    prompt = (
        f"Generate ONLY Python code (no markdown, no explanation) to transform an Excel file.\n\n"
        f"Requirements:\n"
        f"  - Use openpyxl to open 'input.xlsx'\n"
        f"  - Perform this transformation: {code_instruction}\n"
        f"  - Save changes to 'input.xlsx' (or set output_xlsx = Path('output.xlsx'))\n"
        f"  - Print progress updates\n"
        f"  - No error handling (sandbox will catch exceptions)\n\n"
        f"Code must be executable, no imports outside openpyxl/pathlib/json/csv."
    )
    code = await _llm_reply(prompt)
    if not code.strip() or code.startswith(f"{AXI_NAME}:"):
        raise RuntimeError(code.strip() or "LLM returned empty code")

    # Run code in sandbox
    output_path, stdout, stderr = execute_code_on_xlsx(
        xlsx_path,
        code,
        timeout=30,
    )
    if stderr and "Traceback" in stderr:
        raise RuntimeError(f"Code execution failed:\n{stderr}")

    AXI_RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    final_path = AXI_RESULTS_DIR / f"{output_path.stem}_{int(time.time())}.xlsx"
    output_path.rename(final_path)

    return final_path, stdout


# ── DOCX generation ────────────────────────────────────────────────────────────

async def _generate_custom_docx(
    content: str,
    filename_stem: str,
    out_dir: Path,
    *,
    title: str | None = None,
) -> Path:
    """
    Generate a .docx file from LLM markdown-like content.
    Parses # ## ### headings and **bold** markers into Word styles.
    """
    from docx import Document

    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    safe = re.sub(r"[^\w\-]", "_", filename_stem)[:60]
    path = out_dir / f"{safe}_{stamp}.docx"

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    bold_re = re.compile(r"\*\*(.+?)\*\*")

    def _add_para(paragraph, line: str) -> None:
        parts = bold_re.split(line)
        for i, part in enumerate(parts):
            run = paragraph.add_run(part)
            if i % 2 == 1:
                run.bold = True

    for line in content.splitlines():
        stripped = line.rstrip()
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith(("- ", "* ")):
            p = doc.add_paragraph(style="List Bullet")
            _add_para(p, stripped[2:])
        elif stripped == "":
            doc.add_paragraph("")
        else:
            p = doc.add_paragraph()
            _add_para(p, stripped)

    doc.save(str(path))
    return path


def _wants_docx(text: str) -> bool:
    norm = text.lower()
    if _is_docsreg_launch_text(norm):
        # DOCSREG is Omi-owned intent. Axi must not claim it as a docx task.
        return False
    patterns = (
        r"(word|docx|\.docx|в word|в ворд|ворд)",
        r"(return|send|give|верни|пришли|отправь).{0,30}(word|docx)",
        r"(generate|create|сгенерируй|создай).{0,30}(docx|word|report|отчёт|отчет)",
        r"(report|отчёт|отчет|документ).{0,30}(word|docx)",
        r"(generate|create|prepare|draft|combine|review|develop|write).{0,40}(procedure|report|document|manual|strategy)",
        r"(generate|create|prepare|draft|combine|review|develop|write).{0,40}(policy|framework|policy framework|policy_framework)",
        r"(сгенерируй|создай|подготовь|собери|проверь).{0,40}(процедур|отч[её]т|документ|регламент|стратег)",
        r"(создай|подготовь|разработай|сделай).{0,40}(политик|рамк|framework|policy)",
        # Analysis + send/deliver intent (no explicit "word" but clearly wants a document output)
        r"(check|analyze|analyse|review|compare|assess|audit).{0,60}(send|deliver|share|attach|chat|file)",
        r"(send|deliver|share).{0,20}(it|file|result|output).{0,20}(to chat|to us|to me)",
        r"(correct|improve|update).{0,30}file.{0,30}(best practice|standard|based on)",
        r"(проверь|проанализируй|сравни).{0,60}(отправь|пришли|прикрепи)",
    )
    return any(re.search(p, norm) for p in patterns)


def _wants_xlsx_edit(prompt: str, files: list[dict]) -> bool:
    """True when an xlsx is in the batch and the prompt is an edit/update instruction."""
    has_xlsx = any(
        (f.get("name") or "").lower().endswith((".xlsx", ".xlsm"))
        for f in files
    )
    if not has_xlsx:
        return False
    # Docx/analysis intent takes priority over table editing
    if _wants_docx(prompt):
        return False
    low = (prompt or "").lower()
    edit_keywords = (
        "обнови", "измени", "редактируй", "заполни", "добавь", "удали", "исправь",
        "обновить", "изменить", "отредактируй", "поправь", "вставь", "скопируй",
        "update", "edit", "fill", "modify", "change", "add", "delete", "insert",
        "fix", "correct", "replace", "set", "revise",
    )
    return any(kw in low for kw in edit_keywords)


async def _infer_file_intent(prompt: str, files: list[dict]) -> str:
    """Classify file-processing intent.

    Returns one of: 'docx' | 'edit' | 'chat' | 'ask' | 'training_pair'
    Pipeline: keyword fast-path → local Qwen → NIM (primary cloud) → Anthropic (last resort).
    """
    file_names = ", ".join(f.get("name", "?") for f in files[:5])

    async def _cloud_classify(p: str, fn: str) -> str:
        try:
            from router.nim_router import classify_file_intent as _nim
            intent, confidence = _nim(p, fn)
            if confidence >= 0.7 and intent != "ask":
                log.debug("nim_router: intent=%s conf=%.2f", intent, confidence)
                return intent
        except Exception as _e:
            log.debug("nim_router classify error: %s", _e)
        return await _anthropic_classify_intent(p, fn)

    try:
        from router.fallback_router import classify_file_intent as _route
        return await _route(
            prompt,
            file_names,
            anthropic_fallback_fn=_cloud_classify,
        )
    except Exception as e:
        log.warning("_infer_file_intent router error: %s — falling back to cloud", e)
        return await _cloud_classify(prompt, file_names)


def _save_intent_example(prompt: str, files: list[dict], intent: str) -> None:
    """Append confirmed (prompt, intent) pair to self-learning log for future fine-tuning."""
    import json as _json
    record = {
        "ts": int(time.time()),
        "prompt": prompt,
        "files": [f.get("name", "") for f in files],
        "intent": intent,
    }
    try:
        _SELF_LEARN_LOG.parent.mkdir(parents=True, exist_ok=True)
        with _SELF_LEARN_LOG.open("a", encoding="utf-8") as fh:
            fh.write(_json.dumps(record, ensure_ascii=False) + "\n")
    except Exception as e:
        log.debug("_save_intent_example write failed: %s", e)


def _strip_false_attachment_claims(answer: str, *, lang: str) -> str:
    """
    Guardrail for chat-only path: remove model claims that a file is attached/generated
    when no Telegram document was actually sent by runtime.
    """
    text = (answer or "").strip()
    if not text:
        return text
    low = text.lower()
    claim_markers = (
        "is attached",
        "has been generated and is attached",
        "file is attached",
        "document is attached",
        "sending generated file",
        "вложен",
        "файл приложен",
        "документ приложен",
        "прикреплен",
        "прикреплён",
    )
    if not any(m in low for m in claim_markers):
        return text
    lines: list[str] = []
    for ln in text.splitlines():
        ll = ln.lower()
        if any(m in ll for m in claim_markers):
            continue
        lines.append(ln)
    cleaned = "\n".join(lines).strip() or text
    note = (
        "\n\nNote: this response is text-only. No file was attached in this message."
        if lang != "ru"
        else "\n\nПримечание: это текстовый ответ. Файл к этому сообщению не прикреплялся."
    )
    return cleaned + note


def _strip_code_fences(answer: str) -> str:
    """
    Remove fenced code blocks from user-facing chat replies.
    Keeps response focused on conclusions instead of implementation details.
    """
    text = (answer or "").strip()
    if not text:
        return ""
    cleaned = re.sub(r"```[\s\S]*?```", "", text, flags=re.MULTILINE).strip()
    return cleaned or text


def _strip_generate_file_json(text: str) -> str:
    """
    Intercept model output wrapped in {"action":"generate_file",...}.
    Extracts the 'content' field when present; otherwise removes the JSON block.
    Prevents raw pipeline JSON from ending up verbatim in a .docx file.
    """
    text = (text or "").strip()
    if '{"action"' not in text and '"generate_file"' not in text:
        return text
    import json as _json
    json_start = text.find('{"action"')
    if json_start == -1:
        return text
    # Walk forward to find matching closing brace
    brace_depth = 0
    in_str = False
    escape_next = False
    json_end = -1
    for i, ch in enumerate(text[json_start:], start=json_start):
        if escape_next:
            escape_next = False
            continue
        if ch == "\\" and in_str:
            escape_next = True
            continue
        if ch == '"':
            in_str = not in_str
            continue
        if in_str:
            continue
        if ch == "{":
            brace_depth += 1
        elif ch == "}":
            brace_depth -= 1
            if brace_depth == 0:
                json_end = i + 1
                break
    if json_end == -1:
        return text
    json_blob = text[json_start:json_end]
    try:
        obj = _json.loads(json_blob)
        content = obj.get("content", "")
        if content and len(content) > 50:
            before = text[:json_start].strip()
            after = text[json_end:].strip()
            parts = [p for p in [before, content, after] if p]
            return "\n\n".join(parts)
    except Exception as e:
        log.debug("_strip_generate_file_json parse failed: %s", e)
    # Fallback: strip the JSON block, keep surrounding text
    before = text[:json_start].strip()
    after = text[json_end:].strip()
    return (before + ("\n\n" + after if after else "")).strip() or text


def _build_doc_analysis_system() -> str:
    """
    System prompt override for deep structured technical document analysis
    when the source contains equipment tables, preservation matrices, or PPM guidelines.
    Injected for the docx-output path in _process_analyze_batch when xlsx is present.
    """
    return (
        AXI_SYSTEM_PROMPT + "\n\n"
        "DOCUMENT ANALYSIS MODE — MANDATORY RULES:\n"
        "1. OUTPUT: Write directly as a professional engineering document using Markdown "
        "(# ## ### headings, bullet lists, **bold**). "
        "NEVER output JSON, code blocks, or {\"action\":...} constructs — "
        "write the document text directly with no wrapper.\n"
        "2. COMPLETENESS: If the source document lists equipment types, systems, or processes "
        "— enumerate EVERY item by its exact name. Do not summarize, skip, or merge any item.\n"
        "3. PER-EQUIPMENT DEPTH: For each equipment type identified, provide a dedicated subsection with:\n"
        "   - Gap analysis: what is currently missing or insufficient in the source\n"
        "   - 3–5 specific, actionable improvements with engineering rationale\n"
        "   - Applicable standards with full designation "
        "(e.g., API 610 12th ed. §9.3, ISO 55001:2014 §8.1, NACE SP0169-2013, IEC 60034-1)\n"
        "   - Quantitative acceptance criteria where applicable "
        "(e.g., vibration ≤ 4.5 mm/s RMS per ISO 10816-3, insulation ≥ 100 MΩ at 500 V DC)\n"
        "4. DOCUMENT STRUCTURE (use this unless the user specified otherwise):\n"
        "   # [Document Title]\n"
        "   ## Executive Summary\n"
        "   ## Scope and Applicability\n"
        "   ## Equipment-by-Equipment Analysis\n"
        "   ### [Equipment Type 1]\n"
        "   ### [Equipment Type 2]\n"
        "   ...\n"
        "   ## Cross-Cutting Recommendations\n"
        "   ## Standards and Reference Matrix\n"
        "5. STANDARDS: Always cite specific standard numbers and sections. "
        "Never use vague phrases like 'per industry standards' or 'according to best practice' "
        "without naming the actual standard.\n"
        "6. LENGTH: Do not truncate. Every equipment type in the source must appear as its own "
        "subsection with full analysis. A complete analysis is expected."
    )


def _wants_standards_docx_result(text: str) -> bool:
    low = (text or "").lower()
    if _is_docsreg_launch_text(low):
        # DOCSREG launch semantics are delegated to Omi, not Axi.
        return False
    if _wants_docx(low):
        return True
    return any(k in low for k in (
        "correct", "fix", "revise", "update", "redline",
        "исправ", "скоррект", "обнови", "внеси коррект",
    ))


def _build_applied_corrections_section(compliance_review: str) -> str:
    """
    Build a concise delta block to show what changed vs original.
    """
    lines = [(ln or "").strip("-• ").strip() for ln in (compliance_review or "").splitlines()]
    candidates: list[str] = []
    seen: set[str] = set()
    markers = ("gap", "non-com", "critical", "major", "minor", "recommend", "clause", "required")
    for ln in lines:
        low = ln.lower()
        if len(ln) < 18:
            continue
        if not any(m in low for m in markers):
            continue
        if ln in seen:
            continue
        seen.add(ln)
        candidates.append(ln[:220])
        if len(candidates) >= 8:
            break
    if not candidates:
        candidates = [
            "Updated procedure steps to align with cited international standards and clauses.",
            "Added missing control points, verification checkpoints, and records required for compliance.",
            "Clarified acceptance criteria, responsibilities, and evidence requirements.",
        ]
    body = "\n".join(f"- {item}" for item in candidates[:8])
    return "## Applied corrections vs original\n" + body


def _sources_note(
    *,
    lang: str,
    use_search: bool,
    has_dialog_context: bool = False,
    uploaded_files_count: int = 0,
) -> str:
    if not AXI_CHAT_SHOW_SOURCES:
        return ""
    if lang == "ru":
        lines = ["Источники:"]
        lines.append("- запрос пользователя")
        if has_dialog_context:
            lines.append("- контекст предыдущего диалога")
        if uploaded_files_count > 0:
            lines.append(f"- содержимое загруженных файлов ({uploaded_files_count})")
        lines.append("- web search" if use_search else "- без web search")
        return "\n".join(lines)
    lines = ["Sources used:"]
    lines.append("- user request")
    if has_dialog_context:
        lines.append("- recent chat context")
    if uploaded_files_count > 0:
        lines.append(f"- uploaded file content ({uploaded_files_count})")
    lines.append("- web search" if use_search else "- no web search")
    return "\n".join(lines)


# ── Standards / compliance check ─────────────────────────────────────────────

def _wants_standards_check(
    text: str,
    *,
    recent_chat_context: str = "",
    pending_files_count: int = 0,
) -> bool:
    """Detect "check against standards / compliance check" intent."""
    low = (text or "").lower()
    check_kw = any(k in low for k in (
        "проверь на стандарт", "проверь стандарт", "проверь соответствие",
        "check against standard", "check standard", "check compliance",
        "compliance check", "standards check", "standards review",
        "на соответствие стандарт", "verify standard", "сверь со стандарт",
        "аудит стандарт", "standards audit",
        "international standards", "with standards", "against standards",
        "according to standards", "per standard", "as per standard",
        "международн", "по стандарт", "по международн",
    ))
    contextual_pair = (
        ("standard" in low or "standards" in low or "international" in low)
        and any(v in low for v in ("check", "verify", "review", "audit", "correct", "fix", "update"))
    )
    if check_kw or contextual_pair:
        return True

    # Live-chat intent resolution: "this procedure/document" + corrective verbs
    # should route to standards flow when there is recent doc context.
    deictic_ref = any(k in low for k in (
        "this procedure", "this document", "this file", "that procedure", "that document",
        "эту процедуру", "этот документ", "этот файл", "данную процедуру", "этот регламент",
    ))
    corrective_action = any(k in low for k in (
        "correct", "fix", "revise", "update", "improve", "bring in line", "align",
        "исправ", "скоррект", "обнов", "доработ", "приведи в соответств",
    ))
    context_low = (recent_chat_context or "").lower()
    has_recent_doc_context = ("[doc:" in context_low) or ("uploaded file content" in context_low)
    has_pending_files = pending_files_count > 0
    return deictic_ref and corrective_action and (has_recent_doc_context or has_pending_files)


AIMS_ENABLE_CONTEXTUAL_STANDARD_DISCOVERY = os.environ.get(
    "AIMS_ENABLE_CONTEXTUAL_STANDARD_DISCOVERY", "0"
).strip().lower() in ("1", "true", "yes", "on")


def _wants_contextual_discovery(text: str) -> bool:
    """Detect contextual standard discovery / benchmark review intent."""
    low = (text or "").lower()
    markers = (
        "standard discovery",
        "benchmark matrix",
        "gap assessment",
        "compare with public guidance",
        "compare against oem guidance",
        "review procedure against standards",
        "contextual standard discovery",
        "document review benchmark",
        "anonymized standards review",
        "подбор стандартов",
        "матрица требований",
        "анализ разрывов",
        "сравни со стандартами",
        "проверь по лучшим практикам",
        "oem рекомендации",
        "бенчмарк документов",
    )
    if any(m in low for m in markers):
        return True
    return (
        ("benchmark" in low or "бенчмарк" in low or "gap" in low or "разрыв" in low)
        and ("document" in low or "procedure" in low or "документ" in low or "процедур" in low)
    )


async def _handle_contextual_standard_discovery(
    update,
    ctx,
    text: str,
    *,
    lang: str,
    tr_id: str | None,
    progress_msg,
    t0: float,
) -> None:
    """
    Feature-flagged advisory route for contextual standards discovery.
    Safety: no external document text transfer, no slot120 usage, no training export by default.
    """
    advisory: str | None = None
    try:
        from agents.skills.contextual_standard_discovery_and_document_review import run_axi_advisory_mode
        advisory = run_axi_advisory_mode(text, force_fixture=True)
    except Exception as exc:
        log.info("contextual_standard_discovery route fallback: %s", exc)

    if not advisory:
        advisory = (
            "Context recognized: contextual standard discovery / document benchmark review.\n"
            "Suggested benchmark themes:\n"
            "- scope and objective clarity\n"
            "- responsibilities and approval flow\n"
            "- lifecycle phase controls (preservation/shutdown/commissioning)\n"
            "- evidence and acceptance criteria\n"
            "- risk and mitigation coverage\n"
            "- references to public guidance and OEM recommendations\n\n"
            "For full benchmark matrix and gap assessment, route this document review to Doci.\n"
            "Safety: no full document text was sent externally; no slot120 teacher/judge/generator used; "
            "no training export executed unless explicitly requested."
            if lang != "ru"
            else
            "Контекст распознан: contextual standard discovery / benchmark-review документа.\n"
            "Рекомендуемые темы бенчмарка:\n"
            "- scope и цель документа\n"
            "- роли/ответственности и цепочка согласования\n"
            "- контроль по фазе жизненного цикла (preservation/shutdown/commissioning)\n"
            "- критерии приемки и доказательная база\n"
            "- риски и меры снижения\n"
            "- ссылки на публичные guidance и OEM рекомендации\n\n"
            "Для полного benchmark matrix и gap assessment направьте полный review в Doci.\n"
            "Безопасность: полный текст документа внешним провайдерам не отправлялся; "
            "slot120 как teacher/judge/generator не использовался; "
            "training export не выполнялся без явного запроса."
        )

    if progress_msg is not None:
        try:
            await progress_msg.delete()
        except Exception:
            pass
    await update.message.reply_text(advisory)
    if tr_id:
        _tr_done(tr_id, summary="contextual_standard_discovery_advisory")


async def _handle_standards_check(
    update,
    ctx,
    text: str,
    *,
    lang: str,
    tr_id: str | None,
    progress_msg,
    t0: float,
) -> None:
    """
    Pipeline: extract recent files + Omi DB context → LLM compliance analysis → text report.
    """
    _chat_id = update.effective_chat.id if update.effective_chat else None

    # Gather file content from pending analyze state
    file_texts: list[str] = []
    state = _get_pending_analyze_state(_chat_id) if _chat_id else None
    pending_files: list[object] = list((state or {}).get("files") or [])
    for item in pending_files[:3]:
        try:
            if isinstance(item, dict):
                name = str(item.get("name") or "uploaded_document")
                excerpt = str(item.get("excerpt") or "").strip()
                if excerpt:
                    file_texts.append(f"[FILE: {name}]\n{excerpt[:8000]}")
                continue
            fp = Path(str(item))
            if not fp.exists():
                continue
            ext = fp.suffix.lower()
            content = ""
            if ext == ".pdf":
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(str(fp))
                    content = "\n".join(p.extract_text() or "" for p in reader.pages[:8])
                except Exception as e:
                    log.warning("standards_check: PDF extract failed %s: %s", fp.name, e)
            elif ext in (".docx",):
                try:
                    from docx import Document as _Doc
                    doc = _Doc(str(fp))
                    content = "\n".join(p.text for p in doc.paragraphs)[:12000]
                except Exception as e:
                    log.warning("standards_check: DOCX extract failed %s: %s", fp.name, e)
            elif ext in (".txt", ".md"):
                content = fp.read_text(errors="replace")[:12000]
            if content.strip():
                file_texts.append(f"[FILE: {fp.name}]\n{content[:8000]}")
        except Exception as _e:
            log.debug("standards_check: file read error %s: %s", item, _e)

    if not file_texts:
        recent_doc_name, recent_doc_text = _load_recent_generated_doc_context(_chat_id)
        if recent_doc_text:
            file_texts.append(f"[FILE: {recent_doc_name}]\n{recent_doc_text[:8000]}")

    # Gather Omi registry context (top-5 docs by keyword)
    db_context = ""
    try:
        import sqlite3 as _sql
        from aims_paths import workspace_root as _wr
        _db = _wr() / "aims_registry.db"
        if _db.exists():
            kws = _extract_search_keywords(text)
            _conn = _sql.connect(str(_db))
            _conn.row_factory = _sql.Row
            rows: list = []
            for kw in kws[:4]:
                rows += _conn.execute(
                    "SELECT title, summary FROM documents WHERE title LIKE ? OR summary LIKE ? LIMIT 3",
                    (f"%{kw}%", f"%{kw}%"),
                ).fetchall()
            _conn.close()
            seen: set[str] = set()
            parts: list[str] = []
            for r in rows[:6]:
                t = str(r["title"] or "").strip()
                s = (str(r["summary"] or "").strip())[:400]
                if t and t not in seen:
                    seen.add(t)
                    parts.append(f"• {t}: {s}")
            if parts:
                db_context = "AIMS Registry excerpts:\n" + "\n".join(parts)
    except Exception as _dbe:
        log.debug("standards_check: db context error: %s", _dbe)

    # Build compliance prompt
    doc_block = "\n\n".join(file_texts) if file_texts else ""
    compliance_prompt = (
        f"{text}\n\n"
        + (f"[UPLOADED DOCUMENTS]\n{doc_block}\n\n" if doc_block else "")
        + (f"[AIMS REGISTRY CONTEXT]\n{db_context}\n\n" if db_context else "")
        + "INSTRUCTION: Perform a compliance review against applicable standards "
        "(ISO 55001/55002 Asset Management, API RP 570/574/653, ASME, NFPA 70/72, "
        "OSHA 1910, EN 13460, DEP/Shell standards). "
        "Structure output as:\n"
        "1. **Document / scope reviewed**\n"
        "2. **Applicable standards** (explicit standard numbers/revisions + rationale)\n"
        "3. **Compliance findings** (per standard: Compliant / Gap / N/A + evidence)\n"
        "4. **Recommended actions** (priority: Critical / Major / Minor)\n"
        "5. **Standards sources used** (for each standard include source name, publisher, year/revision, and URL if known)\n"
        "6. **Summary score** (0–100%)\n"
        "Be specific; cite clause numbers where relevant. "
        "Do not output Python code, pseudo-code, or fenced markdown code blocks. "
        "Return plain text only (no markdown tables, no markdown formatting markers). "
        "Do not write vague names like 'international standard' without exact identifiers (e.g., ISO 55001:2014 clause 6.2)."
    )

    use_search = AXI_WEB_SEARCH_ENABLED
    if progress_msg is not None:
        try:
            await progress_msg.edit_text(
                f"⏱ {AXI_NAME}: проверяю соответствие стандартам…"
                if lang == "ru" else
                f"⏱ {AXI_NAME}: running standards compliance check…"
            )
        except Exception:
            pass

    answer = await _llm_reply(compliance_prompt, use_search=use_search)
    answer = _strip_false_attachment_claims(answer, lang=lang)
    answer = _strip_code_fences(answer)

    dt = time.perf_counter() - t0
    source_hint = _sources_note(
        lang=lang,
        use_search=use_search,
        has_dialog_context=False,
        uploaded_files_count=len(file_texts),
    )
    task_hint = f" | task `{tr_id}`" if tr_id else ""
    full = answer.strip() + (f"\n\n{source_hint}{task_hint}" if (source_hint or task_hint) else "")

    if _wants_standards_docx_result(text) and answer.strip():
        original_block = "\n\n".join(file_texts) if file_texts else ""
        revised_prompt = (
            f"User request:\n{text}\n\n"
            f"[COMPLIANCE FINDINGS]\n{answer[:14000]}\n\n"
            + (f"[ORIGINAL DOCUMENTS]\n{original_block[:18000]}\n\n" if original_block else "")
            + "INSTRUCTION: Produce a fully corrected and updated procedure document that resolves the identified gaps.\n"
            + "Output only the corrected document content (final version), ready to save as .docx.\n"
            + "Use clear section headings and actionable steps. Include applicable standards with exact identifiers.\n"
            + "Do not output analysis notes, markdown code fences, or pseudo-code."
        )
        revised_doc = await _llm_reply(revised_prompt, use_search=use_search)
        revised_doc = _strip_false_attachment_claims(revised_doc, lang=lang)
        revised_doc = _strip_code_fences(revised_doc)
        generic_request_markers = (
            "please provide the procedure document",
            "once i have the content",
            "send the updated document to this chat",
        )
        if any(m in revised_doc.lower() for m in generic_request_markers):
            revised_doc = ""
        if not revised_doc.strip():
            revised_doc = answer
        if "applied corrections vs original" not in revised_doc.lower():
            revised_doc = revised_doc.strip() + "\n\n" + _build_applied_corrections_section(answer)
        first_line = next((ln.strip().lstrip("#").strip() for ln in revised_doc.splitlines() if ln.strip()), "corrected_procedure")
        docx_path = await _generate_custom_docx(revised_doc, "axi_standards_corrected", AXI_RESULTS_DIR, title=first_line)
        saved_to_kb = _register_corrected_standards_knowledge(
            docx_path=docx_path,
            user_request=text,
            compliance_review=answer,
            revised_document=revised_doc,
        )
        if progress_msg is not None:
            try:
                await progress_msg.delete()
            except Exception:
                pass
        with docx_path.open("rb") as fh:
            await update.message.reply_document(
                document=fh,
                filename=docx_path.name,
                caption=(
                    f"{AXI_NAME}: ✅ {docx_path.name} ({dt:.0f}s)\n"
                    + ("Knowledge base updated for future generations." if saved_to_kb else "Knowledge base update skipped.")
                ),
            )
        if _chat_id:
            _dialog_append(_chat_id, "assistant", f"[doc:{docx_path.name}]")
        _tr_done(
            tr_id,
            summary=(
                f"standards_docx_corrected:{docx_path.name}:"
                f"files={len(file_texts)}:search={use_search}:kb={int(saved_to_kb)}"
            ),
        )
        return

    if progress_msg is not None:
        try:
            await progress_msg.delete()
        except Exception:
            pass

    _TG_LIMIT = 4000
    chunks = [full[i : i + _TG_LIMIT] for i in range(0, len(full), _TG_LIMIT)] if full else ["…"]
    for chunk in chunks:
        await update.message.reply_text(chunk)

    if _chat_id:
        _dialog_append(_chat_id, "assistant", answer[:2000])
    _tr_done(tr_id, summary=f"standards_check:files={len(file_texts)}:search={use_search}")


# ── DB-based document strategy ────────────────────────────────────────────────

def _wants_registry_list(text: str) -> bool:
    """Detect requests to list/show/search the AIMS document registry."""
    low = text.lower()
    list_kw = any(k in low for k in (
        "реестр", "registry", "registered", "зарегистрировано", "зарегистрированы",
        "список документов", "list documents", "what documents", "какие документы",
        "покажи документы", "show documents", "все документы", "all documents",
        "что в базе", "what's in the", "в реестре", "in the registry",
        "сколько документов", "how many documents", "how many files",
    ))
    # Exclude strategy-generation requests — those go to _wants_db_strategy
    gen_kw = any(k in low for k in ("стратегия", "strategy", "подготовь", "generate", "сгенерируй"))
    return list_kw and not gen_kw


def _wants_db_strategy(text: str) -> bool:
    low = text.lower()
    has_db = any(k in low for k in ("database", "from database", "из базы", "из реестра", "из бд", "based on document"))
    has_gen = any(k in low for k in ("strategy", "стратегия", "plan", "план", "report", "отчёт",
                                      "prepare", "подготовь", "generate", "сгенерируй", "create", "создай"))
    return has_db and has_gen


def _extract_search_keywords(text: str) -> list[str]:
    from workers.data_worker import extract_search_keywords
    return extract_search_keywords(text)


def _register_corrected_standards_knowledge(
    *,
    docx_path: Path,
    user_request: str,
    compliance_review: str,
    revised_document: str,
) -> bool:
    """
    Save corrected standards output into AIMS registry for reuse in future generations.
    """
    import sqlite3

    if not OMI_DB_PATH.exists():
        log.warning("knowledge register skipped: OMI_DB_PATH not found: %s", OMI_DB_PATH)
        return False

    now = datetime.now(timezone.utc).isoformat()
    file_name = docx_path.name
    file_path = str(docx_path)
    standards_hits = sorted(set(re.findall(r"\b(?:ISO|API|ASME|NFPA|OSHA|EN)\s*[A-Z0-9./:-]*", compliance_review)))
    standards_kw = ", ".join(s for s in standards_hits if s.strip())[:500]
    summary = (
        "Corrected procedure after international standards compliance review.\n"
        f"Request: {user_request[:300]}\n"
        f"Standards: {standards_kw or 'not explicitly extracted'}\n"
        f"Findings excerpt: {compliance_review[:1200]}"
    )[:3500]
    notes = (
        "[Axi standards correction artifact]\n"
        "This record stores corrected output generated after compliance gap assessment.\n\n"
        f"[REVIEW]\n{compliance_review[:6000]}\n\n"
        f"[REVISED_DOCUMENT]\n{revised_document[:12000]}"
    )[:15000]

    try:
        conn = sqlite3.connect(str(OMI_DB_PATH))
        conn.row_factory = sqlite3.Row
        existing = conn.execute(
            "SELECT id FROM documents WHERE file_path = ? OR file_name = ? LIMIT 1",
            (file_path, file_name),
        ).fetchone()
        if existing:
            conn.execute(
                """
                UPDATE documents
                SET title = ?, summary = ?, keywords = ?, date_modified = ?, language = ?, source = ?, notes = ?
                WHERE id = ?
                """,
                (
                    "Corrected procedure (standards-reviewed)",
                    summary,
                    standards_kw,
                    now,
                    "en",
                    "axi_standards_correction",
                    notes,
                    int(existing["id"]),
                ),
            )
        else:
            conn.execute(
                """
                INSERT INTO documents (
                    file_path, file_name, file_type, title, summary, aims_process,
                    is_master, is_anonymized, language, date_added, date_modified,
                    source, notes, keywords
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    file_path,
                    file_name,
                    ".docx",
                    "Corrected procedure (standards-reviewed)",
                    summary,
                    "P06",
                    0,
                    0,
                    "en",
                    now,
                    now,
                    "axi_standards_correction",
                    notes,
                    standards_kw,
                ),
            )
        conn.commit()
        conn.close()
        return True
    except Exception as e:
        log.warning("knowledge register failed: %s", e)
        return False


def _resolve_doc_path(raw_path: str | None) -> Path | None:
    from workers.data_worker import resolve_doc_path
    return resolve_doc_path(raw_path, workspace_path=AIMS_WORKSPACE)


def _fetch_db_docs(keywords: list[str], max_docs: int = 8) -> list[dict]:
    from workers.data_worker import fetch_db_docs
    return fetch_db_docs(keywords, max_docs, db_path=OMI_DB_PATH, workspace_path=AIMS_WORKSPACE)


async def _handle_registry_list(
    update: "Update",
    ctx: "ContextTypes.DEFAULT_TYPE",
    text: str,
    tr_id: str,
) -> None:
    """Return a formatted list of registered documents from aims_registry.db."""
    from workers.data_worker import query_registry
    chat_id = update.effective_chat.id

    if not OMI_DB_PATH.exists():
        await update.message.reply_text("⚠️ aims_registry.db не найден.")
        _tr_stuck(tr_id, "db not found")
        return

    try:
        loop = asyncio.get_event_loop()
        total, rows = await loop.run_in_executor(
            None, lambda: query_registry(text, 50, db_path=OMI_DB_PATH)
        )
    except Exception as e:
        log.warning("_handle_registry_list db error: %s", e)
        await update.message.reply_text(f"⚠️ Ошибка чтения реестра: {e}")
        _tr_stuck(tr_id, str(e))
        return

    _registry_stop = {"реестр", "registry", "registered", "зарегистрировано", "документы", "documents",
                      "список", "list", "all", "все", "what", "какие", "покажи", "show", "сколько", "how"}
    search_kws = [k for k in _extract_search_keywords(text) if k not in _registry_stop]

    if not rows:
        await update.message.reply_text("📂 В реестре нет документов по вашему запросу.")
        _tr_done(tr_id, summary="empty result")
        return

    label = f"(поиск: {', '.join(search_kws)}, {len(rows)} из {total})" if search_kws else f"последние {len(rows)} из {total}"
    lines = [f"📂 Реестр AIMS — {label}:\n"]
    for r in rows:
        proc = r["aims_process"] or "—"
        title = (r["title"] or r.get("file_name") or f"#{r['id']}").replace("_", " ")
        date = (r["date_added"] or "")[:10]
        lines.append(f"• [{proc}] {title}  ({date})")

    reply = "\n".join(lines)
    # Telegram 4096-char limit — send as file if too long
    if len(reply) > 3800:
        import io
        bio = io.BytesIO(reply.encode())
        bio.name = "registry.txt"
        await ctx.bot.send_document(chat_id=chat_id, document=bio, filename="registry.txt",
                                    caption=f"📂 Реестр AIMS — {label}")
    else:
        await update.message.reply_text(reply)
    _tr_done(tr_id, summary=f"registry list: {len(rows)} docs")


async def _handle_db_strategy(
    update: "Update",
    ctx: "ContextTypes.DEFAULT_TYPE",
    text: str,
    tr_id: str,
) -> None:
    """
    Workflow:
    1. Search AIMS DB for anonymized documents matching the request
    2. Generate draft from internal content (Gemini/Anthropic)
    3. Enrich with external best-practices (web search)
    4. Output .docx
    """
    chat_id = update.effective_chat.id

    await ctx.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    await update.message.reply_text("🔍 Ищу документы в базе данных AIMS...")

    keywords = _extract_search_keywords(text)
    loop = asyncio.get_event_loop()
    docs = await loop.run_in_executor(None, _fetch_db_docs, keywords)

    if not docs:
        await update.message.reply_text(
            "⚠️ Документы по теме не найдены в базе данных AIMS. "
            "Попробуйте уточнить запрос или сначала зарегистрируйте документы через Omi."
        )
        _tr_stuck(tr_id, "no db docs found")
        return

    docs_with_content = [d for d in docs if d["content"]]
    found_list = "\n".join(f"• [{d['process']}] {d['title']}" for d in docs)
    await update.message.reply_text(
        f"📂 Найдено {len(docs)} документов ({len(docs_with_content)} с содержимым):\n{found_list}\n\n"
        f"⚙️ Генерирую стратегию на основе внутренних документов..."
    )
    await ctx.bot.send_chat_action(chat_id=chat_id, action=ChatAction.UPLOAD_DOCUMENT)

    # Build internal context — use full content if available, otherwise metadata only
    doc_context_parts = []
    for d in docs:
        header = f"### [{d['process']}] {d['title']}"
        body = d["content"][:12000] if d["content"] else f"[registered in AIMS DB — file content pending OCR/anonymization]"
        if d.get("keywords"):
            body = f"Keywords: {d['keywords']}\n{body}"
        doc_context_parts.append(f"{header}\n{body}")
    doc_context = "\n\n---\n\n".join(doc_context_parts)

    has_real_content = len(docs_with_content) > 0
    content_note = (
        "Use the document content above as primary source."
        if has_real_content else
        "Document files are registered but full text is pending OCR on server. "
        "Use document titles and metadata as structural anchors. "
        "Generate content based on industry-specific knowledge for aluminum smelting plants."
    )

    # Step 1: generate from internal documents + industry context
    # Extract specific facilities/components from the user request
    internal_prompt = (
        f"{text}\n\n"
        f"[AIMS DATABASE — REGISTERED DOCUMENTS]:\n\n{doc_context}\n\n"
        f"[INSTRUCTION]: {content_note}\n\n"
        "CRITICAL RULES:\n"
        "- This is an ALUMINUM SMELTING PLANT. Apply ONLY standards relevant to:\n"
        "  aluminum production, heavy industrial equipment, smelting/electrolysis, "
        "  gas handling, power systems, and industrial preservation.\n"
        "- RELEVANT standards: ISO 55001 (asset management), NFPA 70/72 (electrical/fire), "
        "  API 570/574 (piping inspection), API 653 (tanks), ISO 19443, DEP (Shell standards), "
        "  ASME standards, OSHA 1910, EN 13460 (maintenance documentation).\n"
        "- DO NOT use: ISO 15643 (cultural heritage), ISO 9241 (ergonomics), "
        "  or any standard unrelated to industrial plant preservation.\n"
        "- Structure: # main title, ## section per component "
        "(cast house / baking plant / rodding plant / power plant / fuel gas supply / utilities / flue gas plant), "
        "### subsections (preservation scope, critical equipment, procedures, inspection intervals), "
        "**bold** key terms, - bullet points.\n"
        "- Output document content only."
    )
    internal_answer = await _llm_reply(internal_prompt)

    # Step 2: enrich with industry-specific external standards
    external_note = ""
    if AXI_WEB_SEARCH_ENABLED:
        await update.message.reply_text("🌐 Сверяю с отраслевыми стандартами для алюминиевой промышленности...")
        enrich_prompt = (
            "What are the specific international standards and industry best practices for "
            "PRESERVATION strategy of ALUMINUM SMELTING PLANT facilities including: "
            "cast house, baking plant (anode baking), rodding plant, power plant, "
            "fuel gas supply, utilities and flue gas treatment plant?\n"
            "Focus ONLY on aluminum industry, metallurgical and heavy industrial standards. "
            "Include: DEP standards, API RPs, NFPA codes, ISO 55001/55002, EN 13460. "
            "Exclude cultural heritage or unrelated standards. "
            "Give 5-7 specific actionable points with standard references."
        )
        external_note = await _llm_reply(enrich_prompt, use_search=True)

    # Combine
    final_content = internal_answer
    if external_note and "не удалось получить ответ" not in external_note:
        final_content += (
            "\n\n## External Standards & Best Practices\n"
            "*Supplementary benchmarks from international standards:*\n\n"
            + external_note
        )

    # Generate .docx
    first_line = next(
        (ln.strip().lstrip("#").strip() for ln in final_content.splitlines() if ln.strip()),
        "Preservation_Strategy",
    )
    docx_path = await _generate_custom_docx(
        final_content, "preservation_strategy", AXI_RESULTS_DIR, title=first_line,
    )

    with docx_path.open("rb") as fh:
        await update.message.reply_document(
            document=fh,
            filename=docx_path.name,
            caption=(
                f"📄 {first_line}\n"
                f"Источники: {len(docs_with_content)} внутренних документов AIMS"
                + (" + внешние стандарты" if external_note and "не удалось" not in external_note else "")
            ),
        )
    _tr_done(tr_id, summary=f"db_strategy:{docx_path.name}:{len(docs)}docs")


# ── VRAM warm-up ───────────────────────────────────────────────────────────────

def _schedule_vram_warmup() -> None:
    """Background warm: 70B on Spark + Qwen on PC (QWEN_PC_ASSIST_STACK); после снятия 70B — прогрев малой."""
    try:
        from ollama_resolve import (
            effective_ollama_base_url,
            heavy_ollama_model_name,
            ollama_ensure_small_warm_after_heavy_gone_watcher,
            ollama_schedule_background_warm,
            ollama_warm_small_when_heavy_absent,
            spark_primary_ollama_base_url,
        )

        ollama_warm_small_when_heavy_absent()
        ollama_ensure_small_warm_after_heavy_gone_watcher()
        if not QWEN_PC_ASSIST_STACK:
            return
        heavy = heavy_ollama_model_name()
        base = spark_primary_ollama_base_url() or effective_ollama_base_url()
        ollama_schedule_background_warm(heavy, base)
    except Exception as e:
        log.debug("vram warmup: %s", e)


# ── Access control ─────────────────────────────────────────────────────────────

def _chat_allowed(update: Update) -> bool:
    if not ALLOWED_CHATS:
        return True
    if update.effective_chat is None:
        return False
    if update.effective_chat.type in ("group", "supergroup") and AXI_GROUP_ALLOW_ALL_MEMBERS:
        return True
    if update.effective_chat.id in ALLOWED_CHATS:
        return True
    # В группе effective_chat.id — это id группы (отрицательный), не user id.
    # Разрешаем участникам из белого списка без дублирования id группы в .env.
    if update.effective_chat.type in ("group", "supergroup"):
        uid = update.effective_user.id if update.effective_user else None
        if uid is not None and uid in ALLOWED_CHATS:
            return True
    return False


def _is_docgen_upgrade_origin(update: Update) -> bool:
    """Private Axi origin gate for DOCGEN upgrade batches."""
    if update.effective_chat is None:
        return False
    if update.effective_chat.type != "private":
        return False
    if not ALLOWED_CHATS and not OWNER_CHATS:
        return True
    return (
        update.effective_chat.id in OWNER_CHATS
        or update.effective_chat.id in ALLOWED_CHATS
    )


def _format_docgen_upgrade_reply(
    *,
    action_id: str,
    action_path: str,
    execution: dict[str, object],
    run_summary: dict[str, object],
    doc_types: list[str],
) -> str:
    loop_state = str(run_summary.get("loop_state") or execution.get("status") or "unknown")
    shelf_state = str(run_summary.get("shelf_state") or "unknown")
    next_step = str(run_summary.get("next_step") or "unknown")
    target_quality = float(run_summary.get("target_quality") or 0.98)
    best_quality = float(run_summary.get("best_quality") or 0.0)
    best_cycles = int(run_summary.get("best_cycles") or 0)
    avg_quality = float(run_summary.get("avg_quality") or 0.0)
    trace = str(run_summary.get("trajectory_trace") or "").strip() or "n/a"
    launcher_status = str(execution.get("status") or "unknown")
    exit_code = str(execution.get("exit_code") or "unknown")
    evidence_dir = str(execution.get("evidence_dir") or "n/a")
    training_review = bool(run_summary.get("training_review_required"))
    promotion_ready = bool(run_summary.get("promotion_ready"))
    blocked_reasons = [
        str(item)
        for item in (run_summary.get("blocked_reasons") or [])
        if str(item).strip()
    ]
    loop_completed = best_cycles > 0 or best_quality > 0.0
    branch_markers = [
        f"{item.get('document_type')}: {float(item.get('achieved_quality') or 0.0):.1%} / {int(item.get('cycles_completed') or 0)}c"
        for item in run_summary.get("batch_summaries") or []
        if isinstance(item, dict)
    ]
    branch_text = " | ".join(branch_markers) if branch_markers else "n/a"
    doc_types_text = ", ".join(doc_types)
    if loop_state in {"BLOCKED", "REVIEW"} and best_cycles <= 0 and best_quality <= 0.0:
        heading = "❌ DOCGEN runtime preflight failed."
    elif loop_state in {"COMPLETE", "TARGET_REACHED"} or best_quality >= target_quality:
        heading = "✅ DOCGEN upgrade completed."
    elif best_cycles > 0:
        heading = "🔵 DOCGEN runtime started."
    else:
        heading = "🟡 DOCGEN upgrade job accepted."
    blocker_text = ", ".join(blocked_reasons) if blocked_reasons else "n/a"
    return (
        f"{heading}\n"
        f"Action: `{action_id}`\n"
        f"Path: `{action_path}`\n"
        f"Types: `{doc_types_text}`\n"
        f"Loop state: `{loop_state}`\n"
        f"Shelf: `{shelf_state}`\n"
        f"Best quality: `{best_quality:.1%}` / target `{target_quality:.1%}`\n"
        f"Avg quality: `{avg_quality:.1%}`\n"
        f"Cycles: `{best_cycles}`\n"
        f"Trace: `{trace}`\n"
        f"Branches: `{branch_text}`\n"
        f"Training review: `{str(training_review).lower()}`\n"
        f"Promotion ready: `{str(promotion_ready).lower()}`\n"
        f"Next: `{next_step}`\n"
        f"Blocker: `{blocker_text}`\n"
        f"Launcher status: `{launcher_status}`\n"
        f"Exit code: `{exit_code}`\n"
        f"Evidence: `{evidence_dir}`"
    )

def _format_docgen_upgrade_completion(
    *,
    action_id: str,
    run_summary: dict[str, object],
) -> str:
    loop_state = str(run_summary.get("loop_state") or "unknown")
    shelf_state = str(run_summary.get("shelf_state") or "unknown")
    next_step = str(run_summary.get("next_step") or "unknown")
    target_quality = float(run_summary.get("target_quality") or 0.98)
    best_quality = float(run_summary.get("best_quality") or 0.0)
    best_cycles = int(run_summary.get("best_cycles") or 0)
    avg_quality = float(run_summary.get("avg_quality") or 0.0)
    branch_markers = [
        f"{item.get('document_type')}: {float(item.get('achieved_quality') or 0.0):.1%} / {int(item.get('cycles_completed') or 0)}c"
        for item in run_summary.get("batch_summaries") or []
        if isinstance(item, dict)
    ]
    branch_text = " | ".join(branch_markers) if branch_markers else "n/a"
    if loop_state in {"BLOCKED", "REVIEW"} and best_cycles <= 0 and best_quality <= 0.0:
        heading = "❌ DOCGEN runtime preflight failed."
    elif loop_state in {"COMPLETE", "TARGET_REACHED"} or best_quality >= target_quality:
        heading = "✅ DOCGEN upgrade completed."
    elif best_cycles > 0:
        heading = "🔵 DOCGEN runtime started."
    else:
        heading = "🟡 DOCGEN upgrade job accepted."
    return (
        f"{heading}\n"
        f"Action: `{action_id}`\n"
        f"Loop state: `{loop_state}`\n"
        f"Shelf: `{shelf_state}`\n"
        f"Best quality: `{best_quality:.1%}` / target `{target_quality:.1%}`\n"
        f"Avg quality: `{avg_quality:.1%}`\n"
        f"Cycles: `{best_cycles}`\n"
        f"Branches: `{branch_text}`\n"
        f"Next: `{next_step}`"
    )



def _is_docsreg_launch_text(text: str) -> bool:
    low = (text or "").lower()
    return any(
        token in low
        for token in (
            "/docsreg",
            "/docsreg_start_media",
            "/gocsreg_start_media",
            "docsreg_start_media",
            "gocsreg_start_media",
            "docsreg ",
            "docsreg:",
        )
    )


async def _is_omi_owned_intent(text: str) -> bool:
    """Return True when the message should be delegated to Omi.

    This is the semantic gate: Axi should not treat Omi-owned document
    workflows as its own generation requests.

    Fast deterministic check runs first so that text-handler messages like
    "Omi, /docsreg <path>" are caught even when the NLP classifier routes
    them to a different intent (e.g. ai_label_removal).
    """
    low = (text or "").strip()
    if not low:
        return False
    low_check = low.lower()
    # If the message is explicitly addressed to Omi it is unconditionally
    # Omi-owned — Axi must not claim it regardless of content.
    if low_check.startswith("omi,") or low_check.startswith("omi "):
        return True
    # Deterministic fast-path: if text contains any docsreg command token,
    # it is unconditionally Omi-owned regardless of NLP result.
    if _is_docsreg_launch_text(low):
        return True
    try:
        from chat_intent_router import OMI_CMDS, classify  # noqa: PLC0415
        routed = await asyncio.to_thread(classify, low, OMI_CMDS)
        return bool(routed and routed[0] == "docsreg")
    except Exception:
        return False


def _recent_ai_label_removal_context(chat_id: int | None, *, max_age_sec: int = 1800) -> bool:
    """Check recent dialog memory for an AI-label-removal request."""
    if chat_id is None:
        return False
    rows = _AXI_DIALOG.get(chat_id) or []
    now = time.time()
    for row in reversed(rows[-10:]):
        if row.get("role") != "user":
            continue
        ts = float(row.get("ts") or 0.0)
        if ts and now - ts > max_age_sec:
            continue
        content = str(row.get("content") or "")
        try:
            from ops.pipelines.ai_label_removal.telegram_handler import detect_ai_label_removal_intent as _ai_rm_intent
            if _ai_rm_intent(content):
                return True
        except Exception:
            if _is_docsreg_launch_text(content):
                return False
    return False


def _should_route_ai_label_removal_text(text: str, chat_id: int | None = None) -> bool:
    """True when the text should seed or continue the AI-label-removal flow."""
    if _wants_docx(text):
        return False
    try:
        from ops.pipelines.ai_label_removal.telegram_handler import detect_ai_label_removal_intent as _ai_rm_intent
        if _ai_rm_intent(text):
            return True
    except Exception:
        pass
    return _recent_ai_label_removal_context(chat_id)


def _should_route_ai_label_removal_upload(message, chat_id: int | None = None) -> bool:
    """
    True when an uploaded document should be processed by the AI-label-removal pipeline.

    This catches three real Telegram shapes:
      - explicit intent in the document caption
      - pending follow-up after a prior text request
      - recent dialog memory when the file arrives as a separate update
    """
    text = getattr(message, "caption", None) or getattr(message, "text", None) or ""
    if _should_route_ai_label_removal_text(text, chat_id):
        return True
    if chat_id is not None and _PENDING_AI_LABEL_REMOVAL.get(chat_id):
        return True
    return False


def _seed_ai_label_removal_pending(chat_id: int | None, text: str) -> None:
    if chat_id is None:
        return
    _PENDING_AI_LABEL_REMOVAL[chat_id] = {"prompt": text, "ts": time.time()}

def _is_owner(update: Update) -> bool:
    if not OWNER_CHATS:
        return False
    if update.effective_chat is None:
        return False
    if update.effective_chat.id in OWNER_CHATS:
        return True
    if update.effective_chat.type in ("group", "supergroup"):
        uid = update.effective_user.id if update.effective_user else None
        if uid is not None and uid in OWNER_CHATS:
            return True
    return False

def _reply_lang(text: str) -> str:
    if AXI_FORCE_REPLY_LANG in ("en", "ru"):
        return AXI_FORCE_REPLY_LANG
    return "ru" if any("\u0400" <= ch <= "\u04ff" for ch in (text or "")) else "en"


def _strip_axi_prefix(text: str) -> str:
    clean = (text or "").strip()
    low = clean.lower()
    for pref in ("axi ", "axi,", "axi:", "акси ", "акси,", "акси:"):
        if low.startswith(pref):
            return clean[len(pref):].strip()
    return clean


# ── Command handlers ───────────────────────────────────────────────────────────

async def cmd_start(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not _chat_allowed(update):
        return
    lang = _reply_lang(update.message.text or "")
    keyboard = InlineKeyboardMarkup([[
        InlineKeyboardButton("🚀 Start working with Axi", callback_data="axi_activate"),
    ]])
    if lang == "ru":
        await update.message.reply_text(
            f"*{AXI_NAME}* — оркестратор AIMS.\n\n"
            "Я обрабатываю внешние задачи: веб-поиск, анализ документов, генерация Word-файлов.\n"
            "Omi присылает результаты напрямую, но задачи идут через Axi.\n\n"
            "Нажмите кнопку или напишите «Axi, …» для начала работы.",
            parse_mode="Markdown",
            reply_markup=keyboard,
        )
    else:
        await update.message.reply_text(
            f"*{AXI_NAME}* — AIMS Orchestrator.\n\n"
            "I handle external tasks: web search, document analysis, Word file generation.\n"
            "Omi can deliver results directly, but all tasks are routed through Axi.\n\n"
            "Press the button or write 'Axi, …' to start.",
            parse_mode="Markdown",
            reply_markup=keyboard,
        )


async def _cmd_start_callback(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle the 'Начать работу с Axi' inline button press."""
    query = update.callback_query
    if query is None:
        return
    await query.answer()
    if not _chat_allowed(update):
        return
    user = update.effective_user
    name = user.first_name if user else "Пользователь"
    await query.message.reply_text(
        f"👋 {name}, Axi активирован.\n"
        "Пишите задачи — я готов к работе.",
    )


async def cmd_help(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not _chat_allowed(update):
        return
    text = (
        f"*{AXI_NAME} — команды*\n\n"
        "/start — приветствие\n"
        "/help — эта справка\n"
        "/aimarks_delete — удалить AI-метки из прикреплённого документа\n"
        "/quality_report [часы] — отчёт качества Task Registry (default: 24ч)\n"
        "/stuck_tasks — список зависших задач прямо сейчас\n"
        "/close_task [task_id] — закрыть зависшую задачу Omi\n"
        "\n"
        "Просто напишите сообщение — Axi ответит или найдёт в интернете.\n"
        "Для генерации Word: 'сделай отчёт в Word' / 'generate report as docx'"
    )
    await update.message.reply_text(text, parse_mode="Markdown")


async def cmd_aimarks_delete(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Explicit alias for the AI-label-removal pipeline."""
    if not _chat_allowed(update):
        return
    message = update.message
    if message is None:
        return

    chat_id = update.effective_chat.id if update.effective_chat else None
    prompt = (message.text or message.caption or "/aimarks_delete").strip()
    _seed_ai_label_removal_pending(chat_id, prompt)

    if getattr(message, "document", None):
        from ops.pipelines.ai_label_removal.telegram_handler import handle_ai_label_removal

        handled = await handle_ai_label_removal(update, ctx)
        if handled:
            return

    await message.reply_text(
        "Прикрепите документ (DOCX, PPTX, XLSX или PDF) к сообщению, и я удалю метки ИИ."
    )


async def cmd_quality_report(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Показать отчёт качества Task Registry."""
    if not _chat_allowed(update):
        return
    hours = 24
    if ctx.args:
        try:
            hours = int(ctx.args[0])
        except ValueError:
            pass
    if _tr_client is None:
        await update.message.reply_text(
            "⚠️ Task Registry недоступен. Убедитесь что task_registry_api.py запущен.",
        )
        return
    try:
        report = _tr_client.quality_report(hours=hours)
        await update.message.reply_text(report or f"Нет данных за {hours}ч.", parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"⚠️ Task Registry ошибка: {e}")


async def cmd_stuck_tasks(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Показать зависшие задачи."""
    if not _chat_allowed(update):
        return
    if _tr_client is None:
        await update.message.reply_text("⚠️ Task Registry недоступен.")
        return
    try:
        tasks = _tr_client.find_stuck(older_than_minutes=TASK_REGISTRY_STUCK_MINUTES)
        if not tasks:
            await update.message.reply_text("✅ Зависших задач нет.")
            return
        lines = [f"⚠️ *Зависшие задачи ({len(tasks)}):*"]
        for t in tasks[:10]:
            lines.append(
                f"  • `{t['task_id']}` [{t['assigned_to']}] "
                f"{int(t['age_minutes'])}м — {t['description'][:60]}"
            )
        await update.message.reply_text("\n".join(lines), parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"⚠️ Ошибка: {e}")


async def cmd_close_task(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Close hanging Omi task in this chat or by explicit task_id."""
    if not _chat_allowed(update):
        return
    chat_id = str(update.effective_chat.id if update.effective_chat else "")
    task_id = (ctx.args[0].strip() if getattr(ctx, "args", None) else "")
    text = task_id or "закрыть задачу"
    handled, msg = _close_omi_task_from_chat(chat_id, text)
    if handled:
        await update.message.reply_text(msg, parse_mode="Markdown")
        return
    await update.message.reply_text("⚠️ Не удалось обработать /close_task.")


async def cmd_cleanup_tasks(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Delete old pending tasks from registry (omi_tasks and bot_cross_handoff)."""
    if not _chat_allowed(update):
        return

    try:
        db_path = Path(AIMS_WORKSPACE) / "aims_registry.db"
        if not db_path.exists():
            await update.message.reply_text("⚠️ Registry database not found.")
            return

        import sqlite3
        conn = sqlite3.connect(str(db_path))
        cursor = conn.cursor()

        # Count old pending tasks (older than 7 days)
        cursor.execute("""
            SELECT COUNT(*) FROM omi_tasks
            WHERE status = 'pending'
            AND datetime(created_at) < datetime('now', '-7 days')
        """)
        omi_count = cursor.fetchone()[0]

        cursor.execute("""
            SELECT COUNT(*) FROM bot_cross_handoff
            WHERE status = 'pending'
            AND datetime(created_at) < datetime('now', '-7 days')
        """)
        handoff_count = cursor.fetchone()[0]

        if omi_count == 0 and handoff_count == 0:
            await update.message.reply_text("✅ Нет старых зависших задач для очистки.")
            conn.close()
            return

        # Delete old pending tasks
        cursor.execute("""
            DELETE FROM omi_tasks
            WHERE status = 'pending'
            AND datetime(created_at) < datetime('now', '-7 days')
        """)

        cursor.execute("""
            DELETE FROM bot_cross_handoff
            WHERE status = 'pending'
            AND datetime(created_at) < datetime('now', '-7 days')
        """)

        conn.commit()
        conn.close()

        await update.message.reply_text(
            f"✅ Очистка завершена:\n"
            f"• Удалено задач из omi_tasks: {omi_count}\n"
            f"• Удалено handoff записей: {handoff_count}\n\n"
            f"Удалены только задачи старше 7 дней в статусе 'pending'.",
            parse_mode="Markdown"
        )

    except Exception as e:
        log.error("cleanup_tasks error: %s", e, exc_info=True)
        await update.message.reply_text(f"⚠️ Ошибка очистки: {e}")


def _set_pending_analyze(
    chat_id: int,
    prompt: str,
    *,
    preserve_existing_files: bool = False,
) -> None:
    prev = _PENDING_ANALYZE.get(chat_id) if preserve_existing_files else None
    prev_files = list(prev.get("files") or []) if isinstance(prev, dict) else []
    prev_last_upload_ts = float(prev.get("last_upload_ts", 0.0)) if isinstance(prev, dict) else 0.0
    prev_lang = str(prev.get("lang", "")).strip() if isinstance(prev, dict) else ""
    _PENDING_ANALYZE[chat_id] = {
        "prompt": prompt,
        "ts": time.time(),
        "last_upload_ts": prev_last_upload_ts,
        "files": prev_files,
        "awaiting_clarify": False,
        "notified_waiting": False,
        "lang": prev_lang or "en",
    }
    _pending_analyze_persist()


def _get_pending_analyze(chat_id: int) -> str | None:
    item = _PENDING_ANALYZE.get(chat_id)
    if not item:
        return None
    ts = float(item.get("ts", 0.0))
    if time.time() - ts > AXI_ANALYZE_WAIT_SEC:
        _PENDING_ANALYZE.pop(chat_id, None)
        _pending_analyze_persist()
        return None
    return str(item.get("prompt", "")).strip() or None


def _get_pending_analyze_state(chat_id: int) -> dict[str, object] | None:
    item = _PENDING_ANALYZE.get(chat_id)
    if not item:
        return None
    ts = float(item.get("ts", 0.0))
    if time.time() - ts > AXI_ANALYZE_WAIT_SEC:
        _PENDING_ANALYZE.pop(chat_id, None)
        _pending_analyze_persist()
        return None
    return item


def _analyze_prompt_is_unclear(prompt: str) -> bool:
    low = (prompt or "").strip().lower()
    if not low:
        return True
    # Special marker: files uploaded but no task description yet
    if "analyze_uploaded_files_waiting_goal" in low:
        return True
    # Training pair workflow: always clear intent, never ask for clarification
    if any(kw in low for kw in ("training pair", "tuning pair", "обучающ пар", "тренировочн пар", "тюнинг пар")):
        return False
    if low in {
        "analyze uploaded document and provide concise findings.",
        "analyze uploaded documents and provide concise findings.",
        "проанализируй загруженные файлы",
    }:
        return True
    informative = ("compare", "summar", "review", "find", "проверь", "сравн", "сводк", "риски", "рекоменд")
    return not any(k in low for k in informative)


def _dialog_append(chat_id: int | None, role: str, content: str) -> None:
    if chat_id is None:
        return
    text = (content or "").strip()
    if not text:
        return
    text = text[:2000]
    bucket = _AXI_DIALOG.setdefault(chat_id, [])
    bucket.append({"role": "user" if role == "user" else "assistant", "content": text})
    _AXI_DIALOG[chat_id] = bucket[-AXI_DIALOG_LOG_MAX:]
    _dialog_persist()


def _claim_update_once(update) -> bool:
    """
    Return True only for the first handler that sees this update within TTL.
    Prevents accidental double-processing for the same Telegram update.
    """
    try:
        uid = getattr(update, "update_id", None)
        key = f"u:{uid}" if uid is not None else ""
        if not key:
            msg = getattr(update, "effective_message", None)
            chat = getattr(update, "effective_chat", None)
            mid = getattr(msg, "message_id", None)
            cid = getattr(chat, "id", None)
            key = f"m:{cid}:{mid}"
        if not key:
            return True
        now = time.time()
        ttl_sec = 120.0
        with _SEEN_UPDATES_LOCK:
            ts = _SEEN_UPDATES.get(key)
            if ts and (now - ts) < ttl_sec:
                return False
            _SEEN_UPDATES[key] = now
            if len(_SEEN_UPDATES) > 4096:
                cutoff = now - ttl_sec
                stale = [k for k, v in _SEEN_UPDATES.items() if v < cutoff]
                for k in stale[:2048]:
                    _SEEN_UPDATES.pop(k, None)
        return True
    except Exception as e:
        log.debug("dedup check error: %s", e)
        return True


def _dialog_context(chat_id: int | None) -> str:
    if chat_id is None:
        return ""
    rows = _AXI_DIALOG.get(chat_id) or []
    if not rows:
        return ""
    lines = ["[Recent chat context | last 10 turns]"]
    for row in rows[-AXI_DIALOG_LOG_MAX:]:
        prefix = "User" if row.get("role") == "user" else "Assistant"
        lines.append(f"{prefix}: {row.get('content', '')}")
    return "\n".join(lines)


def _dialog_load() -> None:
    if not _AXI_DIALOG_PATH.exists():
        return
    try:
        raw = json.loads(_AXI_DIALOG_PATH.read_text(encoding="utf-8"))
        if not isinstance(raw, dict):
            return
        restored: dict[int, list[dict[str, str]]] = {}
        for k, v in raw.items():
            try:
                chat_id = int(k)
            except Exception:
                continue
            if not isinstance(v, list):
                continue
            rows: list[dict[str, str]] = []
            for item in v[-AXI_DIALOG_LOG_MAX:]:
                if not isinstance(item, dict):
                    continue
                role = "user" if item.get("role") == "user" else "assistant"
                content = str(item.get("content", "")).strip()
                if content:
                    rows.append({"role": role, "content": content[:2000]})
            if rows:
                restored[chat_id] = rows[-AXI_DIALOG_LOG_MAX:]
        _AXI_DIALOG.clear()
        _AXI_DIALOG.update(restored)
    except Exception as e:
        log.warning("dialog load failed: %s", e)


def _dialog_persist() -> None:
    with _AXI_DIALOG_LOCK:
        try:
            payload = {str(k): v[-AXI_DIALOG_LOG_MAX:] for k, v in _AXI_DIALOG.items()}
            _AXI_DIALOG_PATH.parent.mkdir(parents=True, exist_ok=True)
            _AXI_DIALOG_PATH.write_text(
                json.dumps(payload, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as e:
            log.debug("dialog persist failed: %s", e)


def _pending_analyze_persist() -> None:
    with _AXI_PENDING_LOCK:
        try:
            payload = {str(k): v for k, v in _PENDING_ANALYZE.items()}
            _AXI_PENDING_ANALYZE_PATH.parent.mkdir(parents=True, exist_ok=True)
            _AXI_PENDING_ANALYZE_PATH.write_text(
                json.dumps(payload, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as e:
            log.debug("pending analyze persist failed: %s", e)


def _pending_analyze_load() -> None:
    if not _AXI_PENDING_ANALYZE_PATH.exists():
        return
    try:
        raw = json.loads(_AXI_PENDING_ANALYZE_PATH.read_text(encoding="utf-8"))
        if not isinstance(raw, dict):
            return
        restored: dict[int, dict[str, object]] = {}
        now = time.time()
        for k, v in raw.items():
            try:
                chat_id = int(k)
            except Exception:
                continue
            if not isinstance(v, dict):
                continue
            ts = float(v.get("ts", 0.0))
            if ts and (now - ts) > AXI_ANALYZE_WAIT_SEC:
                continue
            prompt = str(v.get("prompt", "")).strip()
            if not prompt:
                continue
            files = v.get("files") if isinstance(v.get("files"), list) else []
            restored[chat_id] = {
                "prompt": prompt,
                "ts": ts or now,
                "last_upload_ts": float(v.get("last_upload_ts", 0.0)),
                "files": files,
                "awaiting_clarify": bool(v.get("awaiting_clarify", False)),
                "notified_waiting": bool(v.get("notified_waiting", False)),
                "lang": str(v.get("lang", "en") or "en"),
                "file_ack_sent": bool(v.get("file_ack_sent", False)),
            }
        _PENDING_ANALYZE.clear()
        _PENDING_ANALYZE.update(restored)
    except Exception as e:
        log.warning("pending analyze load failed: %s", e)


def _looks_like_file_dependent_request(text: str) -> bool:
    low = (text or "").strip().lower()
    if not low:
        return False
    # Do not route "send the generated document" into upload/analyze flow.
    resend_markers = (
        "send document",
        "send the document",
        "send file",
        "resent document",
        "resend document",
        "what you generated",
        "generated as attachment",
        "as attachment to chat",
        "пришли документ",
        "отправь документ",
        "повтори отправку",
    )
    if any(m in low for m in resend_markers):
        return False
    if "/analyze" in low:
        return True
    markers = (
        "attached",
        "attach",
        "attachment",
        "file",
        "review the attached",
        "прикреп",
        "вложен",
        "файл",
        "документ",
    )
    return any(m in low for m in markers)


def _parse_visual_and_narration(text: str) -> tuple[str, str | None]:
    """
    Отделяет описание кадра от опционального текста озвучки/декламации.
    Маркеры: озвучка / декламация / текст речи / реплика / голос / narration / voiceover / voice
    """
    raw = (text or "").strip()
    if not raw:
        return "", None
    m = re.search(
        r"(?is)[\n\r]?\s*(?:озвучка|декламация|текст\s+речи|реплика|голос|narration|voiceover|voice)\s*[:：]\s*",
        raw,
    )
    if not m:
        return raw, None
    visual = raw[: m.start()].strip().rstrip("-:—").strip()
    narration = raw[m.end() :].strip()
    if not narration:
        return raw, None
    return visual, narration[:900]


def _build_veo_user_prompt(visual: str, narration: str | None) -> str:
    v = (visual or "").strip()
    n = (narration or "").strip()
    if not v:
        return ""
    if not n:
        return v[:1000]
    return (
        f"{v}\n\n"
        f"[Текст для озвучки или декламации — встроить, если модель поддерживает речь или субтитры]: {n}"
    )[:1000]


def _unlink_axi_temp(p: str | None) -> None:
    if not p:
        return
    try:
        Path(p).unlink(missing_ok=True)
    except OSError:
        pass


def _telegram_doc_is_reference_video(doc) -> bool:
    if doc is None:
        return False
    mt = (doc.mime_type or "").lower()
    if mt.startswith("video/"):
        return True
    suf = Path(doc.file_name or "").suffix.lower()
    return suf in {".mp4", ".mov", ".webm", ".mpeg", ".mpg", ".avi", ".wmv", ".flv"}


def _finalize_video_prompt_from_user_text(text: str) -> tuple[str, bool]:
    """
    Собирает итоговый промпт для Veo. Второй элемент True — нужно ещё описание сцены
    (например указали только блок «озвучка: …» без визуала).
    """
    raw = (text or "").strip()
    if len(raw) < 4:
        return "", True
    visual, narr = _parse_visual_and_narration(raw)
    if visual.strip():
        return _build_veo_user_prompt(visual, narr), False
    if narr and not visual.strip():
        return "", True
    return _build_veo_user_prompt(raw, None), False


def _match_video_trigger_span(low: str) -> tuple[int, int] | None:
    """
    «сгенерируй 2 видио» не содержит подстроку «сгенерируй видио» — нужен regex.
    """
    patterns = (
        r"сгенерируй\s+\d+\s+видио",
        r"сгенерируй\s+\d+\s+видео",
        r"сгенерируй\s+два\s+видио",
        r"сгенерируй\s+два\s+видео",
        r"сгенерируй\s+видио",
        r"сгенерируй\s+видео",
        r"generate\s+\d+\s+videos?",
        r"generate\s+videos?",
    )
    for p in patterns:
        m = re.search(p, low, flags=re.IGNORECASE)
        if m:
            return m.start(), m.end()
    return None


def _split_veo_first_second_blocks(text: str) -> list[str] | None:
    """
    Два промпта в одном сообщении: «Первое … Второе …» или First / Second.
    """
    t = (text or "").strip()
    if not t:
        return None
    for pat in (
        r"(?is)^(.*?)\bпервое\b\s*[:\n]?\s*(.*?)\s*\bвторое\b\s*[:\n]?\s*(.*)\Z",
        r"(?is)^(.*?)\bfirst\b\s*[:\n]?\s*(.*?)\s*\bsecond\b\s*[:\n]?\s*(.*)\Z",
    ):
        m = re.search(pat, t)
        if not m:
            continue
        header = m.group(1).strip()
        body1 = m.group(2).strip()
        body2 = m.group(3).strip()
        if not body1 or not body2:
            continue
        body1 = body1.replace("\u200b", "").replace("\ufeff", "").strip()
        body2 = body2.replace("\u200b", "").replace("\ufeff", "").strip()
        if not body1 or not body2:
            continue
        head = (header + "\n\n") if header else ""
        p1 = f"{head}[Видео 1/2, ~8с, business motion graphics] {body1}".strip()
        p2 = f"{head}[Видео 2/2, ~8с, business motion graphics] {body2}".strip()
        return [p1[:4000], p2[:4000]]
    return None


def _extract_requested_clip_seconds(text: str) -> int | None:
    """
    «по 10 с», «duration: 6» в начале запроса — для подсказки пользователю и clamp в Veo.
    """
    t = (text or "").strip()
    if not t:
        return None
    m = re.search(
        r"(?i)по\s+(\d{1,2})\s*(?:с(?:ек(?:унд(?:ы)?)?)?|sec(?:onds?)?)",
        t,
    )
    if not m:
        m = re.search(
            r"(?is)[\n\r]?\s*(?:длительность|duration)\s*[:：]\s*(\d{1,2})\s*(?:с(?:ек)?|sec|seconds?)?",
            t[:900],
        )
    if not m:
        return None
    try:
        n = int(m.group(1))
        return n if 1 <= n <= 99 else None
    except ValueError:
        return None


def _extract_video_skill_request(text: str) -> tuple[str, str | None, str | None, bool, list[str], int | None] | None:
    original = (text or "").strip()
    low = original.lower()
    if not low:
        return None
    span = _match_video_trigger_span(low)
    if span is None:
        return None
    _start, end = span
    tail = original[end:].strip()
    model_name: str | None = None
    mode_name: str | None = None

    # Optional inline override:
    m = re.search(
        r"(?:на\s+модели|model)\s+([A-Za-z0-9._\-]+)\s*[:,-]?\s*(.*)$",
        tail,
        flags=re.IGNORECASE,
    )
    if m:
        model_name = m.group(1).strip()
        tail = m.group(2).strip()
    else:
        tail = tail.strip(" :,-")

    low_tail = tail.lower()
    if re.search(r"\b(режим\s+качество|quality mode|mode quality|hq mode)\b", low_tail):
        mode_name = "quality"
        tail = re.sub(r"(режим\s+качество|quality mode|mode quality|hq mode)", "", tail, flags=re.IGNORECASE).strip(" :,-")
    elif re.search(r"\b(режим\s+быстро|быстрый\s+режим|fast mode|quick mode|mode fast)\b", low_tail):
        mode_name = "fast"
        tail = re.sub(r"(режим\s+быстро|быстрый\s+режим|fast mode|quick mode|mode fast)", "", tail, flags=re.IGNORECASE).strip(" :,-")

    desc = tail.strip()
    req_sec = _extract_requested_clip_seconds(desc)
    duo = _split_veo_first_second_blocks(desc)
    if duo and len(duo) == 2:
        return duo[0][:1000], model_name, mode_name, False, [duo[1][:1000]], req_sec

    visual, narr = _parse_visual_and_narration(desc)
    if visual.strip():
        prompt = _build_veo_user_prompt(visual, narr)
        needs_description = False
    elif narr and not visual.strip():
        prompt = ""
        needs_description = True
    else:
        needs_description = not bool(desc)
        prompt = desc
    return prompt, model_name, mode_name, needs_description, [], req_sec


def _extract_video_skill_alias_request(text: str) -> tuple[str, str | None, str | None, bool, list[str], int | None] | None:
    """
    Альтернативные формулировки без «сгенерируй видео» — всё равно включают Veo-скилл.
    """
    low = (text or "").strip().lower()
    if not low:
        return None
    aliases = (
        "используй скилл",
        "используй skill",
        "use video skill",
        "use the video skill",
        "скилл видео",
        "video skill",
        "veo skill",
        "режим видео veo",
        "генерация видео скилл",
    )
    if not any(a in low for a in aliases):
        return None
    tail = (text or "").strip()
    for a in aliases:
        idx = low.find(a)
        if idx >= 0:
            tail = tail[idx + len(a) :].strip(" :,-")
            break
    desc = tail.strip()
    req_sec = _extract_requested_clip_seconds(desc)
    visual, narr = _parse_visual_and_narration(desc)
    if visual.strip():
        prompt = _build_veo_user_prompt(visual, narr)
        needs_description = False
    elif narr and not visual.strip():
        prompt = ""
        needs_description = True
    else:
        needs_description = not bool(desc)
        prompt = desc
    return prompt[:1000], None, None, needs_description, [], req_sec


def _is_video_pending_clarification_only(text: str) -> bool:
    """
    Короткие реплики при уже включённом ожидании картинки для видео.
    Не должны уходить в общий LLM (иначе «я не могу генерировать видео»).
    """
    low = (text or "").strip().lower()
    if not low:
        return False
    if len(low) > 120:
        return False
    markers = (
        "возьми эти",
        "возьми это",
        "вот они",
        "вот она",
        "вот оно",
        "сюда",
        "это оно",
        "это они",
        "фото выше",
        "картинк",
        "используй скилл",
        "используй skill",
        "продолж",
        "use the skill",
        "use video skill",
        "here you",
        "here they",
    )
    return any(m in low for m in markers)


def _is_casual_non_task_message(text: str) -> bool:
    low = (text or "").strip().lower()
    if not low:
        return True
    greetings = {
        "hi", "hello", "hey", "yo", "sup",
        "привет", "здравствуй", "здравствуйте", "добрый день", "добрый вечер",
    }
    acknowledgements = {
        "ok", "okay", "thanks", "thank you", "thx",
        "ок", "хорошо", "спасибо", "понял", "принял",
    }
    if low in greetings or low in acknowledgements:
        return True
    # Short phatic/opening messages should not create Task Registry records.
    words = [w for w in re.split(r"\s+", low) if w]
    if len(words) <= 3 and any(w in greetings for w in words):
        return True
    return False


async def _process_analyze_batch(
    bot,
    chat_id: int,
    prompt: str,
    files: list[dict[str, str]],
    *,
    lang: str = "ru",
) -> None:
    t0 = time.perf_counter()
    progress_msg = await bot.send_message(
        chat_id=chat_id,
        text=(
            f"⏱ {AXI_NAME}: взято в работу. Анализирую пакет файлов…"
            if lang == "ru" else
            f"⏱ {AXI_NAME}: taken in progress. Analyzing file batch…"
        ),
    )
    excerpts: list[str] = []
    for f in files:
        name = f.get("name", "file")
        txt = (f.get("excerpt", "") or "").strip()
        if txt:
            excerpts.append(f"[file: {name}]\n{txt}")
    merged = "\n\n---\n\n".join(excerpts)[:100000]
    request_text = (
        f"{prompt}\n\n"
        "[Analyze all uploaded files together as one package. "
        "If standards are requested, compare and consolidate.]"
        f"\n\n{merged}"
    )

    tr_id = _tr_register(
        f"batch analyze: {prompt[:80]}",
        chat_id=str(chat_id),
        source="axi",
    )
    _tr_start(tr_id, assigned_to="axi")
    if tr_id:
        await bot.send_message(
            chat_id=chat_id,
            text=(
                f"🧾 Задача зарегистрирована: `{tr_id}`.\nПакет файлов принят, начинаю общий анализ."
                if lang == "ru" else
                f"🧾 Task registered: `{tr_id}`.\nFile batch accepted, starting consolidated analysis."
            ),
            parse_mode="Markdown",
        )
    _anim_stop = asyncio.Event()
    _anim_base = (
        f"{AXI_NAME}: анализирую пакет…" if lang == "ru"
        else f"{AXI_NAME}: analysing file batch…"
    )
    _anim_task = asyncio.create_task(_animate_progress(progress_msg, _anim_base, _anim_stop))
    try:
        if await _is_omi_owned_intent(prompt):
            log.info("Axi delegating DOCSREG intent to Omi: chat_id=%s", chat_id)
            _anim_stop.set()
            _tr_done(tr_id, summary="delegated_to_omi:docsreg")
            return

        use_search = _should_web_search(prompt)

        # ── AI-label removal: explicit request runs the dedicated pipeline ────────
        try:
            from ops.pipelines.ai_label_removal.telegram_handler import (
                detect_ai_label_removal_intent as _ai_rm_intent,
                run_request as _ai_rm_run,
                build_user_message as _ai_rm_msg,
            )
            if _ai_rm_intent(prompt):
                _sup = next(
                    (f for f in files if (f.get("name") or "").lower().endswith(
                        (".docx", ".pptx", ".xlsx", ".pdf"))),
                    None,
                )
                if _sup and _sup.get("path"):
                    _anim_stop.set()
                    _res = await asyncio.to_thread(
                        _ai_rm_run,
                        Path(_sup["path"]),
                        AXI_RESULTS_DIR / "ai_label_removal",
                        chat_id=chat_id,
                        user_id=None,
                        original_filename=_sup.get("name") or "document",
                    )
                    if _res.status == "SUCCESS" and _res.output_path:
                        with open(_res.output_path, "rb") as _fh:
                            await bot.send_document(
                                chat_id=chat_id, document=_fh,
                                filename=_res.output_path.name, caption=_ai_rm_msg(_res),
                            )
                    else:
                        await bot.send_message(chat_id=chat_id, text=_ai_rm_msg(_res))
                    _tr_done(tr_id, summary=f"ai_label_removal:{_res.status}")
                    return
        except Exception as _ai_rm_err:
            log.debug("ai_label_removal intercept skipped: %s", _ai_rm_err)

        # ── Intent classification: LLM first, keywords as fast-bypass only ────────
        _xlsx_in_batch = any((f.get("name") or "").lower().endswith((".xlsx", ".xlsm")) for f in files)
        _explicit_docx = bool(re.search(r"\b(word|docx|\.docx)\b", prompt.lower()))
        _explicit_edit = _xlsx_in_batch and _wants_xlsx_edit(prompt, files) and not any(
            kw in prompt.lower() for kw in ("check", "compare", "analyze", "analyse", "review", "assess", "report", "send", "share")
        )
        if _explicit_docx:
            file_intent = "docx"
        elif _explicit_edit:
            file_intent = "edit"
        else:
            file_intent = await _infer_file_intent(prompt, files)
            _save_intent_example(prompt, files, file_intent)
        # ──────────────────────────────────────────────────────────────────────────
        if file_intent == "edit":
            xlsx_entry = next(
                (f for f in files if (f.get("name") or "").lower().endswith((".xlsx", ".xlsm"))),
                None,
            )
            xlsx_src = xlsx_entry and xlsx_entry.get("path")
            if xlsx_src and Path(xlsx_src).exists():
                result_xlsx = await _edit_xlsx_with_llm(Path(xlsx_src), prompt)
                dt = time.perf_counter() - t0
                try:
                    await progress_msg.edit_text(
                        f"✅ {AXI_NAME}: таблица обновлена за {dt:.0f}с. Отправляю…"
                        if lang == "ru" else
                        f"✅ {AXI_NAME}: table updated in {dt:.0f}s. Sending…"
                    )
                except Exception:
                    pass
                caption_text = (
                    f"{AXI_NAME}: ✅ {result_xlsx.name}"
                    + (f" | task `{tr_id}`" if tr_id else "")
                )
                for _attempt in range(3):
                    try:
                        with result_xlsx.open("rb") as fh:
                            await bot.send_document(
                                chat_id=chat_id,
                                document=fh,
                                filename=result_xlsx.name,
                                caption=caption_text,
                            )
                        break
                    except Exception as _send_err:
                        if _attempt < 2:
                            await asyncio.sleep(3)
                        else:
                            raise _send_err
                _anim_stop.set()
                _tr_done(tr_id, summary=f"xlsx_edit:{result_xlsx.name}")
                return
        if file_intent == "docx":
            _has_xlsx = _xlsx_in_batch
            _doc_sys = _build_doc_analysis_system() if _has_xlsx else None
            answer = await _llm_reply(request_text, use_search=use_search, system_override=_doc_sys)
            answer = _strip_generate_file_json(answer)
            first_line = next((ln.strip().lstrip("#").strip() for ln in answer.splitlines() if ln.strip()), "axi_document")
            docx_path = await _generate_custom_docx(answer, "axi_analyze_batch", AXI_RESULTS_DIR, title=first_line)
            dt = time.perf_counter() - t0
            _anim_stop.set()
            try:
                await progress_msg.edit_text(
                    f"✅ {AXI_NAME}: задача выполнена за {dt:.0f}с. Отправляю сгенерированный файл…"
                    if lang == "ru" else
                    f"✅ {AXI_NAME}: completed in {dt:.0f}s. Sending generated file…"
                )
            except Exception:
                pass
            caption_text = (
                f"{AXI_NAME}: ✅ {docx_path.name}\n"
                + _sources_note(
                    lang=lang,
                    use_search=use_search,
                    uploaded_files_count=len(files),
                )
                + (f" | task `{tr_id}`" if tr_id else "")
            )
            for _attempt in range(3):
                try:
                    with docx_path.open("rb") as fh:
                        await bot.send_document(
                            chat_id=chat_id,
                            document=fh,
                            filename=docx_path.name,
                            caption=caption_text,
                            parse_mode="Markdown",
                        )
                    break
                except Exception as _send_err:
                    if _attempt < 2:
                        await asyncio.sleep(3)
                    else:
                        raise _send_err
            _tr_done(tr_id, summary=f"batch_docx:{docx_path.name}:{len(files)}")
        elif file_intent == "ask":
            _anim_stop.set()
            _pending_intent[chat_id] = {"prompt": prompt, "files": files, "tr_id": tr_id, "lang": lang, "ts": time.time()}
            # Ambiguous intent → show the 5-option clarification menu:
            # DOCGEN / training pair / standard revision / DOCSREG / AI-label cleanup.
            _menu_shown = False
            try:
                import shutil as _shutil
                from ops.pipelines.ai_label_removal import telegram_intent_menu as _ai_menu
                _sup = next(
                    (f for f in files if (f.get("name") or "").lower().endswith(
                        (".docx", ".pptx", ".xlsx", ".pdf"))),
                    files[0] if files else None,
                )
                if _sup and _sup.get("path") and Path(_sup["path"]).exists():
                    # Copy to a persistent location so the menu callback can still read it.
                    _persist_dir = AXI_RESULTS_DIR / "ai_intent_pending"
                    _persist_dir.mkdir(parents=True, exist_ok=True)
                    _dst = _persist_dir / f"{int(time.time())}_{_sup.get('name') or 'document'}"
                    _shutil.copy2(_sup["path"], _dst)
                    _req = _ai_menu.stash_pending(
                        chat_id=chat_id, user_id=None, file_path=_dst,
                        original_filename=_sup.get("name") or "document", prompt=prompt,
                    )
                    await bot.send_message(
                        chat_id=chat_id, text=_ai_menu.CLARIFICATION_PROMPT,
                        reply_markup=_ai_menu.build_inline_keyboard(_req.token),
                    )
                    _menu_shown = True
                    _tr_done(tr_id, summary="batch_ask_menu")
            except Exception as _menu_err:
                log.debug("intent menu failed, falling back to text clarify: %s", _menu_err)

            if not _menu_shown:
                clarify_text = (
                    f"{AXI_NAME}: получил файлы, но не до конца понял задачу. Уточни:\n"
                    f"— нужен **Word-документ** с анализом/отчётом?\n"
                    f"— или **изменить данные** внутри таблицы?\n"
                    f"— или просто **текстовый ответ**?"
                    if lang == "ru" else
                    f"{AXI_NAME}: received the files but the task is unclear. Please clarify:\n"
                    f"— do you need a **Word document** (analysis / report)?\n"
                    f"— or **edit data** inside the spreadsheet?\n"
                    f"— or just a **text answer**?"
                )
                await bot.send_message(chat_id=chat_id, text=clarify_text, parse_mode="Markdown")
                _tr_done(tr_id, summary="batch_ask_clarify")
        elif file_intent == "training_pair":
            # Training pair creation pipeline
            _anim_stop.set()
            try:
                await progress_msg.edit_text(
                    f"📚 {AXI_NAME}: начинаю подготовку обучающей пары..."
                    if lang == "ru" else
                    f"📚 {AXI_NAME}: starting training pair preparation..."
                )
            except Exception:
                pass

            # Step 1: Extract context from uploaded documents
            await bot.send_message(
                chat_id=chat_id,
                text=(
                    "📄 Шаг 1/3: Извлекаю контекст из загруженных документов..."
                    if lang == "ru" else
                    "📄 Step 1/3: Extracting context from uploaded documents..."
                )
            )

            # Read document texts
            doc_texts = []
            for file_entry in files:
                try:
                    # files is list of dict with 'name' and 'path' keys
                    file_path = Path(file_entry.get("path", "")) if isinstance(file_entry, dict) else file_entry
                    if not file_path.exists():
                        continue
                    doc_text = _read_document_text(file_path)
                    file_name = file_entry.get("name", file_path.name) if isinstance(file_entry, dict) else file_path.name
                    doc_texts.append({"name": file_name, "text": doc_text[:50000], "path": str(file_path)})
                except Exception as e:
                    file_name = file_entry.get("name", "unknown") if isinstance(file_entry, dict) else str(file_entry)
                    log.warning("Failed to read %s: %s", file_name, e)

            if not doc_texts:
                await bot.send_message(
                    chat_id=chat_id,
                    text=(
                        "⚠️ Не удалось прочитать документы"
                        if lang == "ru" else
                        "⚠️ Failed to read documents"
                    )
                )
                _tr_done(tr_id, summary="training_pair_read_failed")
                return

            # Check for duplicate training pair by content hash
            import hashlib
            combined_text = "\n\n".join(d["text"] for d in doc_texts)
            content_hash = hashlib.sha256(combined_text.encode("utf-8")).hexdigest()[:16]

            existing_memo = _find_existing_training_pair(content_hash)
            if existing_memo:
                await _safe_send_message(
                    bot,
                    chat_id,
                    (
                        f"ℹ️ Обучающая пара для этих файлов уже создана ранее:\n"
                        f"Файл: {_tg_plain(existing_memo['memo_file'])}\n"
                        f"Дата: {_tg_plain(existing_memo['created'])}\n"
                        f"Оценка локальной модели: {_tg_plain(existing_memo.get('local_score', 'не оценено'))}\n\n"
                        f"Результат можно посмотреть в мастер-файле."
                        if lang == "ru" else
                        f"ℹ️ Training pair for these files already exists:\n"
                        f"File: {_tg_plain(existing_memo['memo_file'])}\n"
                        f"Date: {_tg_plain(existing_memo['created'])}\n"
                        f"Local model score: {_tg_plain(existing_memo.get('local_score', 'not evaluated'))}\n\n"
                        f"Result available in master file."
                    ),
                )
                _tr_done(tr_id, summary=f"training_pair_duplicate:{existing_memo['memo_file']}")
                return

            # Step 2: Detect template/reference roles from filenames, then extract context
            blank_doc, filled_doc, _role_reason = _detect_template_reference_roles(doc_texts)
            blank_text = blank_doc.get("text", "") if blank_doc else ""
            filled_text = filled_doc.get("text", blank_text) if filled_doc else blank_text
            blank_name_detected = blank_doc.get("name", doc_texts[0]["name"] if doc_texts else "blank") if blank_doc else ""
            log.info("doctuning role detection: %s", _role_reason)
            context_result = await asyncio.to_thread(_extract_doc_context_sync, blank_text, filled_text)

            if not context_result or "error" in context_result:
                err_code = context_result.get("error", "unknown") if context_result else "unknown"
                if err_code == "omniroute_session_limit":
                    user_msg = "External model session limit is still active after cleanup. Please retry in a few minutes."
                elif err_code == "omniroute_rate_limited":
                    user_msg = "External model is rate-limited. Please wait a few minutes and retry."
                else:
                    user_msg = f"Context extraction error: {err_code}"
                await bot.send_message(chat_id=chat_id, text=f"⚠️ {user_msg}")
                _tr_done(tr_id, summary=f"training_pair_context_failed:{err_code}")
                return

            # Track file format for xlsx-aware output
            _blank_suffix = Path(blank_doc.get("name", "")).suffix.lower() if blank_doc else ""
            _filled_suffix = Path(filled_doc.get("name", "")).suffix.lower() if filled_doc else ""
            _preferred_fmt = "xlsx" if _blank_suffix == ".xlsx" or _filled_suffix == ".xlsx" else "docx"

            _PENDING_DOCFILL[chat_id] = {
                "step": "await_standards_approval",
                "source": "batch_training_pair",
                "context_result": context_result,
                "blank_name": blank_name_detected,
                "filled_name": filled_doc.get("name", "") if filled_doc else "",
                "blank_path": blank_doc.get("path", "") if blank_doc else "",
                "filled_path": filled_doc.get("path", "") if filled_doc else "",
                "content_hash": content_hash,
                "blank_suffix": _blank_suffix,
                "filled_suffix": _filled_suffix,
                "preferred_output_format": _preferred_fmt,
                "blank_text": blank_text,
                "filled_text": filled_text,
            }

            standards_list = context_result.get("standards", [])
            standards_str = ", ".join(standards_list) if standards_list else "—"
            await _safe_send_message(
                bot,
                chat_id,
                (
                    f"✅ Context extracted:\n"
                    f"Document type: {_tg_plain(context_result.get('doc_type', 'unknown'))}\n"
                    f"Standards: {_tg_plain(standards_str)}\n"
                    f"Key terms: {len(context_result.get('key_terms', []))}\n\n"
                    "Reply yes or continue to generate the revised reference candidate, "
                    "or cancel to discard."
                ),
            )
            _tr_done(tr_id, summary="training_pair_awaiting_confirmation")
        else:
            answer = await _llm_reply(request_text, use_search=use_search)
            dt = time.perf_counter() - t0
            _anim_stop.set()
            try:
                await progress_msg.edit_text(
                    f"✅ {AXI_NAME}: задача выполнена за {dt:.0f}с."
                    if lang == "ru" else
                    f"✅ {AXI_NAME}: completed in {dt:.0f}s."
                )
            except Exception:
                pass
            out_text = (
                answer[:3600]
                + "\n\n"
                + _sources_note(
                    lang=lang,
                    use_search=use_search,
                    uploaded_files_count=len(files),
                )
            )
            await bot.send_message(chat_id=chat_id, text=out_text[:4000])
            _tr_done(tr_id, summary=f"batch_chat:{len(files)}")
    except Exception as e:
        _anim_stop.set()
        _tr_stuck(tr_id, error=str(e)[:200])
        _emsg = str(e)[:120] if str(e) else type(e).__name__
        try:
            await progress_msg.edit_text(
                f"⚠️ {AXI_NAME}: ошибка обработки пакета — {_emsg}"
                if lang == "ru" else
                f"⚠️ {AXI_NAME}: batch processing error — {_emsg}"
            )
        except Exception:
            pass
        await bot.send_message(
            chat_id=chat_id,
            text=(
                f"{AXI_NAME}: не удалось обработать пакет файлов — {_emsg}"
                if lang == "ru" else
                f"{AXI_NAME}: failed to process file batch — {_emsg}"
            ),
        )
    finally:
        _anim_stop.set()
        _anim_task.cancel()
        try:
            await _anim_task
        except asyncio.CancelledError:
            pass


async def _job_flush_pending_analyze(context: ContextTypes.DEFAULT_TYPE) -> None:
    now = time.time()
    to_process: list[tuple[int, dict[str, object]]] = []
    for chat_id, state in list(_PENDING_ANALYZE.items()):
        if now - float(state.get("ts", 0.0)) > AXI_ANALYZE_WAIT_SEC:
            _PENDING_ANALYZE.pop(chat_id, None)
            _pending_analyze_persist()
            continue
        files = state.get("files") or []
        if not files or state.get("awaiting_clarify"):
            continue
        last_upload = float(state.get("last_upload_ts", 0.0))
        if last_upload and (now - last_upload) >= AXI_ANALYZE_BATCH_WAIT_SEC:
            to_process.append((chat_id, state))
    for chat_id, state in to_process:
        prompt = str(state.get("prompt", "")).strip()
        files = list(state.get("files") or [])
        lang = str(state.get("lang", "ru"))
        if _analyze_prompt_is_unclear(prompt):
            await context.application.bot.send_message(
                chat_id=chat_id,
                text=(
                    "Файлы получены. Уточните, что сделать с пакетом: суммаризация, сравнение, проверка соответствия, риски, или подготовить DOCX?"
                ),
            )
            state["awaiting_clarify"] = True
            _pending_analyze_persist()
            continue
        _PENDING_ANALYZE.pop(chat_id, None)
        _pending_analyze_persist()
        await _process_analyze_batch(context.application.bot, chat_id, prompt, files, lang=lang)


async def _job_axi_deliver_cross_handoffs(context: ContextTypes.DEFAULT_TYPE) -> None:
    """Deliver Omi->Axi handoffs as orchestration status messages in chat."""
    try:
        from cross_bot_handoff import (
            claim_pending_for_target,
            handoff_delivery_enabled,
            mark_delivered,
            mark_failed,
        )
    except Exception as e:
        log.debug("axi handoff import failed: %s", e)
        return
    if not handoff_delivery_enabled():
        return
    try:
        rows = claim_pending_for_target("axi", limit=8)
    except Exception as e:
        log.debug("axi handoff claim failed: %s", e)
        return
    for row in rows:
        hid = int(row.get("id", 0) or 0)
        chat_id = int(row.get("chat_id", 0) or 0)
        payload = str(row.get("payload", "") or "").strip()
        if not hid or not chat_id:
            continue
        try:
            await context.bot.send_message(
                chat_id=chat_id,
                text=payload or "Axi: получен handoff от Omi.",
                parse_mode="Markdown",
            )
            mark_delivered(hid)
        except Exception as e:
            mark_failed(hid, str(e))


def _extract_text_from_uploaded_document(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            parts: list[str] = []
            for page in reader.pages[:30]:
                parts.append((page.extract_text() or "").strip())
            return "\n\n".join([p for p in parts if p]).strip()[:30000]
        if ext == ".docx":
            from docx import Document

            doc = Document(str(path))
            return "\n".join((p.text or "").strip() for p in doc.paragraphs if (p.text or "").strip())[:30000]
        if ext in (".txt", ".md", ".csv", ".log"):
            return path.read_text(errors="replace")[:30000]
        if ext in (".xlsx", ".xlsm"):
            from openpyxl import load_workbook

            wb = load_workbook(str(path), data_only=True, read_only=True)
            lines: list[str] = []
            for ws in wb.worksheets[:3]:
                lines.append(f"[sheet: {ws.title}]")
                for row in ws.iter_rows(min_row=1, max_row=80, max_col=12, values_only=True):
                    vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if vals:
                        lines.append(" | ".join(vals))
            return "\n".join(lines)[:30000]
        if ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(path))
            lines: list[str] = []
            for sidx, slide in enumerate(prs.slides[:20], start=1):
                lines.append(f"[slide {sidx}]")
                for shape in slide.shapes:
                    txt = getattr(shape, "text", "") or ""
                    txt = txt.strip()
                    if txt:
                        lines.append(txt)
            return "\n".join(lines)[:30000]
    except Exception as e:
        log.warning("extract uploaded text failed (%s): %s", path.name, e)
    return ""


def _queue_to_omi_batch_inbox(
    src_path: Path,
    original_name: str,
    *,
    meta: dict[str, object] | None = None,
) -> Path:
    """
    Put file into Omi OCR/batch pipeline inbox.
    Supports cross-device moves (/tmp → /data) via shutil.copy2 fallback.
    If same content already exists, returns existing path and removes src.
    """
    import errno
    import hashlib
    import shutil

    inbox = AIMS_WORKSPACE / "batch_inbox"
    inbox.mkdir(parents=True, exist_ok=True)
    incoming_sha = hashlib.sha256(src_path.read_bytes()).hexdigest()
    dest = inbox / original_name
    if dest.exists():
        if hashlib.sha256(dest.read_bytes()).hexdigest() == incoming_sha:
            src_path.unlink(missing_ok=True)
            return dest
        stem, suffix = Path(original_name).stem, Path(original_name).suffix
        for n in range(2, 100):
            cand = inbox / f"{stem}_v{n}{suffix}"
            if not cand.exists():
                dest = cand
                break
    try:
        src_path.replace(dest)
        log.info("queue_to_omi_batch_inbox: queued %s -> %s", src_path, dest)
    except OSError as e:
        if e.errno == errno.EXDEV:
            shutil.copy2(str(src_path), str(dest))
            if not dest.exists() or dest.stat().st_size == 0:
                raise RuntimeError(
                    f"queue_to_omi_batch_inbox: copy succeeded but dest missing/empty: {dest}"
                )
            src_path.unlink(missing_ok=True)
            log.warning("queue_to_omi_batch_inbox: cross-device move; copied %s -> %s", src_path, dest)
        else:
            raise
    if meta:
        meta_path = dest.with_name(dest.name + ".axi.json")
        try:
            meta_path.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
        except Exception as e:
            log.warning("queue_to_omi meta write failed %s: %s", dest.name, e)
    return dest


def _load_recent_generated_doc_context(chat_id: int | None) -> tuple[str, str]:
    """Load the most recent generated docx content from dialog memory."""
    if chat_id is None:
        return "", ""
    rows = _AXI_DIALOG.get(chat_id) or []
    for row in reversed(rows):
        content = str((row or {}).get("content") or "").strip()
        if not (content.startswith("[doc:") and content.endswith("]")):
            continue
        file_name = content[5:-1].strip()
        if not file_name:
            continue
        p = AXI_RESULTS_DIR / file_name
        if not p.exists() or p.suffix.lower() != ".docx":
            continue
        try:
            from docx import Document as _Doc
            d = _Doc(str(p))
            txt = "\n".join((para.text or "").strip() for para in d.paragraphs if (para.text or "").strip())
            if txt.strip():
                return file_name, txt[:16000]
        except Exception:
            continue
    return "", ""


async def handle_file_upload(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not _chat_allowed(update):
        return
    if update.message is None:
        return
    if not _claim_update_once(update):
        return
    if update.effective_user and bool(update.effective_user.is_bot):
        return
    chat = update.effective_chat
    if chat is None:
        return
    chat_id = chat.id

    # Semantic priority: if the message indicates AI-label removal, route it
    # to the dedicated pipeline before any generic doc-intake/session logic can
    # claim the file.
    try:
        if _should_route_ai_label_removal_upload(update.message, chat_id):
            from ops.pipelines.ai_label_removal.telegram_handler import (
                run_request as _ai_rm_run,
                build_user_message as _ai_rm_msg,
            )
            doc = update.message.document
            if doc is not None:
                import tempfile

                original_filename = doc.file_name or "document"
                with tempfile.TemporaryDirectory(prefix="ai_label_") as tmpdir:
                    tmpdir_path = Path(tmpdir)
                    input_path = tmpdir_path / original_filename
                    tg_file = await ctx.bot.get_file(doc.file_id)
                    await tg_file.download_to_drive(str(input_path))
                    result = _ai_rm_run(
                        input_path,
                        tmpdir_path / "out",
                        chat_id=str(chat_id),
                        user_id=str(getattr(getattr(update.message, "from_user", None), "id", None) or ""),
                        original_filename=original_filename,
                    )
                    _PENDING_AI_LABEL_REMOVAL.pop(chat_id, None)
                    if result.status == "SUCCESS" and result.output_path:
                        with open(result.output_path, "rb") as fh:
                            await update.message.reply_document(
                                document=fh,
                                filename=result.output_path.name,
                                caption=_ai_rm_msg(result),
                            )
                    else:
                        await update.message.reply_text(_ai_rm_msg(result))
                return
    except Exception as _ai_rm_gate_err:
        log.debug("ai_label_removal file-upload gate skipped: %s", _ai_rm_gate_err)

    video_state = _PENDING_VIDEO.get(chat_id)
    if video_state:
        lang = str(video_state.get("lang", "ru"))
        if video_state.get("await_prompt"):
            await update.message.reply_text(
                "Сначала опишите сцену **текстом** одним сообщением, затем пришлите изображение."
                if lang == "ru" else
                "First describe the scene in **one text message**, then send the image.",
                parse_mode="Markdown",
            )
            return
        prompt = str(video_state.get("prompt", "")).strip()
        model_name = str(video_state.get("model", "")).strip() or None
        mode_name = str(video_state.get("mode", "")).strip() or None
        tmp_path: Path | None = None
        out_path: Path | None = None
        veo_busy = False

        doc = update.message.document
        if doc is not None and _telegram_doc_is_reference_video(doc):
            try:
                suf = Path(doc.file_name or ".mp4").suffix.lower() or ".mp4"
                with tempfile.NamedTemporaryFile(prefix="axi_veo_ref_", suffix=suf, delete=False) as tv:
                    ref_local = Path(tv.name)
                vg = await ctx.bot.get_file(doc.file_id)
                await vg.download_to_drive(str(ref_local))
                old_ref = video_state.get("ref_video_path")
                if old_ref:
                    _unlink_axi_temp(str(old_ref))
                video_state["ref_video_path"] = str(ref_local)
                await update.message.reply_text(
                    "Референс-видео сохранено. Теперь пришлите изображение (jpg/png) — стартовый кадр."
                    if lang == "ru" else
                    "Reference clip saved. Now send a start image (jpg/png).",
                )
                return
            except Exception as e:
                log.warning("veo reference video upload failed: %s", e)
                await update.message.reply_text(
                    f"{AXI_NAME}: не удалось сохранить видео — {e}"
                    if lang == "ru" else
                    f"{AXI_NAME}: could not save video — {e}",
                )
                return

        try:
            image_file = None
            suffix = ".jpg"
            if update.message.photo:
                image_file = await ctx.bot.get_file(update.message.photo[-1].file_id)
                suffix = ".jpg"
            elif update.message.document and (
                (update.message.document.mime_type or "").startswith("image/")
                or (Path(update.message.document.file_name or "").suffix.lower() in {".jpg", ".jpeg", ".png"})
            ):
                image_file = await ctx.bot.get_file(update.message.document.file_id)
                suffix = Path(update.message.document.file_name or "").suffix.lower() or ".jpg"
            if image_file is None:
                has_ref = bool(video_state.get("ref_video_path"))
                await update.message.reply_text(
                    (
                        "Пришлите изображение стартового кадра (jpg/png)."
                        + (" Референс-ролик уже принят — осталось только фото." if has_ref else "")
                    )
                    if lang == "ru" else
                    (
                        "Please send a start image (jpg/png)."
                        + (" Reference clip is already set — photo only is missing." if has_ref else "")
                    ),
                )
                return

            await update.message.reply_text(
                "❌ Генерация видео отключена (требует облачного API)."
                if lang == "ru" else
                "❌ Video generation is disabled (requires cloud API)."
            )
            st_done = _PENDING_VIDEO.pop(chat_id, None)
            if st_done and st_done.get("ref_video_path"):
                _unlink_axi_temp(str(st_done["ref_video_path"]))
            return
        finally:
            if tmp_path is not None:
                tmp_path.unlink(missing_ok=True)
            if veo_busy:
                async with _VEO_GEN_BUSY_LOCK:
                    _VEO_GEN_BUSY.discard(chat_id)

    # ── Priority 1: Single-doc review ────────────────────────────────────────
    if chat_id in _PENDING_DOCREVIEW and _PENDING_DOCREVIEW[chat_id].get("step") == "await_doc":
        await _handle_docreview_file(update, ctx, chat_id)
        return

    # ── Priority 2: Docfill upload interception ───────────────────────────────
    if chat_id in _PENDING_DOCFILL and _PENDING_DOCFILL[chat_id].get("step") in ("await_blank", "await_example"):
        await _handle_docfill_file(update, ctx, chat_id)
        return

    # ── Priority 3: Doc intake — add to active session ────────────────────────
    if chat_id in _PENDING_DOC_INTAKE and _PENDING_DOC_INTAKE[chat_id].get("step") == "await_more":
        await _di_add_document(update, ctx, chat_id)
        return

    # ── Priority 4: Doc intake — start new session for document uploads ───────
    if update.message.document:
        await _di_start_session(update, ctx, chat_id)
        return

    state = _get_pending_analyze_state(chat_id)
    if not state:
        # File-first flow: accept files even before explicit /analyze goal.
        _set_pending_analyze(chat_id, "analyze_uploaded_files_waiting_goal")
        state = _get_pending_analyze_state(chat_id)
        if state is None:
            return
        state["awaiting_clarify"] = True
        state["notified_waiting"] = False
        _pending_analyze_persist()
    pending_prompt = str(state.get("prompt", "")).strip()

    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "en"
    if update.message.photo and not update.message.document:
        await update.message.reply_text(
            "Для /analyze отправьте документ (pdf/docx/txt/xlsx/pptx)." if lang == "ru"
            else "For /analyze, send a document (pdf/docx/txt/xlsx/pptx)."
        )
        return
    if not update.message.document:
        return

    doc = update.message.document
    file_name = doc.file_name or f"upload_{doc.file_id[-8:]}"
    suffix = Path(file_name).suffix or ".bin"
    tmp_path: Path | None = None
    if not state.get("file_ack_sent"):
        state["file_ack_sent"] = True
        await update.message.reply_text(
            "Файл получен. Обрабатываю и добавляю в пакет анализа…"
            if lang == "ru" else
            "File received. Processing and adding to analysis batch…"
        )
    try:
        with tempfile.NamedTemporaryFile(prefix="axi_analyze_", suffix=suffix, delete=False) as tmp:
            tmp_path = Path(tmp.name)
        tg_file = await ctx.bot.get_file(doc.file_id)
        await tg_file.download_to_drive(str(tmp_path))
        excerpt = _extract_text_from_uploaded_document(tmp_path)
        if not excerpt:
            omi_task_id = _tr_register(
                f"OCR+analyze:{pending_prompt[:140]} | file:{file_name}",
                chat_id=str(chat_id),
                source="axi",
            )
            _tr_start(omi_task_id, assigned_to="omi_batch")
            queued = _queue_to_omi_batch_inbox(
                tmp_path,
                file_name,
                meta={
                    "source": "axi",
                    "chat_id": chat_id,
                    "prompt": pending_prompt[:4000],
                    "task_id": omi_task_id,
                },
            )
            tmp_path = None
            # Delegate difficult extraction cases to Omi OCR pipeline instead of failing fast.
            await update.message.reply_text(
                (
                    "Текст извлечь сразу не удалось — передал файл в OCR-пайплайн Omi для обработки. "
                    f"Файл: `{queued.name}`"
                    + (f"\nTask: `{omi_task_id}`" if omi_task_id else "")
                )
                if lang == "ru" else
                (
                    "Could not extract text immediately — delegated file to Omi OCR pipeline. "
                    f"File: `{queued.name}`"
                    + (f"\nTask: `{omi_task_id}`" if omi_task_id else "")
                ),
                parse_mode="Markdown",
            )
            return

        files = state.setdefault("files", [])
        if isinstance(files, list):
            file_entry: dict = {"name": file_name, "excerpt": excerpt}
            if suffix.lower() in (".xlsx", ".xlsm") and tmp_path and tmp_path.exists():
                import hashlib
                import shutil as _shutil
                # Use content hash to avoid duplicating identical files
                file_hash = hashlib.sha256(tmp_path.read_bytes()).hexdigest()[:16]
                AXI_RESULTS_DIR.mkdir(parents=True, exist_ok=True)
                saved_xlsx = AXI_RESULTS_DIR / f"axi_xlsx_{file_hash}_{file_name}"
                # Only copy if file doesn't already exist
                if not saved_xlsx.exists():
                    _shutil.copy2(str(tmp_path), str(saved_xlsx))
                file_entry["path"] = str(saved_xlsx)
            files.append(file_entry)
        state["last_upload_ts"] = time.time()
        state["ts"] = time.time()
        state["lang"] = lang
        if pending_prompt == "analyze_uploaded_files_waiting_goal":
            state["awaiting_clarify"] = True
        else:
            state["awaiting_clarify"] = False
        if not state.get("notified_waiting"):
            state["notified_waiting"] = True
            await update.message.reply_text(
                (
                    "Файл принят. Жду формулировку задачи следующим сообщением "
                    "(что именно проверить/сравнить), затем запускайте `/analyze_done`."
                    if pending_prompt == "analyze_uploaded_files_waiting_goal"
                    else "Принял файлы для пакетного /analyze. Жду завершения загрузки, затем обработаю всё вместе."
                )
                if lang == "ru" else
                (
                    "File accepted. Waiting for your task in the next message "
                    "(what exactly to check/compare), then run `/analyze_done`."
                    if pending_prompt == "analyze_uploaded_files_waiting_goal"
                    else "Files accepted for batch /analyze. Waiting for uploads to finish, then processing all together."
                )
            )
        _pending_analyze_persist()
    finally:
        if tmp_path is not None:
            tmp_path.unlink(missing_ok=True)
    # Defensive fallback: never stay silent on handler failure.
    # If something unexpected happens before normal replies, user still gets actionable message.
    if state.get("file_ack_sent") and not state.get("notified_waiting") and not (state.get("files") or []):
        await update.message.reply_text(
            "Файл принят, но не попал в пакет. Повторите загрузку или отправьте `/analyze_done` после повторной загрузки."
            if lang == "ru" else
            "File was received but not added to batch. Re-upload and then send `/analyze_done`."
        )


def _extract_analyze_text(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> str:
    parts: list[str] = []
    if getattr(ctx, "args", None):
        parts.append(" ".join(ctx.args).strip())
    msg = update.message
    if msg and msg.reply_to_message:
        reply = msg.reply_to_message
        if reply.text:
            parts.append(reply.text.strip())
        elif reply.caption:
            parts.append(reply.caption.strip())
    return "\n\n".join([p for p in parts if p]).strip()


async def _process_request_text(
    update: Update,
    ctx: ContextTypes.DEFAULT_TYPE,
    text: str,
    *,
    require_group_mention: bool,
) -> None:
    _chat_id = update.effective_chat.id if update.effective_chat else None
    lang = _reply_lang(text)
    chat_type = update.message.chat.type
    prior_ctx = _dialog_context(_chat_id)
    _dialog_append(_chat_id, "user", text)

    # В группе — только если бот упомянут (для свободного текста)
    bot_username = ctx.bot.username
    if require_group_mention and chat_type in ("group", "supergroup"):
        if f"@{bot_username}" not in text and not text.lower().startswith("axi"):
            return
        text = text.replace(f"@{bot_username}", "").strip()
    text = _strip_axi_prefix(text)

    if await _is_omi_owned_intent(text):
        log.info("Axi delegating DOCSREG intent to Omi: chat_id=%s", _chat_id)
        return

    if not _wants_docx(text) and _should_route_ai_label_removal_text(text, _chat_id):
        _PENDING_AI_LABEL_REMOVAL[_chat_id or -1] = {"prompt": text, "ts": time.time()}
        await update.message.reply_text(
            "Прикрепите документ (DOCX, PPTX, XLSX или PDF), и я удалю метки ИИ."
        )
        return

    # NLP intent routing — map free text to slash commands via local small model
    try:
        from chat_intent_router import classify, AXI_CMDS  # noqa: PLC0415
        _axi_routed = await asyncio.to_thread(classify, text, AXI_CMDS)
        if _axi_routed:
            _axi_cmd, _axi_args = _axi_routed
            _axi_dispatch = {
                "quality_report": cmd_quality_report,
                "stuck_tasks": cmd_stuck_tasks,
                "analyze": cmd_analyze,
            }
            _axi_handler = _axi_dispatch.get(_axi_cmd)
            if _axi_handler:
                ctx.args = _axi_args
                await _axi_handler(update, ctx)
                return
    except Exception as e:
        log.debug("intent_router error: %s", e)

    if _looks_like_task_close_intent(text):
        handled, msg = _close_omi_task_from_chat(str(_chat_id or ""), text)
        if handled:
            await update.message.reply_text(msg, parse_mode="Markdown")
            return

    veo_req = _extract_video_skill_request(text)
    if veo_req is None:
        veo_req = _extract_video_skill_alias_request(text)
    if _chat_id is not None and veo_req is not None:
        veo_prompt, veo_model, veo_mode, needs_description, veo_queue, veo_req_sec = veo_req
        prev_v = _PENDING_VIDEO.get(_chat_id)
        if prev_v and prev_v.get("ref_video_path"):
            _unlink_axi_temp(str(prev_v["ref_video_path"]))
        _PENDING_VIDEO[_chat_id] = {
            "prompt": veo_prompt[:1000],
            "prompt_queue": list(veo_queue or []),
            "await_prompt": needs_description,
            "ref_video_path": None,
            "model": (veo_model or "").strip()[:120],
            "mode": (veo_mode or "").strip()[:30],
            "lang": lang,
            "ts": time.time(),
            "duration_seconds": veo_req_sec,
        }
        if needs_description:
            await update.message.reply_text(
                (
                    "Сначала **одним сообщением** опишите кадр: движение камеры, свет, атмосфера.\n"
                    "По желанию добавьте с новой строки:\n"
                    "`озвучка:` или `декламация:` — короткий текст, который должно произнести видео (если модель поддерживает речь/субтитры).\n"
                    "Затем пришлите изображение (jpg/png)."
                    if lang == "ru" else
                    "First, in **one message**, describe the shot: camera motion, light, mood.\n"
                    "Optionally add a new line:\n"
                    "`voiceover:` or `narration:` — short line to speak (if the model supports speech/subtitles).\n"
                    "Then send an image (jpg/png)."
                ),
                parse_mode="Markdown",
            )
        else:
            if veo_queue:
                eff = 6
                await update.message.reply_text(
                    (
                        "Видео-генерация отключена (требует облачного API)."
                        if lang == "ru" else
                        "Video generation is disabled (requires cloud API)."
                    ),
                    parse_mode="Markdown",
                )
            else:
                await update.message.reply_text(
                    (
                        "Принято. Пришлите изображение (jpg/png), и я сгенерирую видео "
                        f"на модели `{veo_model}`."
                        if veo_model else
                        (
                            "Принято. Пришлите изображение (jpg/png), и я сгенерирую видео "
                            + ("в режиме КАЧЕСТВО." if veo_mode == "quality" else "в режиме БЫСТРО." if veo_mode == "fast" else "на лучшей доступной модели.")
                        )
                    )
                    if lang == "ru" else
                    (
                        f"Accepted. Upload an image (jpg/png), and I will generate a video using model `{veo_model}`."
                        if veo_model else
                        (
                            "Accepted. Upload an image (jpg/png), and I will generate a video "
                            + ("in QUALITY mode." if veo_mode == "quality" else "in FAST mode." if veo_mode == "fast" else "using the best available model.")
                        )
                    )
                )
        return

    if (
        _chat_id is not None
        and _PENDING_VIDEO.get(_chat_id)
        and _PENDING_VIDEO[_chat_id].get("await_prompt")
    ):
        merged, need_more = _finalize_video_prompt_from_user_text(text)
        if need_more or not merged.strip():
            await update.message.reply_text(
                "Нужно описание **кадра** (движение, свет). Можно добавить строку `озвучка:` / `декламация:` с текстом реплики."
                if lang == "ru" else
                "Please add a **visual** description (motion, light). Optional: a `voiceover:` / `narration:` line.",
                parse_mode="Markdown",
            )
            return
        _PENDING_VIDEO[_chat_id]["prompt"] = merged[:1000]
        _PENDING_VIDEO[_chat_id]["await_prompt"] = False
        await update.message.reply_text(
            (
                "Принято. Дальше по шагам:\n"
                "1) (по желанию) пришлите **короткий референс-ролик** примером движения/стиля — файл **mp4/mov** (до ~32 МБ).\n"
                "2) затем пришлите **изображение** — стартовый кадр (jpg/png).\n"
                "Один клип Veo — обычно **4–8 секунд**; полный длинный текст за раз не озвучить — в конец промпта можно дописать строку `длительность: 8` (максимум для одного запроса)."
                if lang == "ru" else
                "Got it. Next:\n"
                "1) (optional) send a **short reference clip** (mp4/mov, ~32MB max) as an example of motion/style.\n"
                "2) then send a **start image** (jpg/png).\n"
                "One Veo clip is typically **4–8 seconds**; you cannot fit a full long voiceover — add a line `duration: 8` at the end of the prompt for the longest single clip allowed."
            ),
            parse_mode="Markdown",
        )
        return

    if (
        _chat_id is not None
        and _PENDING_VIDEO.get(_chat_id)
        and not _PENDING_VIDEO[_chat_id].get("await_prompt")
        and _is_video_pending_clarification_only(text)
    ):
        await update.message.reply_text(
            "Жду одно изображение: отправьте фото или файл jpg/png (режим генерации видео уже включён)."
            if lang == "ru" else
            "Waiting for one image: send a photo or a jpg/png file (video generation mode is already on).",
        )
        return

    # Check for training_pair intent BEFORE file-dependent check
    if _chat_id is not None:
        low_text = text.lower()
        if any(kw in low_text for kw in ("training pair", "tuning pair", "обучающ пар", "тренировочн пар", "тюнинг пар")):
            # Check if files already uploaded
            has_files = bool(((_get_pending_analyze_state(_chat_id) or {}).get("files") or []))
            if has_files:
                # Files already uploaded, proceed with training_pair workflow
                log.info("Detected training_pair intent with files already uploaded, proceeding to doctuning pipeline")
                # This will be handled by the normal message flow below
                pass
            else:
                # No files yet, ask user to upload
                log.info("Detected training_pair intent, requesting file upload")
                _set_pending_analyze(
                    _chat_id,
                    text,
                    preserve_existing_files=False,
                )
                await update.message.reply_text(
                    (
                        "📚 Задача принята: создание обучающей пары для тюнинга.\n\n"
                        "Загрузите 2 документа:\n"
                        "1️⃣ **Пустой шаблон** (blank template)\n"
                        "2️⃣ **Заполненный пример** (filled example)\n\n"
                        "После загрузки отправьте `/analyze_done`"
                        if lang == "ru" else
                        "📚 Task accepted: creating training pair for tuning.\n\n"
                        "Upload 2 documents:\n"
                        "1️⃣ **Blank template**\n"
                        "2️⃣ **Filled example**\n\n"
                        "After upload, send `/analyze_done`"
                    ),
                    parse_mode="Markdown",
                )
                return
        elif _looks_like_file_dependent_request(text):
            has_files = bool(((_get_pending_analyze_state(_chat_id) or {}).get("files") or []))
            _set_pending_analyze(
                _chat_id,
                text,
                preserve_existing_files=has_files,
            )
            await update.message.reply_text(
                (
                    "Task accepted for document analysis.\n"
                    "Please upload attachment(s) to this chat, then send `/analyze_done`.\n"
                    "I will process the batch once uploads are complete."
                ),
                parse_mode="Markdown",
            )
            return

    # Classify: question vs actionable task using LLM
    # If unclear, ask user for clarification
    async def _classify_message_intent(msg: str, dialog_history: list[dict]) -> str:
        """
        Use LLM to classify message intent: 'question', 'task', or 'unclear'.
        Returns: 'question' | 'task' | 'unclear'
        """
        # Build context from last 10 messages
        context_lines = []
        for m in dialog_history[-10:]:
            role = m.get("role", "")
            content = m.get("content", "")[:200]
            if role == "user":
                context_lines.append(f"User: {content}")
            elif role == "assistant":
                context_lines.append(f"Axi: {content}")

        context_block = "\n".join(context_lines) if context_lines else "(no previous context)"

        classification_prompt = f"""Analyze this message and classify intent.

Previous conversation (last 10 messages):
{context_block}

Current message: "{msg}"

Classification rules:
- QUESTION: user asks for information, explanation, clarification (что?, как?, почему?, can you explain?, what is?)
- TASK: user requests action, document creation, analysis, generation (создай, напиши, найди документы, prepare report, analyze, generate, сгенерируй, tuning pair, training pair)
- UNCLEAR: ambiguous, needs clarification

Special cases that are ALWAYS tasks:
- "generate a tuning pair" or "generate a training pair" → TASK
- "сгенерируй обучающую пару" or "создай тюнинг пару" → TASK

Respond with exactly one word: question, task, or unclear"""

        try:
            # ChatLLM 14Bv16 for intent classification — PC Andrey first, DGX fallback
            import httpx
            _classify_payload = {
                "model": AIMS_CHATLLM_MODEL,
                "prompt": classification_prompt,
                "stream": False,
                "options": {"temperature": 0.1, "num_predict": 10},
            }
            _classify_urls = [
                f"{AIMS_CHATLLM_PC_URL.rstrip('/')}/api/generate",
                f"{AIMS_CHATLLM_DGX_URL.rstrip('/')}/api/generate",
            ]
            response = None
            for _url in _classify_urls:
                try:
                    response = httpx.post(_url, json=_classify_payload, timeout=10.0)
                    if response.status_code == 200:
                        break
                except Exception:
                    response = None

            if response is not None and response.status_code == 200:
                result = response.json().get("response", "").strip().lower()
                if "question" in result:
                    return "question"
                elif "task" in result:
                    return "task"
                elif "unclear" in result:
                    return "unclear"
        except Exception as e:
            log.debug("LLM classification failed, using fallback rules: %s", e)

        # Fallback to simple rules if LLM unavailable
        msg_lower = msg.lower().strip()

        # Check for training_pair/tuning_pair keywords first (always TASK)
        if any(kw in msg_lower for kw in ["training pair", "tuning pair", "обучающ пар", "тренировочн пар", "тюнинг пар"]):
            return "task"

        if any(m in msg_lower for m in ["?", "что такое", "как ", "почему", "what is", "how ", "why"]):
            return "question"

        if any(m in msg_lower for m in ["создай", "напиши", "найди документ", "create", "write", "analyze", "generate", "сгенерируй"]):
            return "task"

        # Short messages without clear markers are likely questions
        if len(msg.split()) < 5:
            return "question"

        return "unclear"

    if _is_casual_non_task_message(text):
        reply = (
            "Hi. I am online and ready."
            if lang != "ru"
            else "Привет. Я на связи и готов к задаче."
        )
        await update.message.reply_text(reply)
        _dialog_append(_chat_id, "assistant", reply)
        return

    # Classify message intent using LLM
    dialog_history = _AXI_DIALOG.get(_chat_id, [])
    intent = await _classify_message_intent(text, dialog_history)

    # If unclear, ask for clarification
    if intent == "unclear":
        clarification_msg = (
            "Уточните, пожалуйста:\n"
            "• Это вопрос (нужна информация)?\n"
            "• Или задача (нужно что-то создать/найти/проанализировать)?\n\n"
            "Ответьте 'вопрос' или 'задача', и я продолжу."
            if lang == "ru" else
            "Please clarify:\n"
            "• Is this a question (need information)?\n"
            "• Or a task (need to create/find/analyze something)?\n\n"
            "Reply 'question' or 'task' and I'll proceed."
        )
        await update.message.reply_text(clarification_msg)
        _dialog_append(_chat_id, "assistant", clarification_msg)
        return

    # Register task only if intent is 'task'
    _tr_id = ""
    if intent == "task":
        _tr_id = _tr_register(
            text[:120] or "(пустое сообщение)",
            chat_id=str(_chat_id or ""),
            source="group" if chat_type in ("group", "supergroup") else "axi",
        )
        _tr_start(_tr_id, assigned_to="gemini")
        if _tr_id:
            task_msg = (
                f"🧾 Задача зарегистрирована: `{_tr_id}`.\n"
                "Используйте этот номер для проверки статуса и корректировки результата."
                if lang == "ru" else
                f"🧾 Task registered: `{_tr_id}`.\n"
                "Use this ID to check status and request result corrections."
            )
            await update.message.reply_text(task_msg, parse_mode="Markdown")
        elif AXI_CHAT_SHOW_TASK_REGISTRY_WARNINGS:
            task_msg = (
                "⚠️ Не удалось зарегистрировать задачу в реестре (Task Registry недоступен)."
                if lang == "ru" else
                "⚠️ Task Registry is unavailable, task ID was not created."
            )
            await update.message.reply_text(task_msg, parse_mode="Markdown")
    else:
        log.debug("Message classified as question, skipping task registration: %s", text[:100])

    t0 = time.perf_counter()
    progress_msg = None
    _main_anim_stop: asyncio.Event | None = None
    _main_anim_task = None
    try:
        progress_msg = await update.message.reply_text(
            f"⏱ {AXI_NAME}: выполняю запрос…" if lang == "ru" else f"⏱ {AXI_NAME}: working on your request…"
        )
        _main_anim_stop = asyncio.Event()
        _main_anim_base = (
            f"{AXI_NAME}: обрабатываю…" if lang == "ru" else f"{AXI_NAME}: processing…"
        )
        _main_anim_task = asyncio.create_task(
            _animate_progress(progress_msg, _main_anim_base, _main_anim_stop)
        )
        await ctx.bot.send_chat_action(chat_id=_chat_id, action=ChatAction.TYPING)
        use_search = _should_web_search(text)

        # Workflow: registry listing — Axi handles directly
        if _wants_registry_list(text):
            if progress_msg is not None:
                try:
                    await progress_msg.delete()
                except Exception:
                    pass
            await _handle_registry_list(update, ctx, text, tr_id=_tr_id)
            return

        # Workflow: внутренние ресурсы (БД) → делегируем Omi
        if _wants_db_strategy(text):
            omi_task_id = _tr_register(
                text[:200],
                chat_id=str(_chat_id or ""),
                source="axi",
            )
            # Оставляем omi_task_id в статусе pending — Omi подберёт его по polling
            lang_is_ru = lang == "ru"
            delegate_msg = (
                "🔄 Задача требует внутренней базы AIMS — передаю Оми, она подготовит документ."
                if lang_is_ru else
                "🔄 Task requires internal AIMS database — delegating to Omi, she will prepare the document."
            )
            if progress_msg is not None:
                try:
                    await progress_msg.delete()
                except Exception:
                    pass
            await update.message.reply_text(delegate_msg)
            _dialog_append(_chat_id, "assistant", delegate_msg)
            _tr_done(_tr_id, summary=f"delegated_to_omi:{omi_task_id}")
            return

        # Workflow: contextual standard discovery / advisory check (feature-flagged)
        if (
            AIMS_ENABLE_CONTEXTUAL_STANDARD_DISCOVERY
            and _wants_contextual_discovery(text)
        ):
            await _handle_contextual_standard_discovery(
                update, ctx, text,
                lang=lang,
                tr_id=_tr_id,
                progress_msg=progress_msg,
                t0=t0,
            )
            return

        # Workflow: standards / compliance check
        pending_state = _get_pending_analyze_state(_chat_id) if _chat_id else None
        pending_files_count = len(list((pending_state or {}).get("files") or []))
        if _wants_standards_check(
            text,
            recent_chat_context=prior_ctx,
            pending_files_count=pending_files_count,
        ):
            await _handle_standards_check(
                update, ctx, text,
                lang=lang,
                tr_id=_tr_id,
                progress_msg=progress_msg,
                t0=t0,
            )
            return

        # Workflow: внешние ресурсы — обрабатывает Axi
        if _wants_docx(text):
            await ctx.bot.send_chat_action(chat_id=_chat_id, action=ChatAction.UPLOAD_DOCUMENT)
            if progress_msg is not None:
                try:
                    await progress_msg.edit_text(
                        f"⏱ {AXI_NAME}: формирую документ (72B)…"
                        if lang == "ru" else
                        f"⏱ {AXI_NAME}: generating document (72B)…"
                    )
                except Exception:
                    pass

            # Route through AIMS V2 LogiOrchestrator (7-step pipeline)
            web_ctx: str | None = None
            _aims_ok = False

            _aims_url = os.environ.get("AIMS_LOGI_URL", "http://localhost:8000/task")
            try:
                import httpx as _httpx
                from pathlib import Path as _Path

                _aims_payload = {
                    "content": text,
                    "doc_type": "procedure",
                    "department": "operations",
                    "requester": "axi_bot",
                }
                async with _httpx.AsyncClient(timeout=180.0) as _hc:
                    _aims_resp = await _hc.post(_aims_url, json=_aims_payload)
                    _aims_resp.raise_for_status()
                    _aims_result = _aims_resp.json()

                _aims_status = _aims_result.get("status", "")

                if _aims_status == "success":
                    docx_path = _Path(_aims_result["document"]["doc_path"])
                    _aims_ok = True
                elif _aims_status == "gap_report":
                    _gap_msg = (
                        "⚠️ Документ не прошёл валидацию ISO. Пробелы соответствия:\n"
                        + "\n".join(f"— {g}" for g in _aims_result.get("gaps", []))
                    ) if lang == "ru" else (
                        "⚠️ Document failed ISO validation. Compliance gaps:\n"
                        + "\n".join(f"— {g}" for g in _aims_result.get("gaps", []))
                    )
                    await update.message.reply_text(_gap_msg)
                    _dialog_append(_chat_id, "assistant", _gap_msg)
                    _tr_done(_tr_id, summary="aims:gap_report")
                    return
                elif _aims_status == "blocked_by_policy":
                    _block_msg = (
                        f"🚫 Заблокировано политикой: {_aims_result.get('reason', '—')}"
                        if lang == "ru" else
                        f"🚫 Blocked by policy: {_aims_result.get('reason', '—')}"
                    )
                    await update.message.reply_text(_block_msg)
                    _dialog_append(_chat_id, "assistant", _block_msg)
                    _tr_done(_tr_id, summary="aims:policy_blocked")
                    return
                elif _aims_status == "training_pair_saved":
                    raise RuntimeError("aims: quality below threshold, saved as training pair")
                else:
                    raise RuntimeError(f"aims: unexpected status={_aims_status!r}")

            except Exception as _aims_err:
                log.warning("aims_logi failed (%s), falling back to DocAgent", _aims_err)

            if not _aims_ok:
                try:
                    from doc_agent_api import DocAgentClient
                    _api_url = os.environ.get("DOC_AGENT_API_URL", "http://doc-agent:8767")
                    _client = DocAgentClient(_api_url)
                    import asyncio as _asyncio
                    loop = _asyncio.get_event_loop()
                    _result = await loop.run_in_executor(
                        None,
                        lambda: _client.generate_dual(text, out_dir=str(AXI_RESULTS_DIR)),
                    )
                    if not _result.get("ok"):
                        raise RuntimeError(_result.get("error", "dual pipeline failed"))
                    docx_path = _Path(_result["path"])
                except Exception as _da_err:
                    log.warning("doc_agent failed (%s), falling back to local generation", _da_err)
                    docx_prompt = (
                        text + "\n\n"
                        "[INSTRUCTION: Generate the complete document content. "
                        "Use Markdown headings (# ## ###) for sections, **bold** for emphasis, "
                        "bullet lists with '- '. Output only document content, no meta-commentary.]"
                    )
                    answer = await _llm_reply(docx_prompt, extra_context=prior_ctx, use_search=use_search)
                    first_line = next(
                        (ln.strip().lstrip("#").strip() for ln in answer.splitlines() if ln.strip()),
                        "axi_document",
                    )
                    docx_path = await _generate_custom_docx(answer, "axi_doc", AXI_RESULTS_DIR, title=first_line)
            dt = time.perf_counter() - t0
            source_hint = _sources_note(
                lang=lang,
                use_search=use_search,
                has_dialog_context=bool(prior_ctx.strip()),
            )
            task_hint = f" | task `{_tr_id}`" if _tr_id else ""
            if progress_msg is not None:
                try:
                    await progress_msg.delete()
                except Exception:
                    pass
            with docx_path.open("rb") as fh:
                await update.message.reply_document(
                    document=fh,
                    filename=docx_path.name,
                    caption=f"{AXI_NAME}: ✅ {docx_path.name} ({dt:.0f}с)\n{source_hint}{task_hint}",
                )
            _dialog_append(_chat_id, "assistant", f"[doc:{docx_path.name}]")
            _tr_done(_tr_id, summary=f"docx:{docx_path.name}")
        else:
            answer = await _llm_reply(text, extra_context=prior_ctx, use_search=use_search)
            answer = _strip_false_attachment_claims(answer, lang=lang)
            answer = (
                answer.strip()
                + "\n\n"
                + _sources_note(
                    lang=lang,
                    use_search=use_search,
                    has_dialog_context=bool(prior_ctx.strip()),
                )
            )
            dt = time.perf_counter() - t0
            _TG_LIMIT = 4000
            chunks = [answer[i:i + _TG_LIMIT] for i in range(0, len(answer), _TG_LIMIT)] if answer else ["…"]
            for i, chunk in enumerate(chunks):
                try:
                    await update.message.reply_text(chunk, parse_mode="Markdown")
                except Exception:
                    await update.message.reply_text(chunk)
            if progress_msg is not None:
                try:
                    await progress_msg.delete()
                except Exception:
                    pass
            _dialog_append(_chat_id, "assistant", answer)
            _tr_done(_tr_id, summary=f"chat:dt={dt:.1f}s")

    except Exception as e:
        log.error("handle_chat error: %s", e, exc_info=True)
        _tr_stuck(_tr_id, error=str(e)[:200])
        try:
            from failure_detector import get_detector
            get_detector().from_axi({"ok": False, "error": str(e), "chat_id": _chat_id})
        except Exception:
            pass
        err_msg = (
            f"{AXI_NAME}: не удалось обработать запрос — {type(e).__name__}."
            if lang == "ru" else
            f"{AXI_NAME}: failed to process request — {type(e).__name__}."
        )
        try:
            await update.message.reply_text(err_msg)
        except Exception:
            pass
    finally:
        if _main_anim_stop is not None:
            _main_anim_stop.set()
        if _main_anim_task is not None:
            _main_anim_task.cancel()
            try:
                await _main_anim_task
            except asyncio.CancelledError:
                pass
        if progress_msg is not None:
            try:
                await progress_msg.delete()
            except Exception:
                pass


# ── Docfill — training pair builder ───────────────────────────────────────────

_DOCFILL_NIM_BASE = "https://integrate.api.nvidia.com/v1"
_DOCFILL_NIM_MODEL = "deepseek-ai/deepseek-v3-0324"

def _build_validate_system(standards: list[str]) -> str:
    std_line = (
        ", ".join(standards)
        if standards
        else "applicable industrial maintenance and safety standards"
    )
    return (
        "You are an expert industrial document quality auditor for the AIMS platform. "
        f"Check compliance against these standards: {std_line}. "
        "Assess whether the filled document is factually correct and free of hallucinations. "
        "Respond ONLY with valid JSON — no markdown fences, no extra text:\n"
        '{"hallucination_score": <0.0–1.0>, '
        '"issues": ["<specific issue>" ...], '
        '"standards_met": ["<standard>" ...], '
        '"verdict": "accept" | "reject", '
        '"reason": "<one sentence>"}'
    )


_DOCFILL_CONTEXT_SYSTEM = (
    "You are a technical document classifier. Analyze the given blank form and filled example. "
    "Select standards DIRECTLY relevant to what this document IS — its type and purpose. "
    "For document templates, abstract forms, or reporting documents: prefer ISO 55001 (asset management), "
    "ISO 9001 (quality), or records/lifecycle standards. "
    "Do NOT suggest equipment inspection, hot work permit, confined space entry, or HSE procedure standards "
    "unless the document explicitly covers those operations. "
    "Respond ONLY with valid JSON — no markdown fences, no extra text:\n"
    '{"doc_type": "<type>", "equipment_type": "<equipment or process>", '
    '"industry": "<sector>", "key_terms": ["<term>", ...], "standards": ["<standard>", ...]}'
)


def _detect_template_reference_roles(doc_texts: list[dict]) -> tuple[dict, dict, str]:
    """Return (blank_doc, filled_doc, reason) using filename hints then content density.

    Filename keywords take priority over content analysis:
      template hints : template, blank, form, empty
      reference hints: completed, filled, reference, example, gemini, chatgpt

    Diagnostic snippet:
      blank, filled, reason = _detect_template_reference_roles([
          {"name": "Completed_Downstream_Abstract_Lifecycle_ROI_Asset_Preservation_Chatgpt.xlsx", "text": "a b c"},
          {"name": "Abstract Downstream template.xlsx", "text": "_ _ _"},
      ])
      assert filled["name"].startswith("Completed"), reason
      assert "template" in blank["name"], reason
    """
    TEMPLATE_HINTS = {"template", "blank", "form", "empty"}
    REFERENCE_HINTS = {"completed", "filled", "reference", "example", "gemini", "chatgpt"}

    def _score(name: str) -> tuple[int, int]:
        n = name.lower()
        return (
            sum(1 for h in TEMPLATE_HINTS if h in n),
            sum(1 for h in REFERENCE_HINTS if h in n),
        )

    if len(doc_texts) >= 2:
        t0, r0 = _score(doc_texts[0].get("name", ""))
        t1, r1 = _score(doc_texts[1].get("name", ""))

        if (r0 > t0) and (t1 > r1):
            return (
                doc_texts[1], doc_texts[0],
                f"filename: '{doc_texts[1]['name']}' → template, '{doc_texts[0]['name']}' → reference",
            )
        if (t0 > r0) and (r1 > t1):
            return (
                doc_texts[0], doc_texts[1],
                f"filename: '{doc_texts[0]['name']}' → template, '{doc_texts[1]['name']}' → reference",
            )
        if t0 > r0:
            return doc_texts[0], doc_texts[1], f"filename hint: '{doc_texts[0]['name']}' has template keyword"
        if t1 > r1:
            return doc_texts[1], doc_texts[0], f"filename hint: '{doc_texts[1]['name']}' has template keyword"
        if r0 > t0:
            return doc_texts[1], doc_texts[0], f"filename hint: '{doc_texts[0]['name']}' has reference keyword, treating doc[1] as template"
        if r1 > t1:
            return doc_texts[0], doc_texts[1], f"filename hint: '{doc_texts[1]['name']}' has reference keyword, treating doc[0] as template"

        # No filename signal — fall back to content density
        _SPARSE = {"_", "__", "___", "N/A", "TBD", "[", "]", "...", "—"}

        def _density(text: str) -> float:
            toks = text.split()
            if not toks:
                return 0.0
            return sum(1 for t in toks if t not in _SPARSE) / len(toks)

        d0 = _density(doc_texts[0].get("text", ""))
        d1 = _density(doc_texts[1].get("text", ""))
        if d0 <= d1:
            return doc_texts[0], doc_texts[1], f"content density: '{doc_texts[0]['name']}' is sparser (likely blank)"
        return doc_texts[1], doc_texts[0], f"content density: '{doc_texts[1]['name']}' is sparser (likely blank)"

    if doc_texts:
        return doc_texts[0], doc_texts[0], "single document — used as both blank and filled"
    return {}, {}, "no documents"


# ── OmniRoute session management ─────────────────────────────────────────────

def _omniroute_base_urls() -> tuple[str, str]:
    """Return (base_without_v1, base_with_v1) for the configured OmniRoute endpoint.

    Canonical endpoint: 127.0.0.1:20129 (per AIMS project spec)
    """
    raw = (
        os.environ.get("AIMS_OMNIROUTE_BASE_URL")
        or os.environ.get("OMNIROUTE_BASE_URL")
        or os.environ.get("AIMS_OMNIROUTER_URL")
        or "http://127.0.0.1:20129/v1"
    ).rstrip("/")
    if raw.endswith("/v1"):
        return raw[:-3], raw
    return raw, raw + "/v1"


def _omniroute_get_active_sessions_sync() -> list[dict]:
    """Query known OmniRoute session endpoints. Returns [] on total failure."""
    try:
        import httpx
    except ImportError:
        return []

    base_without, base_with = _omniroute_base_urls()
    auth_token = (
        os.environ.get("OMNIROUTE_API_KEY")
        or os.environ.get("AIMS_CLAUDE_PROXY_TOKEN")
        or "aims-local-repair-token"
    )
    headers = {"Authorization": f"Bearer {auth_token}"}

    candidate_urls = [
        f"{base_without}/sessions",
        f"{base_without}/api/sessions",
        f"{base_without}/v1/sessions",
        f"{base_with}/sessions",
    ]

    for url in candidate_urls:
        try:
            with httpx.Client(timeout=10.0) as client:
                resp = client.get(url, headers=headers)
            if resp.status_code != 200:
                continue
            data = resp.json()
            if isinstance(data, list):
                raw_sessions = data
            elif isinstance(data, dict):
                raw_sessions = (
                    data.get("sessions")
                    or data.get("data")
                    or data.get("items")
                    or []
                )
            else:
                continue

            result = []
            for s in raw_sessions:
                if not isinstance(s, dict):
                    continue
                result.append({
                    "id": s.get("id") or s.get("session_id") or s.get("sessionId", ""),
                    "created_at": s.get("created_at") or s.get("createdAt") or s.get("created") or "",
                    "updated_at": s.get("updated_at") or s.get("updatedAt") or s.get("updated") or "",
                    "status": (s.get("status") or "active"),
                    "raw": s,
                })
            log.info("omniroute sessions: found %d via %s", len(result), url)
            return result
        except Exception:
            continue

    log.info("omniroute sessions: no sessions endpoint responded")
    return []


def _omniroute_close_session_sync(session_id: str) -> bool:
    """Try all known patterns to close one OmniRoute session. Returns True on 2xx or 404."""
    if not session_id:
        return False
    try:
        import httpx
    except ImportError:
        return False

    base_without, base_with = _omniroute_base_urls()
    auth_token = (
        os.environ.get("OMNIROUTE_API_KEY")
        or os.environ.get("AIMS_CLAUDE_PROXY_TOKEN")
        or "aims-local-repair-token"
    )
    headers = {"Authorization": f"Bearer {auth_token}"}

    candidate_actions = [
        ("DELETE", f"{base_without}/sessions/{session_id}"),
        ("POST",   f"{base_without}/sessions/{session_id}/close"),
        ("DELETE", f"{base_without}/api/sessions/{session_id}"),
        ("POST",   f"{base_without}/api/sessions/{session_id}/close"),
        ("DELETE", f"{base_with}/sessions/{session_id}"),
        ("POST",   f"{base_with}/sessions/{session_id}/close"),
    ]

    for method, url in candidate_actions:
        try:
            with httpx.Client(timeout=10.0) as client:
                resp = client.request(method, url, headers=headers)
            if resp.status_code in (200, 201, 202, 204, 404):
                log.info("omniroute close session: %s %s → %d", method, url, resp.status_code)
                return True
        except Exception:
            continue

    log.warning("omniroute close session: all endpoints failed for session_id=%s", session_id)
    return False


def _omniroute_debug_session_tables_sync() -> None:
    """Read-only discovery of OmniRoute SQLite session table names. Logs only, no writes."""
    db_path = Path.home() / ".omniroute" / "storage.sqlite"
    if not db_path.exists():
        log.info("omniroute db discovery: %s not found", db_path)
        return
    try:
        import sqlite3
        with sqlite3.connect(str(db_path)) as conn:
            rows = conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' AND lower(name) LIKE '%session%'"
            ).fetchall()
            log.info("omniroute db discovery: session tables=%s in %s",
                     [r[0] for r in rows], db_path)
    except Exception as exc:
        log.info("omniroute db discovery: could not read %s: %s", db_path, exc)


def _omniroute_preflight_free_one_slot_sync(reason: str = "") -> bool:
    """
    Pre-task preflight: if OmniRoute has >= 2 active sessions, close only the oldest one.
    Never closes the newest session. Closes at most one session per call.
    Returns True if a slot is available (or was freed).
    """
    sessions = _omniroute_get_active_sessions_sync()
    # Include sessions with active/running/open/unknown status; if status absent treat as active.
    active = [
        s for s in sessions
        if (s.get("status") or "").lower() in ("active", "running", "open", "", "unknown")
    ] or sessions  # endpoint may only return active sessions — treat all as active

    if len(active) < 2:
        return True

    def _ts(s: dict) -> str:
        return str(s.get("created_at") or s.get("updated_at") or "")

    sorted_sessions = sorted(active, key=_ts)
    oldest = sorted_sessions[0]
    newest = sorted_sessions[-1]

    log.info(
        "omniroute preflight: active_sessions=%d, closing oldest session=%s,"
        " keeping newest session=%s (reason=%s)",
        len(active),
        oldest.get("id", "?"),
        newest.get("id", "?"),
        reason,
    )
    return _omniroute_close_session_sync(oldest.get("id", ""))


def _track_omniroute_session(state: dict, session_id: str) -> None:
    """Record a session ID owned by this task into the task state dict."""
    if not session_id:
        return
    state.setdefault("omniroute_sessions", [])
    if session_id not in state["omniroute_sessions"]:
        state["omniroute_sessions"].append(session_id)


def _cleanup_omniroute_sessions_from_state(state: dict) -> None:
    """Close all OmniRoute sessions recorded by _track_omniroute_session."""
    for sid in state.get("omniroute_sessions", []):
        _omniroute_close_session_sync(sid)
    state["omniroute_sessions"] = []


def _extract_omniroute_content_from_body(body_text: str) -> str:
    """Extract assistant content string from any OmniRoute response body.

    Handles:
    - Normal OpenAI JSON: {"choices":[{"message":{"content":"..."}}]}
    - SSE/streaming: lines with "data: " prefix — concatenates delta.content chunks
    - Anthropic-like: {"content":[{"type":"text","text":"..."}]}
    - Direct fields: text, response, output
    Falls back to body_text unchanged if nothing extracts. Never raises.
    """
    if not body_text:
        return body_text

    # SSE / streaming path
    if "data:" in body_text and "chat.completion.chunk" in body_text:
        parts: list[str] = []
        for line in body_text.splitlines():
            line = line.strip()
            if not line.startswith("data:"):
                continue
            payload = line[5:].strip()
            if payload == "[DONE]" or not payload:
                continue
            try:
                chunk = json.loads(payload)
            except Exception:
                log.debug("omniroute sse chunk skip: %r", payload[:120])
                continue
            choices = chunk.get("choices")
            if not choices:
                continue
            choice = choices[0] if isinstance(choices, list) else choices
            if not isinstance(choice, dict):
                continue
            delta = choice.get("delta") or choice.get("message") or {}
            content = delta.get("content")
            if isinstance(content, str) and content:
                parts.append(content)
        if parts:
            return "".join(parts)
        log.warning("omniroute sse detected but no content chunks extracted; body_start=%s", body_text[:300])

    # Normal JSON path
    try:
        data = json.loads(body_text)
    except Exception:
        return body_text

    if "choices" in data:
        choice = (data.get("choices") or [{}])[0]
        if isinstance(choice, dict):
            msg = choice.get("message") or choice.get("delta") or {}
            content = msg.get("content", "")
            if isinstance(content, str):
                return content.strip()

    if "content" in data:
        blocks = data.get("content") or []
        if isinstance(blocks, list) and blocks and isinstance(blocks[0], dict):
            text = blocks[0].get("text", "")
            if isinstance(text, str):
                return text.strip()

    for _k in ("text", "response", "output"):
        val = data.get(_k)
        if isinstance(val, str):
            return val.strip()

    return body_text


# ─────────────────────────────────────────────────────────────────────────────

def _call_doctuning_openai_omniroute_sync(
    model: str,
    system: str,
    prompt: str,
    max_tokens: int = 512,
    temperature: float = 0.0,
    timeout: float = 120.0,
) -> tuple[int, str]:
    """POST to OmniRoute /v1/chat/completions. Returns (status_code, raw_text)."""
    try:
        import httpx
    except ImportError:
        return 503, "httpx not installed"

    _, base_with = _omniroute_base_urls()
    url = base_with + "/chat/completions"

    auth_token = (
        os.environ.get("OMNIROUTE_API_KEY")
        or os.environ.get("AIMS_CLAUDE_PROXY_TOKEN")
        or "aims-local-repair-token"
    )

    log.info("doctuning omniroute chat: url=%s model=%s", url, model)

    # Preflight: free a slot if already at the 2-session limit
    _omniroute_preflight_free_one_slot_sync(reason=f"model={model}")

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        "temperature": temperature,
        "max_tokens": max_tokens,
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {auth_token}",
        "Content-Type": "application/json",
    }

    def _do_request() -> tuple[int, str]:
        # Phase 1: network — 503 only if the request itself fails before a response exists
        try:
            with httpx.Client(timeout=timeout) as client:
                resp = client.post(url, headers=headers, json=payload)
        except Exception as net_exc:
            return 503, str(net_exc)

        # Phase 2: capture status and raw body immediately — never downgrade status_code
        status_code = resp.status_code
        body_text = resp.text

        session_id = (
            resp.headers.get("x-omniroute-session-id")
            or resp.headers.get("x-session-id")
            or ""
        )
        if session_id:
            log.info("doctuning omniroute session-id: %s", session_id)

        if status_code == 429:
            log.warning("doctuning omniroute rate limited: model=%s", model)
            return 429, body_text

        # Phase 3: extract assistant content — handles JSON, SSE chunks, direct fields
        raw = _extract_omniroute_content_from_body(body_text)
        if raw is body_text and raw != body_text:
            log.warning(
                "doctuning omniroute parse fallback: status=%s model=%s body_start=%s",
                status_code, model, body_text[:500],
            )

        return status_code, raw

    import time as _time
    status, raw = _do_request()
    if status == 429:
        if "maximum number of active sessions" in raw:
            # Session-limit 429: run preflight cleanup then retry after short wait
            _omniroute_preflight_free_one_slot_sync(reason="429_retry")
            _time.sleep(5)
        else:
            # Generic rate-limit: longer back-off
            _time.sleep(20)
        status, raw = _do_request()
    return status, raw


def _extract_doc_context_sync(blank_text: str, filled_text: str) -> dict:
    """Classify document and extract context via OmniRoute /v1/chat/completions."""
    model = (
        os.environ.get("AIMS_DOCTUNING_EXTRACT_MODEL")
        or os.environ.get("OMNIROUTE_EXTRACT_MODEL")
        or os.environ.get("AIMS_DOCTUNING_MODEL")
        or "doc-extract-combo"
    )

    prompt = (
        f"BLANK TEMPLATE (first 2000 chars):\n{blank_text[:2000]}\n\n"
        f"FILLED EXAMPLE (first 2000 chars):\n{filled_text[:2000]}\n\n"
        "Select standards that match the TYPE of this document, not incidental keywords. "
        "Respond with JSON only."
    )

    status, raw = _call_doctuning_openai_omniroute_sync(
        model=model,
        system=_DOCFILL_CONTEXT_SYSTEM,
        prompt=prompt,
        max_tokens=512,
        temperature=0.0,
        timeout=120.0,
    )
    log.info("doctuning context extract: model=%s status=%d", model, status)

    if status == 429:
        if "maximum number of active sessions" in raw:
            return {"error": "omniroute_session_limit", "raw": raw[:500]}
        return {"error": "omniroute_rate_limited", "raw": raw[:500]}

    if status != 200:
        return {"error": f"omniroute_status_{status}", "raw": raw[:500]}

    # SSE guard — extract content if helper was bypassed
    if raw.lstrip().startswith("data:") and "chat.completion.chunk" in raw:
        raw = _extract_omniroute_content_from_body(raw)

    # strip markdown fences and leading "json" token
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.lstrip("`")
        if cleaned.startswith("json"):
            cleaned = cleaned[4:]
        cleaned = cleaned.strip()
    if cleaned.endswith("```"):
        cleaned = cleaned[:-3].strip()
    cleaned = cleaned.strip("`").strip()
    if cleaned.lower().startswith("json"):
        cleaned = cleaned[4:].strip()

    try:
        return json.loads(cleaned)
    except Exception:
        return {"error": "json_parse_failed", "raw": cleaned[:500]}


def _tg_plain(text: str) -> str:
    """Return text as a plain string (no-op; used to document intent at call sites)."""
    return str(text or "")


async def _safe_send_message(bot, chat_id: int, text: str, **kwargs) -> None:
    """Send a Telegram message, retrying without parse_mode on Markdown entity errors."""
    import telegram.error as _tg_err
    try:
        await bot.send_message(chat_id=chat_id, text=text, **kwargs)
    except _tg_err.BadRequest as _e:
        if "can't parse entities" in str(_e).lower() or "Can't parse entities" in str(_e):
            kwargs.pop("parse_mode", None)
            await bot.send_message(chat_id=chat_id, text=text, **kwargs)
        else:
            raise


async def _di_create_case_package(
    chat_id: int,
    task_description: str,
    mode: str,
    files: list[dict],
    selected_pipeline: str,
    learning_contour: bool = False,
) -> "Path":
    """Create a case package in the intake dir; optionally sync to docs-agent inbox."""
    import shutil as _shutil
    import hashlib as _hashlib
    import uuid as _uuid
    from datetime import datetime as _dt

    case_id = str(_uuid.uuid4())
    case_dir = _DI_INTAKE_DIR / case_id
    docs_dir = case_dir / "documents"
    docs_dir.mkdir(parents=True, exist_ok=True)

    doc_entries = []
    for f in files:
        src = Path(f["local_path"])
        dst = docs_dir / f["filename"]
        _shutil.copy2(src, dst)
        sha256 = _hashlib.sha256(dst.read_bytes()).hexdigest()
        doc_entries.append({
            "filename": f["filename"],
            "original_name": f.get("original_name", f["filename"]),
            "sha256": sha256,
            "size_bytes": dst.stat().st_size,
        })

    import json as _json
    manifest = {
        "case_id": case_id,
        "created_at": _dt.utcnow().isoformat() + "Z",
        "chat_id": chat_id,
        "task_description": task_description,
        "mode": mode,
        "selected_pipeline": selected_pipeline,
        "document_count": len(files),
        "documents": doc_entries,
        "learning_contour": learning_contour,
        "registration_forbidden": True,
        "raw_upload": True,
    }
    (case_dir / "manifest.json").write_text(
        _json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    if _DI_SYNC_ENABLED and mode != "tuning_pair":
        dst_inbox = _DI_INBOX_DIR / case_id
        _shutil.copytree(str(case_dir), str(dst_inbox))

    return case_dir


def _doctuning_ws() -> Path:
    """Resolve the AIMS workspace root for doctuning artifacts.

    Priority: AXI_DATA_ROOT → AIMS_WORKSPACE → /data (if exists) →
              /aims_workspace (if exists) → sibling aims_workspace dir.
    """
    for _env in ("AXI_DATA_ROOT", "AIMS_WORKSPACE"):
        _v = os.environ.get(_env)
        if _v:
            return Path(_v)
    for _candidate in (Path("/data"), Path("/aims_workspace")):
        if _candidate.exists():
            return _candidate
    return Path(__file__).resolve().parent.parent / "aims_workspace"


def _save_doctuning_memo(context: dict, blank_name: str, content_hash: str = None) -> Path:
    ws = _doctuning_ws()
    memo_dir = ws / "training" / "doctuning_memos"
    memo_dir.mkdir(parents=True, exist_ok=True)
    memo_path = memo_dir / f"doctuning_memo_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.json"
    memo_data = {
        "blank_name": blank_name,
        "content_hash": content_hash,
        "created_at": datetime.now(timezone.utc).isoformat(),
        **context
    }
    memo_path.write_text(
        json.dumps(memo_data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return memo_path


def _find_existing_training_pair(content_hash: str) -> dict | None:
    """Check if training pair with this content hash already exists."""
    ws = _doctuning_ws()
    memo_dir = ws / "training" / "doctuning_memos"
    if not memo_dir.exists():
        return None

    for memo_file in memo_dir.glob("doctuning_memo_*.json"):
        try:
            memo_data = json.loads(memo_file.read_text(encoding="utf-8"))
            if memo_data.get("content_hash") == content_hash:
                return {
                    "memo_file": memo_file.name,
                    "created": memo_data.get("created_at", "unknown"),
                    "local_score": memo_data.get("local_model_score", "не оценено"),
                    "blank_name": memo_data.get("blank_name", "unknown"),
                }
        except Exception as e:
            log.warning("Failed to read memo %s: %s", memo_file.name, e)
            continue

    return None


def _save_doctuning_failure_case(
    *,
    stage: str,
    reason: str,
    chat_id: int,
    state: dict | None = None,
    pair_path: str | None = None,
    severity: str = "warning",
    recoverable: bool = True,
) -> Path:
    """Save a doctuning failure case JSON for Repairman investigation."""
    import hashlib as _hashlib
    ws = _doctuning_ws()
    fail_dir = ws / "training" / "doctuning_failures"
    fail_dir.mkdir(parents=True, exist_ok=True)
    ts = datetime.now(timezone.utc)
    fname = f"{ts.strftime('%Y%m%d_%H%M%S')}_{stage}_{chat_id}.json"
    fail_path = fail_dir / fname

    _s = state or {}
    _source_text = str(_s.get("blank_text", "") or "")
    _source_hash = _hashlib.sha256(_source_text.encode()).hexdigest()[:16] if _source_text else ""

    data = {
        "source": "doctuning_batch_training_pair",
        "stage": stage,
        "reason": reason,
        "severity": severity,
        "recoverable": recoverable,
        "chat_id": chat_id,
        "pair_path": pair_path or "",
        "master_document_status": str(_s.get("master_document_status", "")),
        "master_document_error": str(_s.get("master_document_error", "")),
        "repairman_status": "pending",
        "created_at": ts.isoformat(),
        "state_summary": {
            "blank_name": str(_s.get("blank_name", "")),
            "filled_name": str(_s.get("filled_name", "")),
            "preferred_output_format": str(_s.get("preferred_output_format", "docx")),
            "source_hash": _source_hash,
            "task_id": str(_s.get("task_id", "")),
        },
    }
    try:
        fail_path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        log.warning(
            "doctuning failure case saved: %s stage=%s reason=%s",
            fail_path, stage, reason[:120],
        )
    except Exception as _fe:
        log.error("doctuning failure case write failed: %s", _fe)
    return fail_path


def _queue_doctuning_repair_case(failure_path: Path, reason: str) -> None:
    """Append a repair queue entry to doctuning_repair_queue.jsonl."""
    ws = _doctuning_ws()
    queue_path = ws / "training" / "doctuning_repair_queue.jsonl"
    queue_path.parent.mkdir(parents=True, exist_ok=True)
    record = {
        "failure_path": str(failure_path),
        "reason": reason,
        "queued_at": datetime.now(timezone.utc).isoformat(),
        "status": "pending_repairman",
    }
    try:
        with queue_path.open("a", encoding="utf-8") as _qf:
            _qf.write(json.dumps(record, ensure_ascii=False) + "\n")
        log.warning("doctuning repair case queued: %s", failure_path)
    except Exception as _qe:
        log.error("doctuning repair queue write failed: %s", _qe)


def _request_repairman_investigation_for_doctuning_failure(failure_path: Path) -> bool:
    """Submit a doctuning failure case to RepairmanAPI (port 8010).

    Returns True if submitted, False if unavailable/failed (case queued instead).
    Updates repairman_status field in the failure JSON.
    """
    _repairman_url = os.environ.get("AIMS_REPAIRMAN_API_URL", "http://localhost:8010")
    _new_status = "pending"
    try:
        _data = json.loads(failure_path.read_text(encoding="utf-8"))
        _task = (
            f"Doctuning failure: stage={_data.get('stage')} "
            f"reason={_data.get('reason', '')[:200]} "
            f"pair_path={_data.get('pair_path', '')} "
            f"severity={_data.get('severity')}"
        )
        import httpx as _httpx_rm
        with _httpx_rm.Client(timeout=5.0) as _c:
            _resp = _c.post(
                f"{_repairman_url}/repair",
                json={"task": _task, "mode": "inspect", "source": "doctuning_failure"},
            )
        _resp.raise_for_status()
        _new_status = "sent"
        _data["repairman_status"] = _new_status
        failure_path.write_text(json.dumps(_data, indent=2, ensure_ascii=False), encoding="utf-8")
        log.info("doctuning repairman investigation submitted: %s", failure_path)
        return True
    except Exception as _re:
        _err_str = str(_re).lower()
        _new_status = (
            "unavailable"
            if ("connection refused" in _err_str or "connectionrefused" in _err_str)
            else "failed"
        )
        try:
            _data = json.loads(failure_path.read_text(encoding="utf-8"))
            _data["repairman_status"] = _new_status
            failure_path.write_text(json.dumps(_data, indent=2, ensure_ascii=False), encoding="utf-8")
        except Exception:
            pass
        log.warning(
            "doctuning: repairman investigation not submitted (%s): %s", _new_status, _re
        )
        _queue_doctuning_repair_case(failure_path, str(_re)[:200])
        return False


def _read_document_text(file_path: Path) -> str:
    """Read text from .txt, .json, .docx, or .xlsx file."""
    suffix = file_path.suffix.lower()

    if suffix in (".txt", ".json"):
        return file_path.read_text(encoding="utf-8", errors="replace")

    elif suffix == ".docx":
        try:
            from docx import Document
            doc = Document(str(file_path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            log.warning("docx read error: %s", e)
            return file_path.read_text(encoding="utf-8", errors="replace")

    elif suffix == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        lines.append(row_text)
            return "\n".join(lines)
        except Exception as e:
            log.warning("xlsx read error: %s", e)
            return file_path.read_text(encoding="utf-8", errors="replace")

    else:
        return file_path.read_text(encoding="utf-8", errors="replace")


async def _search_standards_for_context(context: dict) -> list[str]:
    """Web search for standards is disabled in local-only mode."""
    return []


def _validate_and_extract_context_sync(blank_text: str, filled_text: str) -> dict:
    """Combined: extract context + validate + produce revised reference via OmniRoute /chat/completions.

    Returns dict with:
    - context: {doc_type, equipment_type, industry, key_terms, standards}
    - validation: {status, hallucination_score, issues, standards_met, reason,
                   missing_fields, unsupported_claims, recommendations}
      hallucination_score is None on hard failure (HTTP error / JSON parse failure).
    - revised_reference: str or None (None on hard failure)
    - error: present only on failure — "omniroute_status_<N>" or "json_parse_failed"
    - raw: first 500 chars of raw response (only on failure)
    """
    model = (
        os.environ.get("AIMS_DOCTUNING_VALIDATE_MODEL")
        or os.environ.get("OMNIROUTE_AUDIT_MODEL")
        or os.environ.get("AIMS_DOCTUNING_MODEL")
        or "doc-training-pair-audit-combo"
    )

    def _failed_validation(reason: str) -> dict:
        return {
            "status": "failed",
            "hallucination_score": None,
            "evidence_risk": None,
            "issues": ["teacher_audit_failed"],
            "standards_met": [],
            "reason": reason,
            "missing_fields": [],
            "unsupported_claims": [],
            "validation_gaps": [],
            "recommendations": [],
        }

    combined_system = (
        "You are an expert industrial document specialist and quality auditor for the AIMS platform. "
        "Perform THREE tasks and respond with a single valid JSON object. "
        "No markdown fences, no extra text outside the JSON.\n\n"
        "Task 1 — CLASSIFY and extract context:\n"
        '  "context": {\n'
        '    "doc_type": "<procedure|specification|report|form|checklist|...>",\n'
        '    "equipment_type": "<equipment or process name>",\n'
        '    "industry": "<sector>",\n'
        '    "key_terms": ["<term>", ...],\n'
        '    "standards": ["<applicable standard>", ...]\n'
        "  }\n\n"
        "Task 2 — VALIDATE the filled reference for hallucinations, fabricated values, "
        "missing required fields, and standards compliance (ISO, API, ASME, OSHA, IEC, NFPA, "
        "IEEE, EN, ANSI, PAS, BS). Do NOT invent page or section references.\n"
        "  Status rules — use exactly one of these values:\n"
        '    "failed"  — any of: hallucination_score ≥ 0.5; fabricated standard codes or\n'
        "               equipment specifications; required normative sections completely absent.\n"
        '    "warning" — issues present but none are fatal; document is usable with corrections.\n'
        "               ALWAYS use 'warning' (not 'passed') when admin/template fields remain\n"
        "               unfilled: TBD, TODO, Placeholder, [Date], [Author], [Dept], ???.\n"
        '    "passed"  — only when ALL of the following hold: no required fields missing,\n'
        "               hallucination_score < 0.2, no fabricated data, no TBD/placeholder\n"
        "               admin fields, no unsupported normative claims.\n"
        "  Hallucination score vs Evidence risk — treat as SEPARATE dimensions:\n"
        "    hallucination_score — rate ONLY fabricated facts: invented specs, made-up standard\n"
        "               clauses, nonexistent equipment models, false equipment parameters.\n"
        "               Do NOT inflate for 'claims without citations' in overview documents.\n"
        "    evidence_risk — rate unsubstantiated assertions relative to the document type.\n"
        "               High evidence_risk is EXPECTED and acceptable for conference abstracts,\n"
        "               executive summaries, and feasibility overviews — flag it but do not\n"
        "               penalise hallucination_score.\n"
        "  For conference abstracts and overview documents: recommendations should explicitly\n"
        "  note the document type and suggest adding methodology/reference sections only if\n"
        "  the submission venue requires them.\n"
        '  "validation": {\n'
        '    "status": "passed" | "warning" | "failed",\n'
        '    "hallucination_score": <0.0–1.0, fabrications only>,\n'
        '    "evidence_risk": <0.0–1.0, unsubstantiated assertions for this doc type>,\n'
        '    "issues": ["<specific issue>", ...],\n'
        '    "standards_met": ["<standard>", ...],\n'
        '    "reason": "<one sentence overall verdict>",\n'
        '    "missing_fields": ["<field or section missing>", ...],\n'
        '    "unsupported_claims": ["<normative claim with no traceable evidence>", ...],\n'
        '    "validation_gaps": ["<check not performed due to doc type or missing data>", ...],\n'
        '    "recommendations": ["<actionable fix>", ...]\n'
        "  }\n\n"
        "Task 3 — PRODUCE a revised reference document:\n"
        '  "revised_reference": "<full revised document text>"\n'
        "  Rules:\n"
        "  - Do NOT create a new document from scratch.\n"
        "  - Start from the supplied filled reference.\n"
        "  - Preserve all valid content and structure.\n"
        "  - Apply only corrections for issues found in Task 2.\n"
        "  - If no changes are needed, return the supplied reference unchanged.\n"
        "  - Always return revised_reference as a non-empty string."
    )

    prompt = (
        f"BLANK TEMPLATE:\n{blank_text[:3000]}\n\n"
        f"FILLED REFERENCE (to audit and revise):\n{filled_text[:4000]}\n\n"
        "Perform all three tasks and return a single JSON object. "
        "Return ONLY raw JSON. Do not use markdown fences. Do not include text outside JSON."
    )

    log.info(
        "doctuning validate+context+revise: model=%s blank_len=%d filled_len=%d",
        model, len(blank_text), len(filled_text),
    )

    status, raw = _call_doctuning_openai_omniroute_sync(
        model=model,
        system=combined_system,
        prompt=prompt,
        max_tokens=8192,
        temperature=0.1,
        timeout=300.0,
    )

    log.info("doctuning validate+context+revise: status=%d raw_len=%d", status, len(raw))

    if status != 200:
        err_key = f"omniroute_status_{status}"
        log.warning("doctuning validate+context+revise: %s raw_excerpt=%r", err_key, raw[:200])
        return {
            "context": {},
            "validation": _failed_validation(f"{err_key}: {raw[:200]}"),
            "revised_reference": None,
            "error": err_key,
            "raw": raw[:500],
        }

    if not raw:
        log.warning("teacher audit empty response: model=%s", model)
        return {
            "context": {},
            "validation": _failed_validation("teacher_audit_empty_response"),
            "revised_reference": None,
            "error": "teacher_audit_empty_response",
            "raw": "",
        }

    # If raw still looks like SSE (e.g. helper was bypassed), extract content now
    if raw.lstrip().startswith("data:") and "chat.completion.chunk" in raw:
        raw = _extract_omniroute_content_from_body(raw)

    # Robust JSON parsing — strip markdown fences and leading "json" token
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.lstrip("`")
        if cleaned.startswith("json"):
            cleaned = cleaned[4:]
        cleaned = cleaned.strip()
    if cleaned.endswith("```"):
        cleaned = cleaned[:-3].strip()
    cleaned = cleaned.strip("`").strip()
    if cleaned.lower().startswith("json"):
        cleaned = cleaned[4:].strip()

    try:
        result = json.loads(cleaned)
    except Exception as exc:
        log.warning(
            "teacher audit json parse failed: model=%s raw_start=%s",
            model, cleaned[:500],
        )
        return {
            "context": {},
            "validation": _failed_validation(f"json_parse_failed: {exc}; raw_start={cleaned[:500]}"),
            "revised_reference": None,
            "error": "json_parse_failed",
            "raw": cleaned[:500],
        }

    # Ensure all three keys exist
    if "context" not in result:
        result["context"] = {}
    if "validation" not in result:
        result["validation"] = _failed_validation("no validation returned")
    else:
        v = result["validation"]
        if "status" not in v:
            v["status"] = "passed" if not v.get("issues") else "warning"
        for field in ("missing_fields", "unsupported_claims", "recommendations", "validation_gaps"):
            if field not in v:
                v[field] = []
        if "evidence_risk" not in v:
            v["evidence_risk"] = None

        # Deterministic guard: downgrade "passed" → "warning" when TBD/placeholder
        # admin fields are detected in filled_text, regardless of LLM verdict.
        import re as _re
        _TBD_PAT = _re.compile(
            r"\b(TBD|TODO|PLACEHOLDER|TO BE DETERMINED|TO BE CONFIRMED)\b"
            r"|\[Date\]|\[DATE\]|\[Author\]|\[AUTHOR\]|\[Dept\]|\[DEPT\]"
            r"|\[Department\]|\[Reviewer\]|\[REVIEWER\]|\[Approver\]|\[APPROVER\]"
            r"|\[Name\]|\[Title\]|\[Division\]|\[Unit\]"
            r"|\?\?\?",
            _re.IGNORECASE,
        )
        if v.get("status") == "passed" and _TBD_PAT.search(filled_text):
            v["status"] = "warning"
            issues = v.setdefault("issues", [])
            if "admin_fields_incomplete" not in issues:
                issues.append("admin_fields_incomplete")
            mf = v.setdefault("missing_fields", [])
            _tbd_note = "Admin/template placeholder fields (TBD, [Date], [Author], etc.) not filled in"
            if _tbd_note not in mf:
                mf.append(_tbd_note)

    # Ensure revised_reference is present and non-empty (success path only)
    if not result.get("revised_reference"):
        log.warning(
            "doctuning validate+context+revise: revised_reference missing — falling back to filled_text"
        )
        result["revised_reference"] = filled_text
        issues = result["validation"].setdefault("issues", [])
        if "teacher_revision_missing" not in issues:
            issues.append("teacher_revision_missing")

    return result


def _nim_validate_document_sync(blank_text: str, filled_text: str, standards: list[str]) -> dict:
    """DEPRECATED: Use _validate_and_extract_context_sync instead.
    Kept for backward compatibility."""
    try:
        from openai import OpenAI as _OAI
    except ImportError:
        return {"verdict": "accept", "hallucination_score": 0.0, "issues": [], "reason": "openai not installed — skipped validation", "standards_met": []}

    api_key = os.environ.get("NVIDIA_API_KEY", "").strip()
    if not api_key:
        return {"verdict": "accept", "hallucination_score": 0.0, "issues": [], "reason": "NVIDIA_API_KEY not set — skipped validation", "standards_met": []}

    client = _OAI(api_key=api_key, base_url=_DOCFILL_NIM_BASE)
    prompt = (
        f"BLANK TEMPLATE:\n{blank_text[:3000]}\n\n"
        f"FILLED EXAMPLE (to audit):\n{filled_text[:4000]}\n\n"
        "Audit the filled example for hallucinations, fabricated values, and standards compliance. "
        "Return JSON as specified."
    )
    try:
        resp = client.chat.completions.create(
            model=_DOCFILL_NIM_MODEL,
            messages=[
                {"role": "system", "content": _build_validate_system(standards)},
                {"role": "user", "content": prompt},
            ],
            temperature=0.1,
            max_tokens=512,
        )
        raw = (resp.choices[0].message.content or "").strip()
        raw = raw.strip("`").strip()
        if raw.startswith("json"):
            raw = raw[4:].strip()
        return json.loads(raw)
    except Exception as e:
        log.warning("docfill nim validate error: %s", e)
        return {"verdict": "accept", "hallucination_score": 0.0, "issues": [], "reason": f"validation error ({e}) — accepted by default", "standards_met": []}


_DOCFILL_FILL_SYSTEM = (
    "You are an expert industrial document specialist. "
    "You receive a blank technical form and must fill in all fields with realistic values "
    "for an oil & gas / industrial maintenance context, following ISO 55001 and API standards. "
    "Output only the filled form text. No explanations, no comments."
)


def _local_generate_fill_sync(blank_text: str) -> str | None:
    """Ask local ollama model to fill the blank form (DPO rejected — model's current output).

    Uses httpx directly so no openai package is required.
    Ollama exposes OpenAI-compatible /v1/chat/completions natively.
    """
    try:
        import httpx
    except ImportError:
        log.warning("local fill: httpx not available")
        return None

    ollama_url = os.environ.get("OLLAMA_LOCAL_URL", "http://localhost:11434").rstrip("/")
    model = AIMS_DOCUMENT_CREATOR_MODEL

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": _DOCFILL_FILL_SYSTEM},
            {"role": "user", "content": f"Fill in this blank form:\n\n{blank_text[:4000]}"},
        ],
        "stream": False,
        "temperature": 0.3,
        "max_tokens": 2048,
    }

    try:
        with httpx.Client(timeout=120.0) as client:
            resp = client.post(
                f"{ollama_url}/v1/chat/completions",
                json=payload,
                headers={"Content-Type": "application/json"},
            )
        if resp.status_code != 200:
            log.warning("local fill: ollama returned status %d", resp.status_code)
            return None
        data = resp.json()
        result = (
            data.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "")
            or ""
        ).strip()
        return result if len(result) > 80 else None
    except Exception as e:
        log.warning("doctuning local generate error: %s", e)
        return None


def _check_ollama_alive_sync() -> bool:
    """Return True if the local Ollama service responds on /api/tags."""
    import urllib.request as _ur
    ollama_url = os.environ.get("OLLAMA_LOCAL_URL", "http://localhost:11434").rstrip("/")
    try:
        with _ur.urlopen(f"{ollama_url}/api/tags", timeout=5) as r:
            return r.status == 200
    except Exception:
        return False


def _check_vram_free_gb_sync() -> float:
    """Return free VRAM in GB on the first GPU via nvidia-smi, or 999.0 if unavailable.

    For GPUs that don't support memory queries (like GB10), use GPU utilization as proxy:
    - If GPU util < 92%, return high value (assume available)
    - If GPU util >= 92%, return 0.0 (assume full)
    """
    import subprocess as _sp
    try:
        # Try memory.free first (works on most GPUs)
        out = _sp.check_output(
            ["nvidia-smi", "--query-gpu=memory.free", "--format=csv,noheader,nounits"],
            timeout=10, text=True,
        )
        first_line = out.strip().splitlines()[0]
        if first_line != "[N/A]":
            mb = float(first_line)
            return mb / 1024.0

        # Fallback: use GPU utilization percentage as proxy
        out = _sp.check_output(
            ["nvidia-smi", "--query-gpu=utilization.gpu", "--format=csv,noheader,nounits"],
            timeout=10, text=True,
        )
        util_pct = float(out.strip().splitlines()[0])

        # If GPU utilization >= 92%, consider VRAM full
        if util_pct >= 92.0:
            return 0.0
        else:
            return 999.0  # Assume plenty of free VRAM

    except Exception:
        return 999.0  # assume OK if nvidia-smi unavailable (non-GPU env)


# Whitelist of commands the repairman is allowed to execute during local-model recovery.
# Never exec arbitrary shell — only ollama operations and read-only diagnostics.
_REPAIR_CMD_WHITELIST = (
    "ollama list",
    "ollama pull ",
    "ollama stop ",
    "ollama run ",
    "nvidia-smi",
    "curl http://localhost:11434",
    "systemctl status ollama",
    "systemctl restart ollama",
)

_REPAIRMAN_SYSTEM_DOCFILL = (
    "You are the AIMS Repairman — an expert repair agent for the AxiOMSphere platform.\n"
    "Always respond with a single JSON object:\n"
    '{"root_cause": "<one clear paragraph>", "files_changed": [], "patch_diff": "none", '
    '"tests_run": ["<safe shell command>", ...], "test_result": "not_run", '
    '"risk_level": "low", "rollback_notes": "<how to revert>"}\n'
    "Only include safe, non-destructive shell commands in tests_run. "
    "Focus on ollama model recovery: check if the model is loaded, VRAM state, "
    "whether the ollama service is running. Do not delete files or restart services "
    "unless risk_level is low."
)


def _repairman_fix_local_model_sync(model: str) -> str:
    """Call repairman gateway to diagnose and recover the local model. Returns root_cause string.

    Uses httpx when the openai package is not installed.
    """
    import subprocess as _sp
    import json as _json

    gateway_url = "http://localhost:8082/v1"
    token = os.environ.get("AIMS_CLAUDE_PROXY_TOKEN", "aims-local-repair-token")
    problem = (
        f"The local Ollama model '{model}' failed to return a completion in the doctuning "
        "pipeline (returned None or raised an exception). Likely causes: model not loaded, "
        "VRAM exhaustion (DGX 128 GB — check if two large models are loaded simultaneously), "
        "ollama service not running, or connection refused on port 11434. "
        "Diagnose and provide safe recovery commands in 'tests_run'."
    )
    payload = {
        "model": "aims-repairman-nemotron",
        "messages": [
            {"role": "system", "content": _REPAIRMAN_SYSTEM_DOCFILL},
            {"role": "user",   "content": problem},
        ],
        "temperature": 0.1,
        "max_tokens": 512,
    }
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }

    content = ""
    try:
        try:
            from openai import OpenAI as _OAI
            client = _OAI(api_key=token, base_url=gateway_url)
            resp_oa = client.chat.completions.create(
                model="aims-repairman-nemotron",
                messages=payload["messages"],
                temperature=0.1,
                max_tokens=512,
            )
            content = (resp_oa.choices[0].message.content or "").strip()
        except ImportError:
            import httpx as _httpx
            with _httpx.Client(timeout=60.0) as _c:
                _r = _c.post(f"{gateway_url}/chat/completions", json=payload, headers=headers)
            content = (
                _r.json().get("choices", [{}])[0]
                .get("message", {})
                .get("content", "")
                or ""
            ).strip()
    except Exception as e:
        log.warning("repairman gateway call failed: %s", e)
        return f"repairman unavailable: {e}"

    try:
        data = _json.loads(content)
        root_cause = data.get("root_cause", content[:200])
        log.info("repairman root_cause: %s", root_cause)
        for cmd in data.get("tests_run", [])[:3]:
            if isinstance(cmd, str) and any(cmd.strip().startswith(p) for p in _REPAIR_CMD_WHITELIST):
                log.info("repairman exec: %s", cmd)
                _sp.run(cmd, shell=True, timeout=60, capture_output=True)
        return root_cause
    except Exception:
        return content[:300]


def _repairman_final_report_sync(model: str, cycles: int, last_root_cause: str) -> str:
    """After exhausting all repair cycles, ask repairman to write a human-readable failure report.

    Uses httpx when the openai package is not installed.
    """
    gateway_url = "http://localhost:8082/v1"
    token = os.environ.get("AIMS_CLAUDE_PROXY_TOKEN", "aims-local-repair-token")
    problem = (
        f"You attempted to repair the local Ollama model '{model}' {cycles} times. "
        f"Last diagnosed root cause: {last_root_cause}\n"
        f"All {cycles} repair cycles failed — generation still returns no valid output. "
        "Write a concise 2–3 sentence report for the engineer explaining:\n"
        "1. What the root cause appears to be.\n"
        "2. What specific manual action is required to resolve it.\n"
        "Plain text only, no JSON, no markdown."
    )
    system_msg = (
        "You are the AIMS Repairman reporting a repair failure to a human engineer. "
        "Be concise and actionable."
    )
    payload = {
        "model": "aims-repairman-nemotron",
        "messages": [
            {"role": "system", "content": system_msg},
            {"role": "user",   "content": problem},
        ],
        "temperature": 0.1,
        "max_tokens": 256,
    }
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }
    try:
        try:
            from openai import OpenAI as _OAI
            client = _OAI(api_key=token, base_url=gateway_url)
            resp_oa = client.chat.completions.create(
                model="aims-repairman-nemotron",
                messages=payload["messages"],
                temperature=0.1,
                max_tokens=256,
            )
            return (resp_oa.choices[0].message.content or "").strip()
        except ImportError:
            import httpx as _httpx
            with _httpx.Client(timeout=60.0) as _c:
                _r = _c.post(f"{gateway_url}/chat/completions", json=payload, headers=headers)
            return (
                _r.json().get("choices", [{}])[0]
                .get("message", {})
                .get("content", "")
                or ""
            ).strip()
    except Exception as e:
        return (
            f"Repairman unreachable after {cycles} cycles. "
            f"Last known cause: {last_root_cause}. Manual intervention required."
        )


async def _local_fill_with_repair(
    blank_text: str, max_repair_cycles: int = 5
) -> tuple[str | None, str]:
    """
    Generate local fill. On first failure enters the repairman repair loop:
      - Repairman diagnoses, executes recovery commands, verifies (ollama health check).
      - If verified → retry generation.
      - Up to max_repair_cycles cycles.
      - After all cycles exhausted → repairman writes a human-readable failure report.

    Returns (filled_text, repairman_report).
    filled_text is None when all cycles failed; repairman_report contains the final diagnosis.
    """
    loop = asyncio.get_event_loop()
    model = AIMS_DOCUMENT_CREATOR_MODEL

    # First attempt (clean, no repairman pre-check)
    result = await loop.run_in_executor(None, _local_generate_fill_sync, blank_text)
    if result:
        return result, ""

    # Quick repairman reachability check — skip the full cycle loop if repairman is down.
    _gateway_url = "http://localhost:8082"
    try:
        import httpx as _httpx_hc
        with _httpx_hc.Client(timeout=5.0) as _hc:
            _hc.get(f"{_gateway_url}/health")
    except Exception as _hc_exc:
        _hc_msg = str(_hc_exc)
        log.warning("doctuning: repairman unreachable, skipping repair loop: %s", _hc_msg)
        return None, f"repairman unavailable: {_hc_msg}"

    log.warning(
        "doctuning: initial generation failed — starting repairman loop (max %d cycles)",
        max_repair_cycles,
    )
    last_root_cause = "unknown"

    for cycle in range(1, max_repair_cycles + 1):
        log.info("doctuning: repairman cycle %d/%d", cycle, max_repair_cycles)

        # Repairman diagnoses and executes repair commands
        last_root_cause = await loop.run_in_executor(
            None, _repairman_fix_local_model_sync, model
        )
        log.info("repairman cycle %d root_cause: %s", cycle, last_root_cause)

        # Wait for repair to take effect
        await asyncio.sleep(20)

        # Repairman verifies: check ollama health
        alive = await loop.run_in_executor(None, _check_ollama_alive_sync)
        if not alive:
            log.warning("repairman cycle %d: ollama still down after repair", cycle)
            continue  # repairman will try again next cycle

        # Service confirmed up → retry generation
        result = await loop.run_in_executor(None, _local_generate_fill_sync, blank_text)
        if result:
            log.info("doctuning: generation succeeded after repairman cycle %d", cycle)
            return result, ""

        log.warning("repairman cycle %d: service up but generation still fails", cycle)

    # All cycles exhausted — repairman writes the failure report for the engineer
    log.error("doctuning: repairman exhausted %d cycles, requesting final report", max_repair_cycles)
    final_report = await loop.run_in_executor(
        None, _repairman_final_report_sync, model, max_repair_cycles, last_root_cause
    )
    log.error("repairman final report: %s", final_report)
    return None, final_report


def _xlsx_safe_text(value: object) -> str:
    """Sanitize a value for an openpyxl string cell — prevent Excel formula injection.

    Excel and LibreOffice treat any cell value starting with =, +, -, @, TAB, or CR
    as a formula, producing #NAME? or unexpected evaluation.  For string-typed cells
    written via openpyxl we prefix the problematic characters with a single apostrophe
    (the canonical CSV/OOXML defense); the apostrophe is stored in the cell but not
    displayed to the end user in most spreadsheet apps.
    """
    s = str(value) if value is not None else ""
    if s and s[0] in ("=", "+", "-", "@", "\t", "\r", "|"):
        s = "'" + s
    return s


def _create_doctuning_xlsx_candidate(
    template_path: "Path | None",
    reference_path: "Path | None",
    chosen_text: str,
    context: dict,
    validation: dict,
    output_path: "Path",
) -> "Path":
    """Build a validated XLSX candidate using the incoming workbook as base.

    Strategy:
    1. Select source: reference .xlsx first, then template .xlsx.
    2. Copy source → output_path (never modifies the original).
    3. Open copy; add/replace only AIMS_Review sheet; leave all other sheets untouched.
    4. Save.
    5. Validate with openpyxl (reopen).
    6. If LibreOffice is available, additionally validate with headless convert.
    7. On any failure raise — caller handles XLSX policy (no DOCX substitution).
    """
    import shutil
    import subprocess
    import tempfile

    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter
    except ImportError as _ie:
        raise RuntimeError("openpyxl not available") from _ie

    # Select source: filled/reference preferred, then blank/template
    source_path: "Path | None" = None
    for candidate_src in (reference_path, template_path):
        if candidate_src and Path(candidate_src).exists() and Path(candidate_src).suffix.lower() == ".xlsx":
            source_path = Path(candidate_src)
            break

    if source_path is None:
        raise FileNotFoundError(
            "No source .xlsx workbook available "
            "(reference_path and template_path both absent or not .xlsx)"
        )

    # Copy source → output_path; original is never touched
    try:
        shutil.copy2(str(source_path), str(output_path))
        log.info("doctuning xlsx: copied %s → %s", source_path, output_path)
    except Exception as _e:
        raise RuntimeError(f"Could not copy source workbook: {_e}") from _e

    # Open the copy
    try:
        wb = load_workbook(str(output_path))
    except Exception as _e:
        Path(output_path).unlink(missing_ok=True)
        raise RuntimeError(f"Could not open copied workbook: {_e}") from _e

    # Add/replace AIMS_Review only — all other sheets left untouched
    if "AIMS_Review" in wb.sheetnames:
        del wb["AIMS_Review"]
    ws = wb.create_sheet("AIMS_Review")

    header_font = Font(bold=True)
    header_fill = PatternFill("solid", fgColor="DCE6F1")
    wrap = Alignment(wrap_text=True, vertical="top")

    def _row(label: str, value: str, row_idx: int) -> None:
        ws.cell(row=row_idx, column=1, value=label).font = header_font
        ws.cell(row=row_idx, column=1).fill = header_fill
        cell = ws.cell(row=row_idx, column=2, value=_xlsx_safe_text(value))
        cell.alignment = wrap

    row = 1
    ws.cell(row=row, column=1, value="AIMS Doctuning Review").font = Font(bold=True, size=13)
    row += 1

    _row("Document type",       context.get("doc_type", "—"),                                   row); row += 1
    _row("Equipment type",      context.get("equipment_type", "—"),                              row); row += 1
    _row("Industry",            context.get("industry", "—"),                                   row); row += 1
    _row("Standards",           ", ".join(context.get("standards", [])) or "—",                  row); row += 1
    _row("Key terms",           ", ".join(context.get("key_terms", [])) or "—",                  row); row += 1
    _row("Validation status",   validation.get("status", "—"),                                   row); row += 1
    h = validation.get("hallucination_score")
    _row("Hallucination score", f"{h:.3f}" if h is not None else "—",                            row); row += 1
    er = validation.get("evidence_risk")
    _row("Evidence/substantiation risk", f"{er:.3f}" if isinstance(er, float) else (str(er) if er is not None else "—"), row); row += 1
    _row("Issues",              "\n".join(validation.get("issues", [])) or "none",               row); row += 1
    _row("Missing fields",      "\n".join(validation.get("missing_fields", [])) or "none",      row); row += 1
    _row("Validation gaps",     "\n".join(validation.get("validation_gaps", [])) or "none",     row); row += 1
    _row("Unsupported claims",  "\n".join(validation.get("unsupported_claims", [])) or "none",  row); row += 1
    _row("Recommendations",     "\n".join(validation.get("recommendations", [])) or "none",     row); row += 1
    _row("Reason",              validation.get("reason", "—"),                                   row); row += 1

    row += 1
    ws.cell(row=row, column=1, value="Revised reference text").font = Font(bold=True)
    row += 1
    cell = ws.cell(row=row, column=2, value=_xlsx_safe_text(chosen_text))
    cell.alignment = wrap

    row += 2
    _row(
        "Note",
        "Original workbook preserved; AIMS review added as separate sheet.",
        row,
    )

    ws.column_dimensions[get_column_letter(1)].width = 26
    ws.column_dimensions[get_column_letter(2)].width = 90

    try:
        wb.save(str(output_path))
        log.info("doctuning xlsx: saved to %s", output_path)
    except Exception as _e:
        Path(output_path).unlink(missing_ok=True)
        raise RuntimeError(f"openpyxl save failed: {_e}") from _e

    # ── Validation 1: openpyxl reopen ────────────────────────────────────────
    try:
        load_workbook(str(output_path), data_only=False)
        log.info("doctuning xlsx: openpyxl validation passed")
    except Exception as exc:
        log.warning("doctuning xlsx validation failed: %s", exc)
        _xlsx_write_validation_failure(output_path, source_path, exc)
        Path(output_path).unlink(missing_ok=True)
        raise ValueError(f"XLSX failed openpyxl validation: {exc}") from exc

    # ── Validation 2: LibreOffice headless convert (if available) ────────────
    lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if lo_bin:
        try:
            with tempfile.TemporaryDirectory() as _tmpdir:
                result = subprocess.run(
                    [lo_bin, "--headless", "--convert-to", "xlsx",
                     "--outdir", _tmpdir, str(output_path)],
                    capture_output=True, text=True, timeout=90,
                )
            if result.returncode != 0:
                raise RuntimeError(
                    f"exit {result.returncode}: {result.stderr[:400]}"
                )
            log.info("doctuning xlsx: LibreOffice validation passed")
        except Exception as exc:
            log.warning("doctuning xlsx libreoffice validation failed: %s", exc)
            _xlsx_write_validation_failure(output_path, source_path, exc)
            Path(output_path).unlink(missing_ok=True)
            raise ValueError(f"XLSX failed LibreOffice validation: {exc}") from exc
    else:
        log.info("doctuning xlsx: libreoffice not found — skipping LO validation")

    return Path(output_path)


def _xlsx_write_validation_failure(output_path: "Path", source_path: "Path | None", exc: Exception) -> None:
    """Write a sidecar failure note to doctuning_debug/. Best-effort, not sent to user."""
    try:
        debug_dir = _doctuning_ws() / "training" / "doctuning_debug"
        debug_dir.mkdir(parents=True, exist_ok=True)
        sidecar = debug_dir / (Path(output_path).stem + "_validation_failure.txt")
        sidecar.write_text(
            f"XLSX validation failed\nsource: {source_path}\nerror: {exc}\n",
            encoding="utf-8",
        )
    except Exception:
        pass


async def _run_doctuning_continue_for_chat(
    bot,
    chat_id: int,
    state: dict,
    lang: str = "en",
) -> None:
    """Batch doctuning continuation: local rejected candidate + teacher audit → await_approval.

    Replicates the normal-VRAM path of _handle_docfill_file step 'await_example'.
    State must contain blank_text, filled_text, blank_name, content_hash, context_result.
    On completion sets state['step'] = 'await_approval' and keeps state in _PENDING_DOCFILL.
    """
    loop = asyncio.get_event_loop()
    blank_text = str(state.get("blank_text", ""))
    filled_text = str(state.get("filled_text", ""))
    blank_name = str(state.get("blank_name", "blank"))
    content_hash = str(state.get("content_hash", ""))
    ctx_res = state.get("context_result", {})
    preferred_fmt = state.get("preferred_output_format", "docx")

    # 1. Save context memo (informational — not the final training pair)
    try:
        memo_path = await loop.run_in_executor(None, _save_doctuning_memo, ctx_res, blank_name, content_hash)
        log.info("doctuning batch: context memo saved: %s", memo_path)
    except Exception as exc:
        log.warning("doctuning batch: context memo save failed: %s", exc)

    await bot.send_message(
        chat_id=chat_id,
        text="⏳ Generating local rejected candidate and running teacher audit in parallel…",
    )

    # 2. Parallel: local fill + teacher validation
    combined_result, (local_fill, repair_report) = await asyncio.gather(
        loop.run_in_executor(None, _validate_and_extract_context_sync, blank_text, filled_text),
        _local_fill_with_repair(blank_text),
    )
    verdict = combined_result.get("validation", {})
    audit_error = combined_result.get("error")
    h_score_raw = verdict.get("hallucination_score")
    h_score = float(h_score_raw) if h_score_raw is not None else None
    v_issues = verdict.get("issues", [])
    v_status = verdict.get("status", "")
    v_standards = verdict.get("standards_met", [])
    v_reason = str(verdict.get("reason", ""))
    audit_failed = (
        bool(audit_error)
        or v_status == "failed"
        or "teacher_audit_failed" in v_issues
        or h_score is None
    )
    if audit_failed:
        _fail_b = _save_doctuning_failure_case(
            stage="teacher_audit",
            reason=v_reason or audit_error or "audit status=failed",
            chat_id=chat_id,
            state=state,
            severity="error",
            recoverable=True,
        )
        await loop.run_in_executor(
            None, _request_repairman_investigation_for_doctuning_failure, _fail_b
        )

    # 3. Separate rejected (local draft) from chosen (teacher-revised or filled_text fallback)
    if local_fill:
        rejected_text = local_fill
    else:
        _local_fail_why = repair_report or "local generation failed"
        rejected_text = f"LOCAL_DRAFT_UNAVAILABLE: {_local_fail_why}\n\n{blank_text}"
        _fail_a = _save_doctuning_failure_case(
            stage="local_rejected_generation",
            reason=_local_fail_why,
            chat_id=chat_id,
            state=state,
            severity="warning",
            recoverable=True,
        )
        await loop.run_in_executor(
            None, _request_repairman_investigation_for_doctuning_failure, _fail_a
        )
    # _validate_and_extract_context_sync returns context+validation+revised_reference.
    # Check all plausible keys in case the schema is extended in the future.
    chosen_text = (
        combined_result.get("revised_reference")
        or combined_result.get("revised_text")
        or combined_result.get("chosen")
        or combined_result.get("candidate")
        or combined_result.get("reference_candidate")
        or combined_result.get("output")
        or None
    )
    teacher_revision_available = chosen_text is not None and not audit_failed
    if not teacher_revision_available:
        chosen_text = filled_text

    # 4. Build candidate file from chosen_text — xlsx if requested, else docx/txt fallback
    stem = Path(blank_name).stem or "candidate"
    candidate_path: Path | None = None
    xlsx_sent = False
    _xlsx_generation_error: str = ""

    if preferred_fmt == "xlsx":
        _blank_path_str = str(state.get("blank_path", "")).strip()
        _filled_path_str = str(state.get("filled_path", "")).strip()
        _template_path = Path(_blank_path_str) if _blank_path_str else None
        _reference_path = Path(_filled_path_str) if _filled_path_str else None
        _ctx = combined_result.get("context") or {}
        _val = combined_result.get("validation") or {}
        _out = Path(f"/tmp/doctuning_candidate_{chat_id}_{stem}.xlsx")
        try:
            _xlsx_result = await loop.run_in_executor(
                None,
                _create_doctuning_xlsx_candidate,
                _template_path if (_template_path and _template_path.exists()) else None,
                _reference_path if (_reference_path and _reference_path.exists()) else None,
                chosen_text,
                _ctx,
                _val,
                _out,
            )
            if _xlsx_result and _xlsx_result.exists():
                candidate_path = _xlsx_result
            else:
                _xlsx_generation_error = "generator returned no path"
                log.warning("doctuning xlsx generation failed; no DOCX substitute sent: generator returned no path")
        except Exception as _xlsx_exc:
            _xlsx_generation_error = str(_xlsx_exc)
            log.warning("doctuning xlsx generation failed; no DOCX substitute sent: %s", _xlsx_exc)

        if candidate_path is None:
            # Policy: do not substitute DOCX when user requested XLSX.
            # Save a debug artifact (not sent to user), report failure, mark state as failed.
            try:
                debug_dir = _doctuning_ws() / "training" / "doctuning_debug"
                debug_dir.mkdir(parents=True, exist_ok=True)
                debug_file = debug_dir / f"doctuning_debug_{chat_id}_{stem}.docx"
                try:
                    from docx import Document as _DocxDocument
                    _d = _DocxDocument()
                    _d.add_paragraph(chosen_text)
                    _d.save(str(debug_file))
                except Exception:
                    debug_file = debug_dir / f"doctuning_debug_{chat_id}_{stem}.txt"
                    debug_file.write_text(chosen_text or "", encoding="utf-8")
            except Exception as _de:
                log.warning("doctuning batch: debug artifact write failed: %s", _de)

            state.update({
                "step": "candidate_generation_failed",
                "candidate_path": "",
                "preferred_output_format": "xlsx",
                "xlsx_generation_error": _xlsx_generation_error,
                "local_fill": rejected_text,
                "rejected_candidate": rejected_text,
                "filled_text": filled_text,
                "chosen_candidate": chosen_text,
                "candidate_text": chosen_text,
                "repair_report": repair_report,
                "nim_score": h_score,
                "nim_issues": v_issues,
                "nim_standards": v_standards,
                "context_combined": combined_result.get("context") or {},
                "validation_combined": combined_result.get("validation") or {},
            })
            _PENDING_DOCFILL[chat_id] = state
            _fail_c = _save_doctuning_failure_case(
                stage="xlsx_candidate_generation",
                reason=_xlsx_generation_error or "xlsx candidate is None after generation",
                chat_id=chat_id,
                state=state,
                severity="error",
                recoverable=True,
            )
            await loop.run_in_executor(
                None, _request_repairman_investigation_for_doctuning_failure, _fail_c
            )
            await bot.send_message(
                chat_id=chat_id,
                text=(
                    "XLSX candidate generation failed validation. "
                    "I did not send a DOCX substitute because the requested output format is Excel. "
                    "The task is kept for repair."
                ),
            )
            return

    else:
        # DOCX fallback (only when preferred_fmt is not xlsx)
        try:
            from docx import Document as _DocxDocument
            doc_out = _DocxDocument()
            doc_out.add_paragraph(chosen_text)
            candidate_path = Path(f"/tmp/doctuning_candidate_{chat_id}_{stem}.docx")
            doc_out.save(str(candidate_path))
        except Exception as exc:
            log.warning("doctuning batch: docx creation failed (%s) — using .txt fallback", exc)
            try:
                candidate_path = Path(f"/tmp/doctuning_candidate_{chat_id}_{stem}.txt")
                candidate_path.write_text(chosen_text or "", encoding="utf-8")
            except Exception as exc2:
                log.error("doctuning batch: txt fallback also failed: %s", exc2)
                candidate_path = None

    # 5. Build caption and send candidate file
    caption_parts = ["📄 Revised reference candidate"]
    if not teacher_revision_available:
        caption_parts.append(
            "Teacher audit failed or revision output was unavailable; sending supplied reference as fallback candidate."
        )
    if repair_report:
        caption_parts.append(f"⚠️ Repairman note: {repair_report[:200]}")
    if candidate_path and candidate_path.suffix.lower() == ".xlsx":
        caption_parts.append(
            "Revised reference candidate is attached as XLSX. "
            "The incoming workbook was used as the base and AIMS review notes were added in a separate sheet."
        )
    caption = "\n".join(caption_parts)

    if candidate_path and candidate_path.exists():
        try:
            with open(candidate_path, "rb") as fh:
                await bot.send_document(
                    chat_id=chat_id,
                    document=fh,
                    filename=candidate_path.name,
                    caption=caption[:1024],
                )
            if candidate_path.suffix.lower() == ".xlsx":
                xlsx_sent = True
        except Exception as exc:
            log.warning("doctuning batch: send_document failed: %s", exc)
            await bot.send_message(
                chat_id=chat_id,
                text="Candidate was generated but could not be sent. Please check logs.",
            )
        finally:
            # Keep XLSX on disk — the approval step will resend and then unlink it.
            if candidate_path and candidate_path.suffix.lower() != ".xlsx":
                try:
                    candidate_path.unlink(missing_ok=True)
                except Exception:
                    pass
    else:
        await bot.send_message(
            chat_id=chat_id,
            text="Candidate was generated but could not be sent. Please check logs.",
        )

    # 6. Update state → await_approval
    state.update({
        "step": "await_approval",
        "local_fill": rejected_text,
        "rejected_candidate": rejected_text,
        "filled_text": filled_text,
        "chosen_candidate": chosen_text,
        "candidate_text": chosen_text,
        "candidate_path": str(candidate_path) if candidate_path else "",
        "repair_report": repair_report,
        "nim_score": h_score,
        "nim_issues": v_issues,
        "nim_standards": v_standards,
        "context_combined": combined_result.get("context") or {},
        "validation_combined": combined_result.get("validation") or {},
    })
    _PENDING_DOCFILL[chat_id] = state

    # 7. Show audit results + approval prompt
    standards_text = ", ".join(v_standards) if v_standards else "—"
    issues_lines = "\n".join(f"  {i+1}. {issue}" for i, issue in enumerate(v_issues))

    if audit_failed:
        nim_block = (
            "⚠️ Teacher audit failed\n"
            f"Reason: {_tg_plain(v_reason or audit_error or 'unknown')}\n\n"
            "Revised reference was not generated by the teacher model.\n"
            "You may still approve the fallback candidate or cancel.\n"
            "Approve: /doctuning_approve\n"
            "Cancel: /doctuning_cancel"
        )
    elif v_issues:
        nim_block = (
            f"📋 Omnirouter (Claude) — reference audit:\n"
            f"Standards: {_tg_plain(standards_text)}\n"
            f"Hallucination score: {h_score:.2f}\n\n"
            f"Recommendations:\n{_tg_plain(issues_lines)}\n\n"
            f"{_tg_plain(v_reason)}\n\n"
            "Consider these points before approving.\n"
            "Approve as master document: /doctuning_approve\n"
            "Cancel: /doctuning_cancel"
        )
    else:
        nim_block = (
            f"✅ Omnirouter (Claude) — no issues found\n"
            f"Standards: {_tg_plain(standards_text)}\n"
            f"Hallucination score: {h_score:.2f} — {_tg_plain(v_reason)}\n\n"
            "Approve as master document: /doctuning_approve\n"
            "Cancel: /doctuning_cancel"
        )

    _dt_keyboard = InlineKeyboardMarkup([[
        InlineKeyboardButton("✅ Approve master document", callback_data="dt_approve"),
        InlineKeyboardButton("❌ Cancel", callback_data="dt_cancel"),
    ]])
    await _safe_send_message(bot, chat_id, nim_block, reply_markup=_dt_keyboard)


def _docfill_save_dpo_pair(blank_text: str, chosen_text: str, rejected_text: str, doc_type: str) -> None:
    """Save doctuning DPO pair (human expert chosen, local model rejected)."""
    ws = Path(__file__).resolve().parent.parent
    dpo_path = ws / "data/training/standard_dpo_pairs.jsonl"
    dpo_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        with dpo_path.open("a", encoding="utf-8") as f:
            f.write(json.dumps({
                "prompt": blank_text,
                "chosen": chosen_text,
                "chosen_score": 1.0,
                "rejected": rejected_text,
                "rejected_score": 0.5,
                "rejected_feedback": "Local model fill — to be improved toward human expert reference",
                "pipeline": f"doctuning_{doc_type}",
                "target_model": AIMS_DOCUMENT_CREATOR_MODEL,
                "target_skill": "document_structure_generation",
                "rejected_model": AIMS_DOCUMENT_CREATOR_MODEL,
                "teacher_model": AIMS_DOCUMENT_TEACHER_MODEL,
                "audit_model": AIMS_DOCUMENT_AUDIT_MODEL,
                "teacher_models": {
                    "standards_discovery": AIMS_DOCUMENT_TEACHER_MODEL,
                    "audit": AIMS_DOCUMENT_AUDIT_MODEL,
                    "quality_judge": AIMS_DOCUMENT_TEACHER_MODEL,
                },
            }, ensure_ascii=False) + "\n")
        log.info("doctuning: DPO pair saved to standard_dpo_pairs.jsonl")
    except Exception as e:
        log.warning("doctuning dpo pair save failed: %s", e)


def _docfill_save_training_pair(blank_text: str, filled_text: str, doc_type: str, source_name: str) -> Path:
    """Append a (blank → filled) training pair to the docfill dataset."""
    import sys as _sys
    ws = Path(__file__).resolve().parent.parent
    out_dir = ws / "ops/ft/data/docfill_v1"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / "train_docfill_v1.jsonl"

    system_msg = (
        "You are an AIMS document specialist. When given a blank technical form, "
        "fill in all fields accurately for the specified equipment type based on "
        "industrial maintenance and safety standards."
    )
    pair = {
        "messages": [
            {"role": "system", "content": system_msg},
            {"role": "user", "content": f"Fill in this form:\n\n{blank_text}"},
            {"role": "assistant", "content": filled_text},
        ],
        "_meta": {
            "source": "axi_docfill_upload",
            "doc_type": doc_type,
            "blank_name": source_name,
            "saved_at": datetime.now(timezone.utc).isoformat(),
            "target_model": AIMS_DOCUMENT_CREATOR_MODEL,
            "target_skill": "document_structure_generation",
            "rejected_model": AIMS_DOCUMENT_CREATOR_MODEL,
            "teacher_model": AIMS_DOCUMENT_TEACHER_MODEL,
            "audit_model": AIMS_DOCUMENT_AUDIT_MODEL,
            "teacher_models": {
                "standards_discovery": AIMS_DOCUMENT_TEACHER_MODEL,
                "audit": AIMS_DOCUMENT_AUDIT_MODEL,
                "quality_judge": AIMS_DOCUMENT_TEACHER_MODEL,
            },
        },
    }
    with out_file.open("a", encoding="utf-8") as f:
        f.write(json.dumps(pair, ensure_ascii=False) + "\n")
    return out_file


def _doc_creator_120b_save(
    *,
    pair_type: str,
    blank_text: str,
    filled_text: str,
    doc_type: str,
    source_name: str,
    extra_meta: dict | None = None,
) -> Path:
    """Save a training record to the 120B document-creator training pool.

    pair_type must be one of: template_fill_pairs, approved_master_pairs,
    audit_revision_pairs, failed_generation_cases, structure_pairs.
    """
    subdir = AIMS_DOC_CREATOR_TRAINING_DIR / pair_type
    subdir.mkdir(parents=True, exist_ok=True)
    out_file = subdir / f"train_{pair_type}.jsonl"
    meta: dict = {
        "source": source_name,
        "doc_type": doc_type,
        "pair_type": pair_type,
        "saved_at": datetime.now(timezone.utc).isoformat(),
        "target_model": AIMS_DOCUMENT_CREATOR_MODEL,
        "target_skill": "document_structure_generation",
        "teacher_model": AIMS_DOCUMENT_TEACHER_MODEL,
        "audit_model": AIMS_DOCUMENT_AUDIT_MODEL,
    }
    if extra_meta:
        meta.update(extra_meta)
    record = {
        "messages": [
            {
                "role": "system",
                "content": (
                    "You are an AIMS document specialist. When given a blank technical form, "
                    "fill in all fields accurately for the specified equipment type based on "
                    "industrial maintenance and safety standards."
                ),
            },
            {"role": "user", "content": f"Fill in this form:\n\n{blank_text}"},
            {"role": "assistant", "content": filled_text},
        ],
        "_meta": meta,
    }
    with out_file.open("a", encoding="utf-8") as f:
        f.write(json.dumps(record, ensure_ascii=False) + "\n")
    return out_file


def _doc_creator_120b_save_failed(
    *,
    blank_text: str,
    doc_type: str,
    source_name: str,
    repair_report: str,
) -> None:
    """Record a failed 120B generation attempt for later analysis."""
    subdir = AIMS_DOC_CREATOR_TRAINING_DIR / "failed_generation_cases"
    subdir.mkdir(parents=True, exist_ok=True)
    out_file = subdir / "failed_generation_cases.jsonl"
    record = {
        "blank_text": blank_text[:2000],
        "doc_type": doc_type,
        "source_name": source_name,
        "repair_report": repair_report,
        "saved_at": datetime.now(timezone.utc).isoformat(),
        "target_model": AIMS_DOCUMENT_CREATOR_MODEL,
    }
    try:
        with out_file.open("a", encoding="utf-8") as f:
            f.write(json.dumps(record, ensure_ascii=False) + "\n")
    except Exception as e:
        log.warning("doc_creator_120b failed-case save error: %s", e)


def _docfill_save_to_omi(title: str, content: str, doc_type: str) -> str | None:
    """Optional future path: register master document via OmiAgent POST /documents.

    Only call this when a real /documents endpoint is confirmed to exist.
    As of 2026-05: omi-api returns 404 on /documents — use _queue_to_omi_batch_inbox instead.
    Returns doc_id string on success.
    Raises RuntimeError on HTTP/network/parse failure.
    """
    import uuid as _uuid
    import urllib.request as _ur
    import urllib.error as _ue

    base_url = (
        os.environ.get("OMI_API_URL")
        or os.environ.get("OMI_URL")
        or os.environ.get("AIMS_OMI_URL")
        or "http://localhost:8008"
    ).rstrip("/")
    url = f"{base_url}/documents"
    task_id = f"doctuning-{_uuid.uuid4().hex[:12]}"

    payload = json.dumps({
        "task_id": task_id,
        "title": title,
        "doc_type": doc_type,
        "status": "master",
        "quality_score": 0.0,
        "summary": (content or "")[:500],
        "metadata": {
            "source": "doctuning_batch_training_pair",
            "created_at": datetime.now(timezone.utc).isoformat(),
        },
    }).encode()

    log.info("docfill_save_to_omi: POST %s task_id=%s title=%r", url, task_id, title[:80])
    try:
        req = _ur.Request(
            url,
            data=payload,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with _ur.urlopen(req, timeout=15) as resp:
            status_code = resp.status
            body = resp.read().decode("utf-8", errors="replace")
            log.info("docfill_save_to_omi: HTTP %s body=%r", status_code, body[:500])
            result = json.loads(body)
            doc_id = (
                result.get("doc_id")
                or result.get("id")
                or result.get("document_id")
            )
            if doc_id:
                return str(doc_id)
            raise ValueError(f"OmiAgent returned no id field: body={body[:300]!r}")
    except _ue.HTTPError as http_exc:
        body_text = (http_exc.read() or b"").decode("utf-8", errors="replace")[:500]
        raise RuntimeError(
            f"OmiAgent HTTP {http_exc.code}: url={url} body={body_text!r}"
        ) from http_exc


async def _handle_docfill_file(
    update: Update,
    ctx: ContextTypes.DEFAULT_TYPE,
    chat_id: int,
) -> None:
    """Handle file upload in docfill mode."""
    state = _PENDING_DOCFILL[chat_id]
    step = str(state.get("step", ""))
    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "ru"

    if not update.message or not update.message.document:
        await update.message.reply_text(
            "Пожалуйста, отправьте файл (.txt, .json, .docx, .xlsx)."
            if lang == "ru" else
            "Please send a file (.txt, .json, .docx, .xlsx)."
        )
        return

    doc = update.message.document
    file_name = doc.file_name or f"upload_{doc.file_id[-8:]}"
    suffix = Path(file_name).suffix.lower() or ".txt"

    # Validate file type
    if suffix not in (".txt", ".json", ".docx", ".xlsx"):
        await update.message.reply_text(
            f"❌ Неподдерживаемый формат: `{suffix}`\n"
            "Поддерживаются: .txt, .json, .docx, .xlsx"
            if lang == "ru" else
            f"❌ Unsupported format: `{suffix}`\n"
            "Supported: .txt, .json, .docx, .xlsx",
            parse_mode="Markdown",
        )
        return

    with tempfile.NamedTemporaryFile(prefix="axi_docfill_", suffix=suffix, delete=False) as tmp:
        tmp_path = Path(tmp.name)
    tg_file = await ctx.bot.get_file(doc.file_id)
    await tg_file.download_to_drive(str(tmp_path))

    if step == "await_blank":
        state["blank_path"] = str(tmp_path)
        state["blank_name"] = file_name
        state["blank_suffix"] = suffix
        state["step"] = "await_example"
        _PENDING_DOCFILL[chat_id] = state
        await update.message.reply_text(
            f"✅ Шаблон принят: {_tg_plain(file_name)}\n\n"
            "Теперь загрузите эталонный заполненный документ (файл 2 из 2)."
            if lang == "ru" else
            f"✅ Template accepted: {_tg_plain(file_name)}\n\n"
            "Now send the filled reference document (file 2 of 2).",
        )
        return

    if step == "await_example":
        blank_path = Path(str(state.get("blank_path", "")))
        blank_name = str(state.get("blank_name", "blank"))
        doc_type = str(state.get("doc_type", "procedure"))
        filled_name = Path(file_name).stem
        _blank_suffix = state.get("blank_suffix", Path(blank_name).suffix.lower())
        state["filled_suffix"] = suffix
        state["preferred_output_format"] = "xlsx" if _blank_suffix == ".xlsx" or suffix == ".xlsx" else "docx"
        state["step"] = "validating"
        _PENDING_DOCFILL[chat_id] = state

        try:
            blank_text = _read_document_text(blank_path)
            filled_text = _read_document_text(tmp_path)
        except Exception as e:
            await update.message.reply_text(
                f"Ошибка чтения файла: {e}" if lang == "ru" else f"File read error: {e}"
            )
            _PENDING_DOCFILL.pop(chat_id, None)
            tmp_path.unlink(missing_ok=True)
            blank_path.unlink(missing_ok=True)
            return

        loop = asyncio.get_event_loop()

        # ── Preflight: check Ollama availability and VRAM ────────────────────
        ollama_alive = await loop.run_in_executor(None, _check_ollama_alive_sync)
        vram_free_gb = await loop.run_in_executor(None, _check_vram_free_gb_sync)
        model = AIMS_DOCUMENT_CREATOR_MODEL

        if not ollama_alive:
            await update.message.reply_text(
                f"⚠️ Ollama недоступна. Запускаю репairman для восстановления сервиса…\n"
                f"NIM-валидация вашего эталона запущена параллельно."
                if lang == "ru" else
                f"⚠️ Ollama is down. Launching repairman to restore the service…\n"
                f"NIM validation of your reference is running in parallel."
            )
            # Repairman restores ollama; NIM path runs independently below
            repair_task = loop.run_in_executor(None, _repairman_fix_local_model_sync, model)
        else:
            repair_task = None

        if ollama_alive and vram_free_gb < DOCFILL_VRAM_MIN_FREE_GB:
            await update.message.reply_text(
                f"⏸ VRAM перегружена ({vram_free_gb:.1f} GB свободно, нужно ≥{DOCFILL_VRAM_MIN_FREE_GB:.0f} GB).\n"
                f"Задача локальной генерации поставлена в очередь — уведомлю, когда освободится.\n"
                f"Omnirouter (Claude) валидация эталона запущена."
                if lang == "ru" else
                f"⏸ VRAM is full ({vram_free_gb:.1f} GB free, need ≥{DOCFILL_VRAM_MIN_FREE_GB:.0f} GB).\n"
                f"Local generation queued — I'll notify you when VRAM is available.\n"
                f"Omnirouter (Claude) validation is running."
            )
            _DOCFILL_VRAM_QUEUE[chat_id] = blank_text
            # Combined context extraction + validation via Omnirouter
            combined_result = await loop.run_in_executor(
                None, _validate_and_extract_context_sync, blank_text, filled_text
            )
            context_vram = combined_result.get("context", {})
            verdict = combined_result.get("validation", {})
            local_fill = None
            repair_report = ""
        else:
            await update.message.reply_text(
                f"✅ Эталон принят: {_tg_plain(file_name)}\n\n"
                "⏳ Запускаю два параллельных пути:\n"
                "  • Локальная модель заполняет шаблон (rejected)\n"
                "  • Omnirouter (Claude) проверяет ваш эталон по стандартам…"
                if lang == "ru" else
                f"✅ Reference accepted: {_tg_plain(file_name)}\n\n"
                "⏳ Running two parallel paths:\n"
                "  • Local model fills the blank (rejected candidate)\n"
                "  • Omnirouter (Claude) audits your reference against standards…",
            )
            # Wait for repair task to complete (if repairman was launched for ollama-down)
            if repair_task is not None:
                await repair_task
                await asyncio.sleep(20)  # give ollama time to stabilise

            # Combined context extraction + validation via Omnirouter, parallel with local fill
            combined_result, (local_fill, repair_report) = await asyncio.gather(
                loop.run_in_executor(None, _validate_and_extract_context_sync, blank_text, filled_text),
                _local_fill_with_repair(blank_text),
            )
            context = combined_result.get("context", {})
            verdict = combined_result.get("validation", {})

            # Save context memo
            if context:
                memo_path_ctx = _save_doctuning_memo(context, blank_name)
                memo_path_ctx.unlink(missing_ok=True)

        h_score = float(verdict.get("hallucination_score", 0.0))
        v_issues = verdict.get("issues", [])
        v_standards = verdict.get("standards_met", [])
        v_reason = str(verdict.get("reason", ""))

        standards_text = ", ".join(v_standards) if v_standards else "—"
        issues_lines = "\n".join(f"  {i+1}. {issue}" for i, issue in enumerate(v_issues))

        # Save approved state — master doc saved only after /doctuning_approve
        state.update({
            "step": "await_approval",
            "blank_text": blank_text,
            "filled_text": filled_text,
            "filled_name": filled_name,
            "local_fill": local_fill,
            "repair_report": repair_report,
            "nim_score": h_score,
            "nim_issues": v_issues,
            "nim_standards": v_standards,
            "blank_path": str(blank_path),  # Keep for cleanup
            "filled_path": str(tmp_path),   # Keep for cleanup
        })
        _PENDING_DOCFILL[chat_id] = state

        if v_issues:
            nim_block = (
                f"📋 Omnirouter (Claude) — замечания по эталону:\n"
                f"Стандарты: {_tg_plain(standards_text)}\n"
                f"Hallucination score: {h_score:.2f}\n\n"
                f"Замечания:\n{_tg_plain(issues_lines)}\n\n"
                f"{_tg_plain(v_reason)}\n\n"
                "Учтите замечания при финальном решении.\n"
                "Подтвердить как мастер-документ: /doctuning_approve\n"
                "Отмена: /doctuning_cancel"
                if lang == "ru" else
                f"📋 Omnirouter (Claude) — reference audit:\n"
                f"Standards: {_tg_plain(standards_text)}\n"
                f"Hallucination score: {h_score:.2f}\n\n"
                f"Recommendations:\n{_tg_plain(issues_lines)}\n\n"
                f"{_tg_plain(v_reason)}\n\n"
                "Consider these points before approving.\n"
                "Approve as master document: /doctuning_approve\n"
                "Cancel: /doctuning_cancel"
            )
        else:
            nim_block = (
                f"✅ Omnirouter (Claude) — замечаний нет\n"
                f"Стандарты: {_tg_plain(standards_text)}\n"
                f"Hallucination score: {h_score:.2f} — {_tg_plain(v_reason)}\n\n"
                "Подтвердить как мастер-документ: /doctuning_approve\n"
                "Отмена: /doctuning_cancel"
                if lang == "ru" else
                f"✅ Omnirouter (Claude) — no issues found\n"
                f"Standards: {_tg_plain(standards_text)}\n"
                f"Hallucination score: {h_score:.2f} — {_tg_plain(v_reason)}\n\n"
                "Approve as master document: /doctuning_approve\n"
                "Cancel: /doctuning_cancel"
            )

        await _safe_send_message(ctx.bot, chat_id, nim_block)
        # Don't delete tmp_path yet — keep for cleanup after approval
        return

    # Unknown state — reset
    _PENDING_DOCFILL.pop(chat_id, None)
    tmp_path.unlink(missing_ok=True)


async def cmd_doctuning(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Start docfill training pair creation: upload blank template + filled reference."""
    if not _chat_allowed(update):
        return
    if update.message is None or update.effective_chat is None:
        return
    chat_id = update.effective_chat.id
    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "ru"

    doc_type = " ".join(ctx.args).strip().lower() if ctx.args else "procedure"

    _PENDING_DOCFILL[chat_id] = {"step": "await_blank", "doc_type": doc_type}

    await update.message.reply_text(
        f"📄 **Режим создания тюнинг-пары** (тип: `{doc_type}`)\n\n"
        "Шаг 1 из 2 — загрузите **шаблон** (пустой бланк) в формате .txt, .json, .docx или .xlsx.\n\n"
        "Порядок:\n"
        "1. Загружаете шаблон\n"
        "2. Загружаете заполненный эталон\n"
        "3. Omnirouter (Claude) извлекает контекст и валидирует эталон на галлюцинации и соответствие стандартам\n"
        "4. Если проверка пройдена → пара сохраняется как мастер-документ и тюнинг-пара\n\n"
        "Отмена: `/doctuning_cancel`"
        if lang == "ru" else
        f"📄 **Training pair mode** (type: `{doc_type}`)\n\n"
        "Step 1 of 2 — upload the **blank template** (.txt, .json, .docx, or .xlsx).\n\n"
        "Flow:\n"
        "1. Upload blank template\n"
        "2. Upload filled reference\n"
        "3. Omnirouter (Claude) extracts context and validates for hallucinations and standards compliance\n"
        "4. Pair is always saved as master document + training pair\n\n"
        "Cancel: `/doctuning_cancel`",
        parse_mode="Markdown",
    )


async def cmd_doctuning_cancel(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not _chat_allowed(update) or update.message is None or update.effective_chat is None:
        return
    chat_id = update.effective_chat.id
    state = _PENDING_DOCFILL.pop(chat_id, None)
    if state:
        # Cleanup temporary files
        blank_path = state.get("blank_path")
        filled_path = state.get("filled_path")
        if blank_path:
            Path(str(blank_path)).unlink(missing_ok=True)
        if filled_path:
            Path(str(filled_path)).unlink(missing_ok=True)
    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "ru"
    await update.message.reply_text(
        "Режим doctuning отменён." if lang == "ru" else "Doctuning mode cancelled."
    )


async def cmd_docgen_upgrade(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Queue a DOCGEN self-improvement upgrade batch. Private Axi only."""
    if not _chat_allowed(update) or update.message is None:
        return
    if not _is_docgen_upgrade_origin(update):
        await update.message.reply_text(
            "⚠️ `/DOCGEN_UPGRADE` is private-only for Axi. "
            "Use `/docgen` for generation-only runs.",
            parse_mode="Markdown",
        )
        return

    try:
        from logi.docgen_skill_scope_policy import DOCGEN_DOCUMENT_TYPES  # noqa: PLC0415
        from docgen.universal_overlay.nightly_improvement_orchestrator import (  # noqa: PLC0415
            build_docgen_upgrade_batch_action,
            make_docgen_upgrade_run_tag,
            split_docgen_upgrade_request,
            run_docgen_upgrade_batch_now,
            stream_docgen_upgrade_progress,
        )
    except Exception as exc:  # pragma: no cover - import depends on runtime layout
        await update.message.reply_text(
            f"❌ DOCGEN upgrade launcher unavailable: {type(exc).__name__}",
            parse_mode="Markdown",
        )
        return

    topic_prefix, requested_types = split_docgen_upgrade_request(
        ctx.args or [],
        DOCGEN_DOCUMENT_TYPES,
    )
    doc_types = requested_types or list(DOCGEN_DOCUMENT_TYPES)
    run_tag = make_docgen_upgrade_run_tag()
    action_preview = build_docgen_upgrade_batch_action(
        document_types=doc_types,
        created_by="axi",
        self_improvement_enabled=True,
        max_cycles=0,
        topic_prefix=" ".join(topic_prefix).strip() or None,
        run_tag=run_tag,
    )

    progress_msg = await update.message.reply_text("⏳ DOCGEN upgrade started...", parse_mode="Markdown")
    chat_id = update.effective_chat.id
    stop_event = asyncio.Event()

    async def _run_upgrade_and_report() -> None:
        try:
            log.info("DOCGEN upgrade: starting batch run for %s", ",".join(doc_types))
            monitor_task = asyncio.create_task(
                stream_docgen_upgrade_progress(
                    action=action_preview,
                    document_types=doc_types,
                    chat_id=chat_id,
                    bot=ctx.bot,
                    stop_event=stop_event,
                )
            )
            action = await asyncio.to_thread(
                run_docgen_upgrade_batch_now,
                document_types=doc_types,
                created_by="axi",
                self_improvement_enabled=True,
                max_cycles=0,
                topic_prefix=" ".join(topic_prefix).strip() or None,
                run_tag=run_tag,
            )
            action_payload = action.get("action") or {}
            action_id = str(action_payload.get("action_id") or "")
            action_path = str(action.get("action_path") or "")
            execution = action.get("execution_result") or {}
            run_summary = action.get("run_summary") or {}
            summary = run_summary if isinstance(run_summary, dict) else {}
            final_msg = _format_docgen_upgrade_reply(
                action_id=action_id,
                action_path=action_path,
                execution=execution,
                run_summary=summary,
                doc_types=doc_types,
            )
            completion_msg = _format_docgen_upgrade_completion(
                action_id=action_id,
                run_summary=summary,
            )
            try:
                await progress_msg.edit_text(final_msg, parse_mode="Markdown")
            except Exception:
                pass
            await ctx.bot.send_message(chat_id=chat_id, text=completion_msg, parse_mode="Markdown")
            log.info("DOCGEN upgrade: completion message sent for %s", action_id or "unknown")
            stop_event.set()
            try:
                await monitor_task
            except Exception:
                pass
        except Exception as exc:
            stop_event.set()
            err = f"❌ DOCGEN upgrade failed: {type(exc).__name__}: {exc}"
            try:
                await progress_msg.edit_text(err, parse_mode="Markdown")
            except Exception:
                await ctx.bot.send_message(chat_id=chat_id, text=err, parse_mode="Markdown")

    await _run_upgrade_and_report()
    return


async def _job_docfill_vram_monitor(ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Periodic job: resume queued docfill local-fill tasks when VRAM is available."""
    if not _DOCFILL_VRAM_QUEUE:
        return

    loop = asyncio.get_event_loop()
    free_gb = await loop.run_in_executor(None, _check_vram_free_gb_sync)
    if free_gb < DOCFILL_VRAM_MIN_FREE_GB:
        log.debug("docfill vram monitor: %.1f GB free < %.0f GB threshold, waiting", free_gb, DOCFILL_VRAM_MIN_FREE_GB)
        return

    # Process one queued task per tick to avoid VRAM spike from concurrent fills
    chat_id = next(iter(_DOCFILL_VRAM_QUEUE))
    blank_text = _DOCFILL_VRAM_QUEUE.pop(chat_id)

    state = _PENDING_DOCFILL.get(chat_id)
    if not state or state.get("step") != "await_approval":
        # State was cleared (user cancelled) — discard
        log.info("docfill vram monitor: chat_id=%d state gone, discarding queued task", chat_id)
        return

    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "ru"

    local_fill, repair_report = await _local_fill_with_repair(blank_text)

    # Update state with the now-available local fill result
    state["local_fill"] = local_fill
    state["repair_report"] = repair_report
    _PENDING_DOCFILL[chat_id] = state

    if local_fill:
        # Done state — notify
        await ctx.bot.send_message(
            chat_id,
            "✅ DPO-пара готова. Локальная генерация завершена.\n"
            "Подтвердите: `/doctuning_approve`"
            if lang == "ru" else
            "✅ DPO pair ready. Local fill complete.\n"
            "Confirm with: `/doctuning_approve`",
            parse_mode="Markdown",
        )
    else:
        # Full failure state — notify with repairman report
        report_suffix = f"\n\n_{repair_report}_" if repair_report else ""
        await ctx.bot.send_message(
            chat_id,
            f"🔴 Repairman исчерпал 5 циклов ремонта — локальная модель недоступна. DPO-пара будет пропущена.{report_suffix}\n\n"
            "SFT-пара и мастер-документ доступны: `/doctuning_approve`"
            if lang == "ru" else
            f"🔴 Repairman exhausted 5 repair cycles — local model unavailable. DPO pair will be skipped.{report_suffix}\n\n"
            "SFT pair and master document are ready: `/doctuning_approve`",
            parse_mode="Markdown",
        )


async def _approve_doctuning_for_chat(
    bot,
    chat_id: int,
    state: dict,
    lang: str = "en",
) -> None:
    """Shared approval handler for batch_training_pair source.

    Works from both /doctuning_approve command and the dt_approve inline callback.
    Uses bot.send_message / bot.send_document so it is not tied to a specific Update.
    """
    loop = asyncio.get_event_loop()

    blank_text    = str(state.get("blank_text", ""))
    filled_text   = str(state.get("filled_text", ""))
    blank_name    = str(state.get("blank_name", "blank"))
    doc_type      = str(state.get("doc_type", "procedure"))
    nim_issues    = state.get("nim_issues") or []
    nim_standards = state.get("nim_standards") or []
    nim_score     = state.get("nim_score")
    ctx_result    = state.get("context_result") or {}
    ctx_combined  = state.get("context_combined") or {}
    val_combined  = state.get("validation_combined") or {}
    preferred_output_format = str(state.get("preferred_output_format", "docx")).lower()
    candidate_path_raw = str(state.get("candidate_path", "") or "")
    candidate_path = Path(candidate_path_raw) if candidate_path_raw else None

    chosen_text = (
        state.get("chosen_candidate")
        or state.get("candidate_text")
        or filled_text
    )
    rejected_text = (
        state.get("rejected_candidate")
        or state.get("local_fill")
        or f"LOCAL_DRAFT_UNAVAILABLE: approval path fallback\n\n{blank_text}"
    )

    await bot.send_message(
        chat_id=chat_id,
        text="⏳ Saving training pair and queuing master document for Omi ingestion…",
    )

    # 1. Queue master document to Omi batch inbox (primary path).
    #    POST /documents is not available on current omi-api; inbox is the correct handoff.
    doc_title = f"{doc_type.title()} — {blank_name}"
    omi_id: str | None = None
    master_doc_status: str = "not_attempted"
    master_doc_path: str = "local_only"
    master_doc_error: str = ""

    log.info("batch_training_pair: queuing master document to Omi batch inbox: title=%r", doc_title)
    try:
        import functools as _functools
        _inbox_ts = datetime.now(timezone.utc)
        _inbox_name = (
            f"doctuning_{chat_id}_{_inbox_ts.strftime('%Y%m%d_%H%M%S')}.txt"
        )
        _omi_queue_tmp_dir = _doctuning_ws() / "training" / "doctuning_omi_queue_tmp"
        _omi_queue_tmp_dir.mkdir(parents=True, exist_ok=True)
        _inbox_tmp = _omi_queue_tmp_dir / _inbox_name
        _inbox_tmp.write_text(chosen_text or "", encoding="utf-8")
        _inbox_meta = {
            "source": "doctuning_batch_training_pair",
            "title": doc_title,
            "doc_type": doc_type,
            "chat_id": chat_id,
            "created_at": _inbox_ts.isoformat(),
            "pair_path": None,
            "status": "master",
            "master_document_status": "queued_for_omi_ingestion",
            "ingest_intent": "master_document_registration",
        }
        _queued_path = await loop.run_in_executor(
            None,
            _functools.partial(
                _queue_to_omi_batch_inbox, _inbox_tmp, _inbox_name, meta=_inbox_meta
            ),
        )
        master_doc_status = "queued_for_omi_ingestion"
        master_doc_path = f"omi_inbox:{_queued_path}"
        log.info(
            "batch_training_pair: master doc queued to omi inbox: %s queue_path=%s",
            _inbox_name, _queued_path,
        )
    except Exception as _q_exc:
        master_doc_status = "omi_failed"
        master_doc_path = "local_only"
        master_doc_error = str(_q_exc)
        log.error("batch_training_pair: omi inbox queue failed: %s", _q_exc, exc_info=True)

    # 2. Save JSON training pair (must not be blocked by Omi failure)
    now_ts = datetime.now(timezone.utc)
    pair_dir = _doctuning_ws() / "training" / "doctuning_pairs"
    pair_dir.mkdir(parents=True, exist_ok=True)
    pair_fname = f"{now_ts.strftime('%Y%m%d_%H%M%S')}_{chat_id}.json"
    pair_path = pair_dir / pair_fname
    _ctx_for_pair = ctx_combined or ctx_result or {}
    pair_data = {
        "source": "batch_training_pair",
        "created_at": now_ts.isoformat(),
        "approved_at": now_ts.isoformat(),
        "input": {
            "template": blank_text,
            "context": _ctx_for_pair,
            "standards": nim_standards,
        },
        "chosen": chosen_text,
        "rejected": rejected_text,
        "context": _ctx_for_pair,
        "validation": val_combined or {
            "hallucination_score": nim_score,
            "issues": nim_issues,
            "standards_met": nim_standards,
        },
        "candidate_path": str(pair_path),
        "master_document_status": master_doc_status,
        "master_document_path": master_doc_path,
        "master_document_error": master_doc_error,
    }
    pair_saved = False
    pair_save_error: str | None = None
    try:
        pair_path.write_text(json.dumps(pair_data, indent=2, ensure_ascii=False), encoding="utf-8")
        if not pair_path.exists() or pair_path.stat().st_size == 0:
            raise ValueError("pair file missing or empty after write")
        json.loads(pair_path.read_text(encoding="utf-8"))
        pair_saved = True
        log.info("batch_training_pair: pair saved verified: %s (%d bytes)", pair_path, pair_path.stat().st_size)
    except Exception as _e:
        pair_save_error = str(_e)
        log.error("batch_training_pair: pair save failed: %s | path=%s", _e, pair_path, exc_info=True)

    # Point D — Omi registration failure case
    if master_doc_status == "omi_failed":
        _fail_d = _save_doctuning_failure_case(
            stage="omi_registration",
            reason=master_doc_error or "OmiAgent registration failed",
            chat_id=chat_id,
            state=state,
            pair_path=str(pair_path) if pair_saved else None,
            severity="warning" if pair_saved else "error",
            recoverable=True,
        )
        await asyncio.get_event_loop().run_in_executor(
            None, _request_repairman_investigation_for_doctuning_failure, _fail_d
        )

    # Point E — training pair save failure case
    if not pair_saved:
        _fail_e = _save_doctuning_failure_case(
            stage="training_pair_save",
            reason=pair_save_error or "pair save returned False",
            chat_id=chat_id,
            state=state,
            severity="critical",
            recoverable=True,
        )
        await asyncio.get_event_loop().run_in_executor(
            None, _request_repairman_investigation_for_doctuning_failure, _fail_e
        )

    # 3. Send approved document file back to user
    stem = Path(blank_name).stem or "approved"
    approved_path: Path | None = None
    file_sent = False
    log.info(
        "batch_training_pair: approved output format=%s candidate_path=%s",
        preferred_output_format,
        candidate_path,
    )

    if preferred_output_format == "xlsx":
        if (
            candidate_path is not None
            and candidate_path.exists()
            and candidate_path.suffix.lower() == ".xlsx"
        ):
            approved_path = candidate_path
            log.info("batch_training_pair: sending approved xlsx=%s", approved_path)
            try:
                with open(approved_path, "rb") as _fh:
                    await bot.send_document(
                        chat_id=chat_id,
                        document=_fh,
                        filename=approved_path.name,
                        caption="✅ Approved master document",
                    )
                file_sent = True
            except Exception as _e:
                log.warning("batch_training_pair: send approved xlsx failed: %s", _e)
                _fail_f_xlsx = _save_doctuning_failure_case(
                    stage="telegram_send_document",
                    reason=str(_e),
                    chat_id=chat_id,
                    state=state,
                    pair_path=str(pair_path) if pair_saved else None,
                    severity="error",
                    recoverable=True,
                )
                await asyncio.get_event_loop().run_in_executor(
                    None, _request_repairman_investigation_for_doctuning_failure, _fail_f_xlsx
                )
                await bot.send_message(
                    chat_id=chat_id,
                    text="Candidate was generated but could not be sent. Please check logs.",
                )
        else:
            await bot.send_message(
                chat_id=chat_id,
                text=(
                    "Approved training pair was saved, but the approved XLSX document "
                    "could not be resent because the candidate file is missing."
                ),
            )
    else:
        try:
            from docx import Document as _DocxDocument
            doc_out = _DocxDocument()
            doc_out.add_paragraph(chosen_text)
            approved_path = Path(f"/tmp/doctuning_approved_{chat_id}_{stem}.docx")
            doc_out.save(str(approved_path))
        except Exception:
            try:
                approved_path = Path(f"/tmp/doctuning_approved_{chat_id}_{stem}.txt")
                approved_path.write_text(chosen_text, encoding="utf-8")
            except Exception as _e2:
                log.error("batch_training_pair: approved file creation failed: %s", _e2)
                approved_path = None

        if approved_path and approved_path.exists():
            try:
                with open(approved_path, "rb") as _fh:
                    await bot.send_document(
                        chat_id=chat_id,
                        document=_fh,
                        filename=approved_path.name,
                        caption="✅ Approved master document",
                    )
                file_sent = True
            except Exception as _e:
                log.warning("batch_training_pair: send approved doc failed: %s", _e)
                _fail_f_doc = _save_doctuning_failure_case(
                    stage="telegram_send_document",
                    reason=str(_e),
                    chat_id=chat_id,
                    state=state,
                    pair_path=str(pair_path) if pair_saved else None,
                    severity="error",
                    recoverable=True,
                )
                await asyncio.get_event_loop().run_in_executor(
                    None, _request_repairman_investigation_for_doctuning_failure, _fail_f_doc
                )
                await bot.send_message(
                    chat_id=chat_id,
                    text="Candidate was generated but could not be sent. Please check logs.",
                )
            finally:
                try:
                    approved_path.unlink(missing_ok=True)
                except Exception:
                    pass

    # 4. Cleanup and confirmation
    _PENDING_DOCFILL.pop(chat_id, None)

    lines = ["✅ Approved and saved (batch training pair):"]
    if pair_saved:
        lines.append(f"  Training pair: {_tg_plain(pair_path.resolve())}")
        lines.append(f"  training_pair_saved:{_tg_plain(pair_path.resolve())}")
    else:
        lines.append("  ⚠️ Training pair was not saved. See logs.")

    if master_doc_status == "registered":
        lines.append(f"  Master document registered in OmiAgent (id={_tg_plain(omi_id)}): {_tg_plain(doc_title)}")
    elif master_doc_status == "queued_for_omi_ingestion":
        lines.append(f"  Master document queued for Omi ingestion: {_tg_plain(master_doc_path)}")
    elif master_doc_status == "omi_failed":
        lines.append("  ⚠️ OmiAgent registration failed and a repair case was queued.")

    if not file_sent and approved_path is None:
        lines.append("  ⚠️ Approved document file could not be generated.")

    await _safe_send_message(bot, chat_id, "\n".join(lines))


async def _doctuning_callback_handler(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle dt_approve and dt_cancel inline button callbacks."""
    query = update.callback_query
    if query is None:
        return
    await query.answer()

    chat_id = query.message.chat_id if query.message else None
    if chat_id is None:
        return

    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "ru"

    if query.data == "dt_approve":
        state = _PENDING_DOCFILL.get(chat_id)
        if not state or state.get("step") != "await_approval":
            await ctx.bot.send_message(
                chat_id=chat_id,
                text="No pending approval. Start with /doctuning.",
            )
            return
        if chat_id in _DOCFILL_VRAM_QUEUE:
            await ctx.bot.send_message(
                chat_id=chat_id,
                text=(
                    "⏳ Local generation is still queued (waiting for VRAM).\n"
                    "I'll notify you when it's done — approve then.\n"
                    "Or /doctuning_cancel to cancel (DPO pair will be skipped)."
                ),
            )
            return
        if state.get("source") == "batch_training_pair":
            await _approve_doctuning_for_chat(ctx.bot, chat_id, state, lang)
        else:
            await ctx.bot.send_message(
                chat_id=chat_id,
                text="Use /doctuning_approve to approve this document type.",
            )

    elif query.data == "dt_cancel":
        state = _PENDING_DOCFILL.pop(chat_id, None)
        if state:
            await ctx.bot.send_message(chat_id=chat_id, text="❌ Doctuning cancelled.")
        else:
            await ctx.bot.send_message(chat_id=chat_id, text="No active doctuning session to cancel.")


async def _handle_docreview_file(
    update: Update, ctx: ContextTypes.DEFAULT_TYPE, chat_id: int
) -> None:
    """Download the uploaded doc for a single-document review session."""
    if update.message is None or update.message.document is None:
        return
    doc = update.message.document
    file_name = doc.file_name or f"upload_{doc.file_id[-8:]}"
    staging_dir = _DI_INTAKE_DIR / "staging" / str(chat_id)
    staging_dir.mkdir(parents=True, exist_ok=True)
    local_path = staging_dir / file_name
    tg_file = await ctx.bot.get_file(doc.file_id)
    await tg_file.download_to_drive(str(local_path))

    state = _PENDING_DOCREVIEW[chat_id]
    state.setdefault("files", []).append({
        "filename": file_name,
        "original_name": file_name,
        "local_path": str(local_path),
        "file_id": doc.file_id,
    })
    state["step"] = "ready"

    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("Evaluate", callback_data="dr_standards_yes"),
        InlineKeyboardButton("No", callback_data="dr_standards_no"),
    ]])
    await update.message.reply_text(
        "Do you want to verify the document against standards?",
        reply_markup=kb,
    )


async def _di_start_session(
    update: Update, ctx: ContextTypes.DEFAULT_TYPE, chat_id: int
) -> None:
    """Start a new doc intake session: download the first document and ask for more."""
    if update.message is None or update.message.document is None:
        return
    doc = update.message.document
    file_name = doc.file_name or f"upload_{doc.file_id[-8:]}"
    intake_session_id = __import__("uuid").uuid4().hex[:10]
    prompt_version = 1

    # Use session-specific staging, not chat-wide staging, so old test uploads
    # cannot leak into a new document package.
    staging_dir = _DI_INTAKE_DIR / "staging" / intake_session_id
    staging_dir.mkdir(parents=True, exist_ok=True)
    local_path = staging_dir / file_name
    tg_file = await ctx.bot.get_file(doc.file_id)
    await tg_file.download_to_drive(str(local_path))
    _PENDING_DOC_INTAKE[chat_id] = {
        "step": "await_more",
        "intake_session_id": intake_session_id,
        "prompt_version": prompt_version,
        "files": [{
            "filename": file_name,
            "original_name": file_name,
            "local_path": str(local_path),
            "file_id": doc.file_id,
        }],
    }

    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("Yes", callback_data=f"di_more:{intake_session_id}:{prompt_version}"),
        InlineKeyboardButton("No", callback_data=f"di_enough:{intake_session_id}:{prompt_version}"),
    ]])
    sent_msg = await update.message.reply_text(
        f"Received: {file_name}\nAny more documents?",
        reply_markup=kb,
    )
    _PENDING_DOC_INTAKE[chat_id]["last_prompt_message_id"] = sent_msg.message_id


async def _di_add_document(
    update: Update, ctx: ContextTypes.DEFAULT_TYPE, chat_id: int
) -> None:
    """Add another document to an active intake session."""
    if update.message is None or update.message.document is None:
        return
    doc = update.message.document
    file_name = doc.file_name or f"upload_{doc.file_id[-8:]}"

    state = _PENDING_DOC_INTAKE[chat_id]
    if not state.get("intake_session_id"):
        state["intake_session_id"] = __import__("uuid").uuid4().hex[:10]

    # Add files to the current intake session folder only.
    staging_dir = _DI_INTAKE_DIR / "staging" / str(state["intake_session_id"])
    staging_dir.mkdir(parents=True, exist_ok=True)
    local_path = staging_dir / file_name
    tg_file = await ctx.bot.get_file(doc.file_id)
    await tg_file.download_to_drive(str(local_path))

    # Delete previous Yes/No prompt so the chat keeps only the latest intake prompt.
    # Fallback to removing only the keyboard if Telegram refuses deletion.
    old_prompt_message_id = state.get("last_prompt_message_id")
    if old_prompt_message_id:
        try:
            await ctx.bot.delete_message(
                chat_id=chat_id,
                message_id=int(old_prompt_message_id),
            )
        except Exception as exc:
            log.debug("doc intake: could not delete old prompt message: %s", exc)
            try:
                await ctx.bot.edit_message_reply_markup(
                    chat_id=chat_id,
                    message_id=int(old_prompt_message_id),
                    reply_markup=None,
                )
            except Exception as exc2:
                log.debug("doc intake: could not clear old prompt keyboard: %s", exc2)

    if not state.get("intake_session_id"):
        state["intake_session_id"] = __import__("uuid").uuid4().hex[:10]
    state["prompt_version"] = int(state.get("prompt_version") or 0) + 1
    intake_session_id = state["intake_session_id"]
    prompt_version = state["prompt_version"]
    state["files"].append({
        "filename": file_name,
        "original_name": file_name,
        "local_path": str(local_path),
        "file_id": doc.file_id,
    })
    _PENDING_DOC_INTAKE[chat_id] = state

    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("Yes", callback_data=f"di_more:{intake_session_id}:{prompt_version}"),
        InlineKeyboardButton("No", callback_data=f"di_enough:{intake_session_id}:{prompt_version}"),
    ]])
    count = len(state["files"])
    sent_msg = await update.message.reply_text(
        f"Received: {file_name} ({count} document{'s' if count > 1 else ''} total)\nAny more documents?",
        reply_markup=kb,
    )
    state["last_prompt_message_id"] = sent_msg.message_id
    _PENDING_DOC_INTAKE[chat_id] = state


async def _di_receive_task(
    update: Update, ctx: ContextTypes.DEFAULT_TYPE, chat_id: int
) -> None:
    """Receive task description text and present mode selection."""
    if update.message is None:
        return
    task_text = (update.message.text or "").strip()
    if not task_text:
        return
    state = _PENDING_DOC_INTAKE[chat_id]
    state["task_description"] = task_text
    state["step"] = "await_mode"
    state["task_prompt_sent"] = True
    _PENDING_DOC_INTAKE[chat_id] = state
    await update.message.reply_text("Task received. Choose processing mode.")
    await _di_ask_mode(update, ctx, chat_id)


async def _di_ask_mode(
    update: Update, ctx: ContextTypes.DEFAULT_TYPE, chat_id: int
) -> None:
    """Present mode-selection buttons based on document count."""
    state = _PENDING_DOC_INTAKE.get(chat_id, {})
    count = len(state.get("files", []))

    if count == 1:
        kb = InlineKeyboardMarkup([[
            InlineKeyboardButton("Evaluate", callback_data="di_one_evaluate"),
            InlineKeyboardButton("No", callback_data="di_cancel"),
        ]])
        text = "1 document received. Evaluate against standards?"
    elif count == 2:
        kb = InlineKeyboardMarkup([[
            InlineKeyboardButton("Tuning", callback_data="di_two_tuning"),
            InlineKeyboardButton("Consolidate", callback_data="di_two_consolidate"),
        ]])
        text = "2 documents received. Create a tuning pair or compare and consolidate documents?"
    else:
        kb = InlineKeyboardMarkup([[
            InlineKeyboardButton("Register", callback_data="di_multi_register"),
            InlineKeyboardButton("Consolidate", callback_data="di_multi_consolidate"),
        ]])
        text = f"{count} documents received. Register documents or compare and consolidate into one document?"

    if update.message:
        await update.message.reply_text(text, reply_markup=kb)
    elif update.callback_query and update.callback_query.message:
        await update.callback_query.message.reply_text(text, reply_markup=kb)


async def _di_callback_handler(
    update: Update, ctx: ContextTypes.DEFAULT_TYPE
) -> None:
    """Handle all di_* and dr_* inline button callbacks."""
    query = update.callback_query
    if query is None:
        return
    await query.answer()
    chat_id = query.message.chat_id if query.message else None
    if chat_id is None:
        return
    # De-duplicate Telegram callback queries. This prevents double replies when
    # Telegram/client retries the same callback or one button event is processed twice.
    callback_id = str(getattr(query, "id", "") or "")
    if callback_id:
        if callback_id in _DOC_INTAKE_HANDLED_CALLBACKS:
            try:
                await query.answer("Already handled.", show_alert=False)
            except Exception:
                pass
            return
        _DOC_INTAKE_HANDLED_CALLBACKS.add(callback_id)
        if len(_DOC_INTAKE_HANDLED_CALLBACKS) > 2000:
            _DOC_INTAKE_HANDLED_CALLBACKS.clear()

    data = query.data or ""

    # ── dr_* — single-doc review ──────────────────────────────────────────────
    if data == "dr_standards_yes":
        state = _PENDING_DOCREVIEW.pop(chat_id, None)
        if not state or not state.get("files"):
            await query.edit_message_text("Session expired. Please upload the document again.")
            return
        task = state.get("task_description", "Evaluate against standards")
        case_dir = await _di_create_case_package(
            chat_id=chat_id,
            task_description=task,
            mode="evaluate_against_standards",
            files=state["files"],
            selected_pipeline="single_document_docreview",
            learning_contour=False,
        )
        reply_text = None
        try:
            from doc_agent_api import DocAgentClient  # type: ignore
            _api_url = os.environ.get("DOC_AGENT_API_URL", "http://doc-agent:8767")
            _client = DocAgentClient(_api_url)
            result = await asyncio.get_event_loop().run_in_executor(
                None, lambda: _client.review(str(case_dir))
            )
            reply_text = f"Document review submitted. Case: {case_dir.name}\nResult: {result}"
        except Exception:
            pass
        if reply_text is None:
            reply_text = (
                f"Document received for standards evaluation.\n"
                f"Case ID: {case_dir.name}\n"
                "The review pipeline will process it shortly."
            )
        await query.edit_message_text(reply_text)

    elif data == "dr_standards_no":
        _PENDING_DOCREVIEW.pop(chat_id, None)
        await query.edit_message_text("Document received. No standards check requested.")

    # ── di_more / di_enough ───────────────────────────────────────────────────
    elif data.startswith("di_more"):
        state = _PENDING_DOC_INTAKE.get(chat_id)
        if not state:
            await query.edit_message_text("Session expired. Please upload your documents again.")
            return
        parts = data.split(":")
        if len(parts) >= 3:
            if (
                str(state.get("intake_session_id")) != parts[1]
                or str(state.get("prompt_version")) != parts[2]
            ):
                await query.answer("Old prompt ignored. Use the latest document prompt.", show_alert=False)
                return
        elif state.get("intake_session_id"):
            await query.answer("Old prompt ignored. Use the latest document prompt.", show_alert=False)
            return
        state["step"] = "await_more"
        _PENDING_DOC_INTAKE[chat_id] = state
        await query.edit_message_text(
            f"{len(state['files'])} document(s) collected. Upload the next document."
        )
        return

    elif data.startswith("di_enough"):
        log.warning("DOC_INTAKE_TRACE di_enough chat_id=%s data=%s state=%s", chat_id, data, _PENDING_DOC_INTAKE.get(chat_id))
        state = _PENDING_DOC_INTAKE.get(chat_id)
        if not state:
            await query.edit_message_text("Session expired.")
            return
        parts = data.split(":")
        if len(parts) >= 3:
            if (
                str(state.get("intake_session_id")) != parts[1]
                or str(state.get("prompt_version")) != parts[2]
            ):
                await query.answer("Old prompt ignored. Use the latest document prompt.", show_alert=False)
                return
        elif state.get("intake_session_id"):
            await query.answer("Old prompt ignored. Use the latest document prompt.", show_alert=False)
            return

        # Avoid duplicate identical task prompts if the same callback/update is processed twice
        # or if the user presses "No" from more than one intake message.
        if state.get("task_description"):
            state["step"] = "await_mode"
            _PENDING_DOC_INTAKE[chat_id] = state
            await query.answer("Task already received.")
            await _di_ask_mode(update, ctx, chat_id)
            return

        if state.get("task_prompt_sent") and state.get("step") == "await_task":
            await query.answer("Task prompt already sent.")
            return

        state["step"] = "await_task"
        state["task_prompt_sent"] = True
        _PENDING_DOC_INTAKE[chat_id] = state
        count = len(state["files"])
        await query.edit_message_text(
            f"{count} document(s) ready. Describe the task for the uploaded documents.",
            reply_markup=None,
        )
        return

    # ── Mode callbacks ────────────────────────────────────────────────────────
    elif data == "di_one_evaluate":
        state = _PENDING_DOC_INTAKE.pop(chat_id, None)
        if not state:
            await query.edit_message_text("Session expired.")
            return
        task = state.get("task_description", "Evaluate against standards")
        case_dir = await _di_create_case_package(
            chat_id=chat_id,
            task_description=task,
            mode="evaluate_against_standards",
            files=state["files"],
            selected_pipeline="single_document_docreview",
            learning_contour=False,
        )
        reply_text = None
        try:
            from doc_agent_api import DocAgentClient  # type: ignore
            _api_url = os.environ.get("DOC_AGENT_API_URL", "http://doc-agent:8767")
            _client = DocAgentClient(_api_url)
            result = await asyncio.get_event_loop().run_in_executor(
                None, lambda: _client.review(str(case_dir))
            )
            reply_text = f"Review submitted. Case: {case_dir.name}\nResult: {result}"
        except Exception:
            pass
        if reply_text is None:
            reply_text = (
                f"Document queued for standards evaluation.\n"
                f"Case ID: {case_dir.name}"
            )
        await query.edit_message_text(reply_text)

    elif data == "di_two_tuning":
        state = _PENDING_DOC_INTAKE.pop(chat_id, None)
        if not state:
            await query.edit_message_text("Session expired.")
            return
        case_dir = await _di_create_case_package(
            chat_id=chat_id,
            task_description=state.get("task_description", "Create tuning pair"),
            mode="tuning_pair",
            files=state["files"],
            selected_pipeline="doctuning_pair_generation",
            learning_contour=True,
        )
        await query.edit_message_text(
            f"Tuning pair package created.\nCase ID: {case_dir.name}"
        )

    elif data in ("di_two_consolidate", "di_multi_consolidate"):
        state = _PENDING_DOC_INTAKE.pop(chat_id, None)
        if not state:
            await query.edit_message_text("Session expired.")
            return
        case_dir = await _di_create_case_package(
            chat_id=chat_id,
            task_description=state.get("task_description", "Compare and consolidate"),
            mode="compare_consolidate",
            files=state["files"],
            selected_pipeline="multi_document_consolidate",
            learning_contour=False,
        )
        await query.edit_message_text(
            f"Consolidation package created.\nCase ID: {case_dir.name}"
        )

    elif data == "di_multi_register":
        state = _PENDING_DOC_INTAKE.pop(chat_id, None)
        if not state:
            await query.edit_message_text("Session expired.")
            return
        case_dir = await _di_create_case_package(
            chat_id=chat_id,
            task_description=state.get("task_description", "Register documents"),
            mode="register",
            files=state["files"],
            selected_pipeline="multi_document_register",
            learning_contour=False,
        )
        await query.edit_message_text(
            f"Registration package created.\nCase ID: {case_dir.name}\n"
            "Documents will be reviewed before registration."
        )

    elif data == "di_cancel":
        _PENDING_DOC_INTAKE.pop(chat_id, None)
        _PENDING_DOCREVIEW.pop(chat_id, None)
        await query.edit_message_text("Document intake cancelled.")


async def cmd_doctuning_approve(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Approve the pending doctuning reference — saves SFT pair, DPO pair, and master document."""
    if not _chat_allowed(update) or update.message is None or update.effective_chat is None:
        return
    chat_id = update.effective_chat.id
    lang = AXI_FORCE_REPLY_LANG if AXI_FORCE_REPLY_LANG in ("en", "ru") else "ru"

    state = _PENDING_DOCFILL.get(chat_id)
    if not state or state.get("step") != "await_approval":
        await update.message.reply_text(
            "Нет активного ожидания одобрения. Запустите /doctuning для начала."
            if lang == "ru" else
            "No pending approval. Start with /doctuning.",
        )
        return

    # Guard: local fill still being generated (VRAM was queued)
    if chat_id in _DOCFILL_VRAM_QUEUE:
        await update.message.reply_text(
            "⏳ Локальная генерация ещё в очереди (ждём VRAM).\n"
            "Уведомлю, когда завершится — тогда и подтверждайте.\n"
            "Или /doctuning_cancel чтобы отменить (DPO-пара будет пропущена)."
            if lang == "ru" else
            "⏳ Local generation is still queued (waiting for VRAM).\n"
            "I'll notify you when it's done — approve then.\n"
            "Or /doctuning_cancel to cancel (DPO pair will be skipped).",
        )
        return

    # ── batch_training_pair: dedicated approval path ──────────────────────────
    if state.get("source") == "batch_training_pair":
        await _approve_doctuning_for_chat(ctx.bot, chat_id, state, lang)
        return
    # ── normal /doctuning approval path ───────────────────────────────────────

    blank_text    = str(state.get("blank_text", ""))
    filled_text   = str(state.get("filled_text", ""))
    filled_name   = str(state.get("filled_name", "reference"))
    doc_type      = str(state.get("doc_type", "procedure"))
    local_fill    = state.get("local_fill")   # may be None
    repair_report = str(state.get("repair_report", ""))

    await update.message.reply_text(
        "⏳ Сохраняю мастер-документ и тренировочные пары…"
        if lang == "ru" else
        "⏳ Saving master document and training pairs…"
    )

    loop = asyncio.get_event_loop()
    doc_title = f"{doc_type.title()} — {filled_name}"

    # Run saves sequentially to reduce CPU/GPU load (was parallel with asyncio.gather)
    # Step 1: Save SFT pair (legacy path) + 120B template_fill_pairs + approved_master_pairs
    try:
        sft_result = await loop.run_in_executor(
            None, _docfill_save_training_pair, blank_text, filled_text, doc_type, filled_name
        )
        sft_ok = True
    except Exception as e:
        sft_result = e
        sft_ok = False

    # Save to 120B training pool: template_fill (blank→approved_reference) and approved_master
    try:
        await loop.run_in_executor(
            None,
            lambda: _doc_creator_120b_save(
                pair_type="template_fill_pairs",
                blank_text=blank_text,
                filled_text=filled_text,
                doc_type=doc_type,
                source_name=filled_name,
            ),
        )
        await loop.run_in_executor(
            None,
            lambda: _doc_creator_120b_save(
                pair_type="approved_master_pairs",
                blank_text=blank_text,
                filled_text=filled_text,
                doc_type=doc_type,
                source_name=filled_name,
                extra_meta={"approved_by": "human_expert"},
            ),
        )
    except Exception as e:
        log.warning("doc_creator_120b save error: %s", e)

    # Step 2: Save DPO pair (if local fill available); record failed case otherwise
    dpo_saved = False
    if local_fill:
        try:
            await loop.run_in_executor(
                None, _docfill_save_dpo_pair, blank_text, filled_text, local_fill, doc_type
            )
            dpo_saved = True
        except Exception as e:
            log.warning("DPO save error: %s", e)
    else:
        # 120B generation failed — record for analysis
        try:
            await loop.run_in_executor(
                None,
                lambda: _doc_creator_120b_save_failed(
                    blank_text=blank_text,
                    doc_type=doc_type,
                    source_name=filled_name,
                    repair_report=repair_report,
                ),
            )
        except Exception as e:
            log.warning("doc_creator_120b failed-case save error: %s", e)

    # Step 3: Save master document to OmiAgent
    try:
        omi_result = await loop.run_in_executor(
            None, _docfill_save_to_omi, doc_title, filled_text, doc_type
        )
        omi_ok = True
        omi_id = omi_result if isinstance(omi_result, str) else None
    except Exception as e:
        omi_result = e
        omi_ok = False
        omi_id = None

    # Cleanup temporary files
    blank_path = state.get("blank_path")
    filled_path = state.get("filled_path")
    if blank_path:
        Path(str(blank_path)).unlink(missing_ok=True)
    if filled_path:
        Path(str(filled_path)).unlink(missing_ok=True)

    _PENDING_DOCFILL.pop(chat_id, None)

    # Build confirmation
    lines_ru = ["✅ Одобрено и сохранено:"]
    lines_en = ["✅ Approved and saved:"]

    if sft_ok:
        lines_ru.append("  SFT-пара: ops/ft/data/docfill_v1/train_docfill_v1.jsonl")
        lines_en.append("  SFT pair: ops/ft/data/docfill_v1/train_docfill_v1.jsonl")
    else:
        lines_ru.append(f"  ⚠️ SFT-пара: ошибка — {_tg_plain(sft_result)}")
        lines_en.append(f"  ⚠️ SFT pair: error — {_tg_plain(sft_result)}")

    if dpo_saved:
        lines_ru.append("  DPO-пара: data/training/standard_dpo_pairs.jsonl")
        lines_en.append("  DPO pair: data/training/standard_dpo_pairs.jsonl")
    elif local_fill:
        lines_ru.append(f"  ⚠️ DPO-пара: ошибка — {_tg_plain(results[2])}")
        lines_en.append(f"  ⚠️ DPO pair: error — {_tg_plain(results[2])}")
    else:
        report_note = f" ({_tg_plain(repair_report)})" if repair_report else ""
        lines_ru.append(f"  DPO-пара: Repairman исчерпал 5 циклов ремонта — пропущено{report_note}")
        lines_en.append(f"  DPO pair: Repairman exhausted 5 repair cycles — skipped{report_note}")

    if omi_ok:
        id_str = f" (id={_tg_plain(omi_id)})" if omi_id else ""
        lines_ru.append(f"  Мастер-документ → OmiAgent{id_str}: {_tg_plain(doc_title)}")
        lines_en.append(f"  Master document → OmiAgent{id_str}: {_tg_plain(doc_title)}")
    else:
        lines_ru.append(f"  ⚠️ Мастер-документ: ошибка OmiAgent — {_tg_plain(omi_result)}")
        lines_en.append(f"  ⚠️ Master document: OmiAgent error — {_tg_plain(omi_result)}")

    msg = "\n".join(lines_ru) if lang == "ru" else "\n".join(lines_en)
    await _safe_send_message(ctx.bot, chat_id, msg)




async def cmd_analyze(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not _chat_allowed(update):
        return
    text = _extract_analyze_text(update, ctx)
    if not text:
        _set_pending_analyze(update.effective_chat.id, "Analyze uploaded document and provide concise findings.")
        await update.message.reply_text(
            "Режим /analyze включён. Загрузите файлы группой.\n"
            "Я дождусь окончания загрузки и обработаю пакет целиком.\n"
            "Можно добавить цель: `/analyze <что проверить>`."
        )
        return
    _set_pending_analyze(update.effective_chat.id, text)
    await update.message.reply_text("Принято. Жду загрузку файлов для пакетного анализа.")


async def cmd_analyze_done(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not _chat_allowed(update):
        return
    if update.effective_chat is None:
        return
    chat_id = update.effective_chat.id
    state = _get_pending_analyze_state(chat_id)
    if not state:
        await update.message.reply_text("Нет активного режима /analyze.")
        return
    files = list(state.get("files") or [])
    if not files:
        await update.message.reply_text("Файлы для пакетного анализа ещё не загружены.")
        return
    override_prompt = " ".join(ctx.args).strip() if getattr(ctx, "args", None) else ""
    prompt = override_prompt or str(state.get("prompt", "")).strip()
    if _analyze_prompt_is_unclear(prompt):
        await update.message.reply_text(
            "Уточните цель: `/analyze_done <что сделать с файлами>`."
        )
        return
    lang = _reply_lang(update.message.text or prompt)
    _PENDING_ANALYZE.pop(chat_id, None)
    _pending_analyze_persist()
    await update.message.reply_text(
        "Запускаю пакетный анализ без ожидания таймера."
        if lang == "ru" else
        "Starting batch analysis immediately (without timer wait)."
    )
    await _process_analyze_batch(ctx.bot, chat_id, prompt, files, lang=lang)


async def cmd_analyze_end(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    """Alias for /analyze_done to avoid command mismatch in chats."""
    await cmd_analyze_done(update, ctx)


# ── Main message handler ───────────────────────────────────────────────────────

async def handle_chat(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> None:
    if not update.message or not update.message.text:
        return
    if not _claim_update_once(update):
        return
    _chat_id = update.effective_chat.id if update.effective_chat else None
    text = update.message.text.strip()
    if not _chat_allowed(update):
        return
    if await _is_omi_owned_intent(text):
        return
    # ── Intent clarification response ─────────────────────────────────────────
    if _chat_id and _chat_id in _pending_intent:
        pending = _pending_intent.pop(_chat_id)
        answer_text = text.lower()
        if any(w in answer_text for w in ("word", "docx", "document", "report", "анализ", "отчёт", "документ", "word-")):
            confirmed = "docx"
        elif any(w in answer_text for w in ("edit", "update", "table", "spreadsheet", "таблиц", "измен", "обнов")):
            confirmed = "edit"
        else:
            confirmed = "chat"
        _save_intent_example(pending["prompt"], pending["files"], confirmed)
        lang = pending.get("lang", "en")
        await _process_analyze_batch(ctx.bot, _chat_id, pending["prompt"], pending["files"], lang=lang)
        return
    # ── Doctuning batch_training_pair: standards approval ─────────────────────
    if _chat_id and _chat_id in _PENDING_DOCFILL:
        df_state = _PENDING_DOCFILL[_chat_id]
        raw_lower = text.lower()
        if df_state.get("step") == "await_standards_approval" and df_state.get("source") == "batch_training_pair":
            # Guard: no file was generated yet
            if any(w in raw_lower for w in ("file", "xlsx", "docx", "generate", "regenerate", "attachment", "send me")):
                await update.message.reply_text(
                    "The revised reference candidate is not generated yet. "
                    "Please continue the current step first."
                )
                return
            # Confirmation — trigger real continuation pipeline
            if any(w in raw_lower for w in ("yes", "continue", "proceed", "ok", "looks good", "go ahead")):
                await _run_doctuning_continue_for_chat(
                    ctx.bot, _chat_id, df_state,
                    lang=_reply_lang(update.message.text),
                )
                return
            # Cancel
            if any(w in raw_lower for w in ("cancel", "stop", "no")):
                _PENDING_DOCFILL.pop(_chat_id, None)
                await update.message.reply_text("Doctuning cancelled.")
                return
            # Upload standards intent
            if any(w in raw_lower for w in ("upload standards", "provide standards", "i will provide standards")):
                await update.message.reply_text(
                    "Please upload your standards document as a file."
                )
                return
            # Unknown text — hold, don't fall through to _process_request_text
            await update.message.reply_text(
                "Doctuning approval pending. Reply **yes** to generate the candidate or **cancel** to discard.",
                parse_mode="Markdown",
            )
            return
        # Other pending doctuning steps (not file-upload steps): guard hallucination claims
        if df_state.get("step") not in ("await_blank", "await_example"):
            if any(w in raw_lower for w in ("file", "xlsx", "docx", "generate", "regenerate", "attachment")):
                await update.message.reply_text(
                    "The revised reference candidate is not generated yet. "
                    "Please continue the current step first."
                )
                return
    # ── Doc intake: capture task description ──────────────────────────────────
    if _should_route_ai_label_removal_text(text, _chat_id):
        _PENDING_AI_LABEL_REMOVAL[_chat_id or -1] = {"prompt": text, "ts": time.time()}
        await update.message.reply_text(
            "Прикрепите документ (DOCX, PPTX, XLSX или PDF), и я удалю метки ИИ."
        )
        return
    if (
        _chat_id is not None
        and _chat_id in _PENDING_DOC_INTAKE
        and _PENDING_DOC_INTAKE[_chat_id].get("step") in ("await_task", "await_more", "await_mode")
        and (_PENDING_DOC_INTAKE[_chat_id].get("files") or [])
    ):
        await _di_receive_task(update, ctx, _chat_id)
        return
    # ──────────────────────────────────────────────────────────────────────────
    state = _get_pending_analyze_state(_chat_id) if _chat_id is not None else None
    if state and state.get("awaiting_clarify") and (state.get("files") or []):
        state["prompt"] = text
        state["awaiting_clarify"] = False
        files = list(state.get("files") or [])
        lang = _reply_lang(text)
        _PENDING_ANALYZE.pop(_chat_id, None)
        _pending_analyze_persist()
        await _process_analyze_batch(ctx.bot, _chat_id, text, files, lang=lang)
        return
    await _process_request_text(
        update,
        ctx,
        text,
        require_group_mention=AXI_GROUP_REQUIRE_MENTION,
    )


# ── Periodic stuck-task monitor (Axi как монитор качества) ───────────────────

async def _job_monitor_stuck_tasks(context: ContextTypes.DEFAULT_TYPE) -> None:
    """
    Периодически (TASK_REGISTRY_MONITOR_INTERVAL сек) проверяет зависшие задачи.
    При обнаружении — уведомляет владельца в Axi-чат.
    Задачи старше TASK_REGISTRY_AUTO_CLEANUP_MINUTES отправляются на расследование repairman,
    который анализирует причину зависания, регистрирует проблему и удаляет задачу.
    Роль Axi: monitors quality (docs/project-owner-answers.md §9).
    """
    if _tr_client is None:
        return
    try:
        tasks = _tr_client.find_stuck(older_than_minutes=TASK_REGISTRY_STUCK_MINUTES)
    except Exception as e:
        log.warning("monitor_stuck_tasks: %s", e)
        return

    if not tasks:
        log.debug("monitor: no stuck tasks")
        return

    # Investigate and cleanup very old stuck tasks via repairman
    auto_cleanup_threshold = TASK_REGISTRY_AUTO_CLEANUP_MINUTES
    very_old_tasks = [t for t in tasks if t.get("age_minutes", 0) >= auto_cleanup_threshold]

    if very_old_tasks:
        log.warning("monitor: found %d tasks stuck for >%.0f minutes - sending to repairman for investigation",
                    len(very_old_tasks), auto_cleanup_threshold)

        for t in very_old_tasks:
            task_id = t.get("task_id")
            assigned_to = t.get("assigned_to", "unassigned")
            age_min = t.get("age_minutes", 0)
            description = t.get("description", "")[:500]
            status = t.get("status", "unknown")

            # Send to repairman for root cause analysis
            investigation_prompt = (
                f"STUCK TASK INVESTIGATION\n\n"
                f"Task ID: {task_id}\n"
                f"Assigned to: {assigned_to}\n"
                f"Status: {status}\n"
                f"Age: {age_min:.0f} minutes\n"
                f"Description: {description}\n\n"
                f"Analyze why this task is stuck:\n"
                f"1. Is the assigned agent running?\n"
                f"2. Is there a missing dependency or configuration?\n"
                f"3. Is this a scheduled task with no executor?\n"
                f"4. Should this task type be disabled or reassigned?\n\n"
                f"Provide root_cause analysis and recommend fix."
            )

            try:
                # Call repairman for investigation
                import httpx
                _omnirouter_base = os.environ.get("AIMS_OMNIROUTER_URL", "http://127.0.0.1:20129").rstrip("/")
                gateway_url = f"{_omnirouter_base}/v1/chat/completions"
                token = os.environ.get("AIMS_CLAUDE_PROXY_TOKEN", "aims-local-repair-token")

                response = httpx.post(
                    gateway_url,
                    json={
                        "model": "claude-sonnet-4",
                        "messages": [{"role": "user", "content": investigation_prompt}],
                        "max_tokens": 500,
                    },
                    headers={"Authorization": f"Bearer {token}"},
                    timeout=30.0,
                )

                if response.status_code == 200:
                    result = response.json()
                    root_cause = result.get("choices", [{}])[0].get("message", {}).get("content", "unknown")
                    log.warning(
                        "REPAIRMAN INVESTIGATION: task_id=%s\n"
                        "Root cause: %s",
                        task_id, root_cause[:500]
                    )
                else:
                    root_cause = f"repairman_unavailable (status {response.status_code})"
                    log.warning("Repairman investigation failed for %s: %s", task_id, root_cause)

            except Exception as e:
                root_cause = f"repairman_error: {e}"
                log.warning("Repairman investigation error for %s: %s", task_id, e)

            # Register the problem and cleanup the task
            try:
                _tr_client.done(task_id, result_summary=f"stuck_investigated_cleanup: {root_cause[:200]}")
                log.info("monitor: investigated and cleaned up task %s (age: %.0f min)", task_id, age_min)
            except Exception as e:
                log.warning("monitor: failed to cleanup task %s: %s", task_id, e)

        # Remove cleaned tasks from notification list
        tasks = [t for t in tasks if t.get("age_minutes", 0) < auto_cleanup_threshold]

    if not tasks:
        log.debug("monitor: all stuck tasks investigated and cleaned")
        return

    # Suppress OCR handoff backlog from Axi alerts: these are handled by Omi batch
    # pipeline and create noisy false positives in Axi chat monitoring.
    filtered: list[dict] = []
    suppressed = 0
    for t in tasks:
        assigned = str(t.get("assigned_to", "")).lower()
        desc = str(t.get("description", "")).lower()
        if assigned == "omi_batch" and "ocr+analyze" in desc:
            suppressed += 1
            continue
        # Suppress repetitive "status of repeated anonymization" tracker tasks.
        if assigned == "omi" and ("статус повторной ан" in desc or "status of repeated anonym" in desc):
            suppressed += 1
            continue
        filtered.append(t)
    tasks = filtered
    if not tasks:
        if suppressed:
            log.info("monitor: suppressed %d omi_batch OCR backlog task(s)", suppressed)
        return

    # Deduplicate notifications: alert only on real task state changes.
    # Pure age growth must not produce repeated alerts.
    signature_parts: list[str] = []
    for t in sorted(tasks, key=lambda x: str(x.get("task_id", ""))):
        tid = str(t.get("task_id", "")).strip()
        if not tid:
            continue
        assigned = str(t.get("assigned_to", "")).strip().lower() or "-"
        status = str(t.get("status", "")).strip().lower() or "-"
        desc = str(t.get("description", "")).strip()[:120]
        signature_parts.append(f"{tid}|{assigned}|{status}|{desc}")
    signature = "||".join(signature_parts)
    last_sig = context.application.bot_data.get("_axi_stuck_last_signature")
    if signature == last_sig:
        log.info("monitor: skipped duplicate stuck-task alert (signature unchanged)")
        return
    context.application.bot_data["_axi_stuck_last_signature"] = signature

    lines = [f"⚠️ *Axi: зависшие задачи ({len(tasks)}):*"]
    for t in tasks[:5]:
        age_min = int(t['age_minutes'])
        lines.append(
            f"  • `{t['task_id']}` [{t['assigned_to'] or '—'}] "
            f"{age_min}м — {t['description'][:60]}"
        )
    if len(tasks) > 5:
        lines.append(f"  _…и ещё {len(tasks) - 5}_")
    lines.append(f"\n_Задачи старше {int(auto_cleanup_threshold)}м расследуются repairman и удаляются._")
    report_text = "\n".join(lines)

    bot = context.application.bot
    for owner_id in OWNER_CHATS:
        try:
            await bot.send_message(chat_id=owner_id, text=report_text, parse_mode="Markdown")
            log.info("monitor: notified owner %s about %d stuck tasks", owner_id, len(tasks))
        except Exception as e:
            log.warning("monitor: failed to notify owner %s: %s", owner_id, e)


# ── Bot setup ──────────────────────────────────────────────────────────────────



def _di_docs_agent_processed() -> Path:
    """Directory where docs-agent moves processed case folders."""
    return Path(os.environ.get("DOCS_AGENT_PROCESSED_DIR", "/data/docs_agent/processed"))


def _di_docs_agent_notify_interval_sec() -> float:
    try:
        return max(2.0, float(os.environ.get("DOCS_AGENT_RESULT_NOTIFY_INTERVAL_SEC", "10")))
    except Exception:
        return 10.0


async def _di_docs_agent_result_notifier(application: Application) -> None:
    """
    Poll docs-agent processed cases and send primary_output back to Telegram.

    Contract:
    processed/<case_id>/
      manifest.json  -> contains chat_id
      result.json    -> contains primary_output, status, selected_pipeline, note
      .telegram_notified marker prevents duplicate notifications
    """
    import json as _json
    import asyncio as _asyncio
    import logging as _logging
    from datetime import datetime as _datetime, timezone as _timezone

    _log = globals().get("log") or _logging.getLogger(__name__)
    bot = application.bot

    while True:
        try:
            root = _di_docs_agent_processed()
            root.mkdir(parents=True, exist_ok=True)

            for case_dir in sorted([x for x in root.iterdir() if x.is_dir()]):
                marker = case_dir / ".telegram_notified"
                if marker.exists():
                    continue

                manifest_path = case_dir / "manifest.json"
                result_path = case_dir / "result.json"
                if not manifest_path.exists() or not result_path.exists():
                    continue

                try:
                    manifest = _json.loads(manifest_path.read_text(encoding="utf-8"))
                    result = _json.loads(result_path.read_text(encoding="utf-8"))
                except Exception as exc:
                    _log.warning("docs-agent notifier: cannot read %s: %s", case_dir, exc)
                    continue

                chat_id = manifest.get("chat_id")
                if not chat_id:
                    _log.warning("docs-agent notifier: missing chat_id in %s", manifest_path)
                    continue

                case_id = result.get("case_id") or manifest.get("case_id") or case_dir.name
                pipeline = result.get("selected_pipeline") or manifest.get("selected_pipeline") or manifest.get("mode") or "docs-agent"
                status = result.get("status", "unknown")
                document_count = result.get("document_count") or manifest.get("document_count") or len(manifest.get("documents", []))
                note = result.get("note", "")
                primary_output = result.get("primary_output")

                summary = (
                    "✅ Пакет документов обработан docs-agent.\n\n"
                    f"Case ID: {case_id}\n"
                    f"Pipeline: {pipeline}\n"
                    f"Status: {status}\n"
                    f"Documents: {document_count}"
                )
                if note:
                    summary += f"\nNote: {note}"

                sent = False
                primary_path = None
                if primary_output:
                    primary_path = case_dir / str(primary_output)
                    if primary_path.exists() and primary_path.is_file():
                        try:
                            with primary_path.open("rb") as fh:
                                await bot.send_document(
                                    chat_id=chat_id,
                                    document=fh,
                                    filename=primary_path.name,
                                    caption=summary[:1024],
                                )
                            sent = True
                        except Exception as exc:
                            _log.warning("docs-agent notifier: send_document failed for %s: %s", primary_path, exc)

                if not sent:
                    if primary_output and primary_path is not None and not primary_path.exists():
                        summary += f"\n\n⚠️ Output file listed but not found: {primary_output}"
                    await bot.send_message(chat_id=chat_id, text=summary)
                    sent = True

                if sent:
                    marker.write_text(
                        _datetime.now(_timezone.utc).isoformat() + "\n",
                        encoding="utf-8",
                    )
                    _log.info("docs-agent notifier: Telegram notified for case %s", case_id)

        except _asyncio.CancelledError:
            raise
        except Exception as exc:
            _log.warning("docs-agent notifier loop error: %s", exc)

        await _asyncio.sleep(_di_docs_agent_notify_interval_sec())


async def _post_init(application: Application) -> None:
    try:
        application.create_task(_di_docs_agent_result_notifier(application))
    except Exception as e:
        log.warning("docs-agent notifier start failed: %s", e)

    commands = [
        BotCommand("start",          "Запуск / приветствие"),
        BotCommand("help",           "Справка по командам"),
        BotCommand("analyze",        "Ожидать и анализировать файл"),
        BotCommand("analyze_done",   "Запустить пакет /analyze сейчас"),
        BotCommand("analyze_end",    "Алиас завершения загрузки файлов"),
        BotCommand("aimarks_delete", "Удалить AI-метки из документа"),
        BotCommand("quality_report", "Отчёт качества задач"),
        BotCommand("stuck_tasks",    "Зависшие задачи"),
        BotCommand("close_task",     "Закрыть зависшую задачу"),
        BotCommand("doctuning",         "Создать тюнинг-пару из шаблона + эталона"),
        BotCommand("doctuning_approve", "Одобрить эталон — сохранить мастер-документ и пары"),
        BotCommand("doctuning_cancel",  "Отменить режим doctuning"),
        BotCommand("docgen_upgrade",    "Запуск DOCGEN upgrade batch (private only)"),
    ]
    try:
        await application.bot.set_my_commands(commands)
    except Exception as e:
        log.warning("set_my_commands failed: %s", e)

    try:
        _schedule_vram_warmup()
    except Exception as e:
        log.debug("post_init vram warm: %s", e)

    log.info(
        "Axi bot started. allowed=%s owner=%s "
        "aims_knowledge=%s local_llm_first=%s task_registry=%s",
        sorted(ALLOWED_CHATS),
        sorted(OWNER_CHATS),
        AXI_AIMS_KNOWLEDGE_MODE,
        _axi_llm_local_first_enabled(),
        _tr_client.base_url if _tr_client else "unavailable",
    )


def main() -> None:
    if not AXI_BOT_TOKEN:
        raise RuntimeError(
            "AXI_BOT_TOKEN (or AXI_TELEGRAM_TOKEN) is not set. "
            "Set it in .env or docker-compose environment."
        )

    _dialog_load()
    _pending_analyze_load()

    app = (
        Application.builder()
        .token(AXI_BOT_TOKEN)
        .post_init(_post_init)
        .build()
    )

    # Commands
    app.add_handler(CommandHandler("start",          cmd_start))
    app.add_handler(CommandHandler("help",           cmd_help))
    app.add_handler(CommandHandler("analyze",        cmd_analyze))
    app.add_handler(CommandHandler("analyze_done",   cmd_analyze_done))
    app.add_handler(CommandHandler("analyze_end",    cmd_analyze_end))
    app.add_handler(CommandHandler("aimarks_delete", cmd_aimarks_delete))
    app.add_handler(CommandHandler("quality_report", cmd_quality_report))
    app.add_handler(CommandHandler("stuck_tasks",    cmd_stuck_tasks))
    app.add_handler(CommandHandler("close_task",     cmd_close_task))
    app.add_handler(CommandHandler("cleanup_tasks",  cmd_cleanup_tasks))
    app.add_handler(CommandHandler("doctuning",         cmd_doctuning))
    app.add_handler(CommandHandler("doctuning_approve", cmd_doctuning_approve))
    app.add_handler(CommandHandler("doctuning_cancel",  cmd_doctuning_cancel))
    app.add_handler(CommandHandler("docgen_upgrade",    cmd_docgen_upgrade))

    # Inline button callbacks
    app.add_handler(CallbackQueryHandler(_cmd_start_callback, pattern="^axi_activate$"))
    app.add_handler(CallbackQueryHandler(_doctuning_callback_handler, pattern="^dt_"))
    app.add_handler(CallbackQueryHandler(_di_callback_handler, pattern="^di_"))
    app.add_handler(CallbackQueryHandler(_di_callback_handler, pattern="^dr_"))

    # AI-label removal clarification menu (5 options on ambiguous document intent)
    try:
        from ops.pipelines.ai_label_removal import telegram_intent_menu as _ai_menu
        from ops.pipelines.ai_label_removal.telegram_handler import (
            handle_intent_menu_callback as _ai_menu_cb,
        )

        async def _route_docgen(req, update, ctx):
            _files = [{"name": req.original_filename, "path": str(req.file_path)}]
            await _process_analyze_batch(
                ctx.bot, req.chat_id,
                (req.prompt or "") + " сделай word документ с анализом", _files, lang="ru",
            )

        async def _route_standard(req, update, ctx):
            _files = [{"name": req.original_filename, "path": str(req.file_path)}]
            await _process_analyze_batch(
                ctx.bot, req.chat_id,
                (req.prompt or "") + " проверь документ на соответствие стандартам и верни результат",
                _files, lang="ru",
            )

        async def _route_training(req, update, ctx):
            await ctx.bot.send_message(
                chat_id=req.chat_id,
                text="🎓 Для обучающей пары пришлите два документа (до и после) "
                     "одним сообщением с подписью «обучающая пара».",
            )

        async def _route_docsreg(req, update, ctx):
            await ctx.bot.send_message(
                chat_id=req.chat_id,
                text="🗂 Для регистрации пачки в DOCSREG загрузите комплект документов "
                     "с запросом регистрации в базу.",
            )

        _ai_menu.register_route("docgen", _route_docgen)
        _ai_menu.register_route("standard_revision", _route_standard)
        _ai_menu.register_route("training_pair", _route_training)
        _ai_menu.register_route("docsreg", _route_docsreg)
        app.add_handler(CallbackQueryHandler(_ai_menu_cb, pattern=r"^axi_intent:"))
        log.info("ai_label_removal: clarification menu + cleanup wired")
    except Exception as _ai_wire_err:
        log.warning("ai_label_removal wiring skipped: %s", _ai_wire_err)

    # Text messages
    app.add_handler(MessageHandler(filters.Document.ALL | filters.PHOTO, handle_file_upload))
    app.add_handler(MessageHandler(filters.Regex(r"^/DOCGEN_UPGRADE(?:\s|$)"), cmd_docgen_upgrade))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_chat))

    # Periodic stuck-task monitor (Axi = качественный монитор §9)
    if app.job_queue is not None:
        app.job_queue.run_repeating(
            _job_monitor_stuck_tasks,
            interval=TASK_REGISTRY_MONITOR_INTERVAL,
            first=60,
            name="monitor_stuck_tasks",
        )

        # VRAM queue monitor: resume docfill local-fill tasks when GPU memory is available
        app.job_queue.run_repeating(
            _job_docfill_vram_monitor,
            interval=60,
            first=90,
            name="docfill_vram_monitor",
        )
        app.job_queue.run_repeating(
            _job_flush_pending_analyze,
            interval=5,
            first=5,
            name="flush_pending_analyze",
        )
        app.job_queue.run_repeating(
            _job_axi_deliver_cross_handoffs,
            interval=2,
            first=3,
            name="axi_cross_handoff",
        )
    else:
        log.warning("job queue is unavailable; background monitors are disabled")

    log.info("Starting Axi bot polling…")
    open("/tmp/bot_alive", "w").close()
    app.run_polling(drop_pending_updates=True)


if __name__ == "__main__":
    main()
