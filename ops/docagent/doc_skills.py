"""
AIMS Document Skills — Python bridge for project agents.

Maps the AIMS document commands (doc-generate, doc-ingest, doc-search,
doc-edit, doc-analyze, doc-tables, doc-presentation, doc-knowledge-ops) to
importable Python functions that project agents can call through their scoped
runner.

Usage:
    from docagent.doc_skills import DocSkillRunner
    runner = DocSkillRunner()
    result = runner.invoke("doc-analyze", source="path/to/doc.pdf", analysis_type="iso")
    result = runner.invoke("doc-search", query="порядок действий при аварии", mode="hybrid")
    result = runner.invoke("doc-generate", topic="Инструкция по ТБ", doc_type="ИНС")
"""
from __future__ import annotations

import hashlib
import json
import logging
import os
import re
import sqlite3
import sys
import urllib.request
from datetime import datetime, timezone, timedelta
from pathlib import Path
from typing import Any

log = logging.getLogger("docagent.doc_skills")

# ── Paths ──────────────────────────────────────────────────────────────────────
_OPS = Path(__file__).resolve().parent.parent
_ROOT = _OPS.parent
_OUTBOX = _ROOT / "aims_workspace" / "outbox"
_OUTBOX.mkdir(parents=True, exist_ok=True)
_TRAINING_CANDIDATES = (
    _ROOT
    / "aims_workspace"
    / "axi_ft_log"
    / "document_training_candidates.jsonl"
)
_TRAINING_CANDIDATES.parent.mkdir(parents=True, exist_ok=True)
_DB_PATH = Path(os.environ.get("AIMS_REGISTRY_DB", str(_ROOT / "data" / "aims_registry.db")))

OLLAMA_URL = os.environ.get("OLLAMA_LOCAL_URL", "http://localhost:11434")
DRAFT_MODEL = os.environ.get("AIMS_DRAFT_MODEL", "axi_omi_sphere")

TZ_DUBAI = timezone(timedelta(hours=4))

# ISO 55001 clause keywords for compliance checks
_ISO_CLAUSES: dict[str, list[str]] = {
    "4.1 Context": ["контекст организации", "внешние факторы", "внутренние факторы", "заинтересованные стороны"],
    "5.1 Leadership": ["руководство", "политика", "роли", "ответственность"],
    "6.1 Risks": ["риски", "возможности", "управление рисками", "оценка рисков"],
    "6.2 Objectives": ["цели управления активами", "показатели", "kpi", "метрики"],
    "7.1 Resources": ["ресурсы", "компетентность", "осведомлённость"],
    "7.5 Documentation": ["документированная информация", "реестр", "актуализация"],
    "8.1 Operations": ["оперативное планирование", "жизненный цикл", "критичность активов"],
    "8.3 Change Mgmt": ["управление изменениями", "moc"],
    "9.1 Monitoring": ["мониторинг", "измерение", "анализ", "контроль"],
    "10.1 Improvement": ["постоянное улучшение", "несоответствия", "корректирующие действия"],
}

_DOC_TYPE_RULES: dict[str, list[str]] = {
    "ТО": ["технологическая операция", "порядок выполнения", "технологическая карта"],
    "ИНС": ["инструкция", "порядок действий", "обязан", "запрещается"],
    "РД": ["руководящий документ", "требования к", "нормы и правила"],
    "АКТ": ["акт", "составлен", "комиссия", "подписи"],
    "ОТЧ": ["отчёт", "результаты", "анализ", "выводы"],
    "ПЛАН": ["план", "мероприятия", "срок", "ответственный"],
    "СП": ["стандарт предприятия", "СП-", "утверждён приказом"],
}


# ── Internal helpers ──────────────────────────────────────────────────────────

def _now() -> str:
    return datetime.now(TZ_DUBAI).isoformat()


def _ollama_generate(prompt: str, temperature: float = 0.25, num_predict: int = 4000) -> str:
    body = json.dumps({
        "model": DRAFT_MODEL,
        "prompt": prompt,
        "stream": False,
        "options": {"temperature": temperature, "num_predict": num_predict},
    }).encode()
    req = urllib.request.Request(
        f"{OLLAMA_URL}/api/generate",
        data=body,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as r:
            return json.loads(r.read())["response"]
    except Exception as exc:
        log.warning("Ollama unavailable: %s", exc)
        return f"[Ollama unavailable: {exc}]"


def _load_text(source: str) -> tuple[str, str]:
    """Return (text, method). source = file path or str(doc_id)."""
    p = Path(source)
    if p.exists():
        if p.suffix == ".pdf":
            try:
                import pdfplumber  # type: ignore
                with pdfplumber.open(str(p)) as pdf:
                    text = "\n".join(pg.extract_text() or "" for pg in pdf.pages)
                return text, "pdfplumber"
            except ImportError:
                pass
        elif p.suffix == ".docx":
            try:
                from docx import Document  # type: ignore
                doc = Document(str(p))
                text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
                return text, "python-docx"
            except ImportError:
                pass
        elif p.suffix in (".txt", ".md"):
            return p.read_text(encoding="utf-8", errors="replace"), "plaintext"
    # Try registry by doc_id (UUID string or integer-looking string)
    try:
        if _DB_PATH.exists():
            conn = sqlite3.connect(str(_DB_PATH))
            row = conn.execute(
                "SELECT file_path, summary FROM documents WHERE id=?", (source,)
            ).fetchone()
            conn.close()
            if row:
                # Prefer loading from file_path; fall back to summary
                fp = row[0]
                if fp and Path(fp).exists():
                    return Path(fp).read_text(encoding="utf-8", errors="replace"), "registry-file"
                if row[1]:
                    return row[1], "registry-summary"
    except Exception:
        pass
    return "", "none"


def _db_register(title: str, doc_type: str, content: str, source_file: str) -> str | None:
    """Register document in aims_registry.db; returns doc_id (UUID string) or None."""
    try:
        from docs_pipeline.docs_registry_mutation_guard import legacy_direct_db_register_allowed
        allowed, reason = legacy_direct_db_register_allowed()
        if not allowed:
            log.warning("doc_skills: direct DB register blocked: %s", reason)
            return None
    except Exception as exc:
        log.warning("doc_skills: registry mutation guard unavailable: %s", exc)
    if not _DB_PATH.exists():
        return None
    try:
        import uuid as _uuid
        doc_id = str(_uuid.uuid4())
        content_hash = hashlib.sha256(content.encode()).hexdigest()
        summary = content[:500].replace("\n", " ").strip()
        conn = sqlite3.connect(str(_DB_PATH))
        conn.execute("""
            INSERT INTO documents
                (id, title, doc_type, status, quality_score, file_path,
                 content_hash, summary, created_at, updated_at, metadata)
            VALUES (?, ?, ?, 'draft', NULL, ?, ?, ?, ?, ?, '{}')
        """, (doc_id, title, doc_type, source_file, content_hash, summary, _now(), _now()))
        conn.commit()
        conn.close()
        return doc_id
    except Exception as exc:
        log.warning("DB register failed: %s", exc)
        return None


def _quarantine_training_candidate(
    operation: str,
    original: str,
    edited: str,
    score: float = 0.0,
) -> bool:
    try:
        pair = {
            "prompt": f"[{operation.upper()}] {original[:500]}",
            "completion": edited,
            "operation": operation,
            "local_score": score,
            "generated_at": _now(),
            "source": "doc-skills",
            "status": "CANDIDATE_NOT_TRAINING_DATA",
            "required_gates": [
                "complete_context_provenance",
                "deterministic_acceptance_checks",
                "claude_teacher_quality_pass",
                "recommendation_lineage_pass",
                "baseline_holdout_benchmark_win"
            ],
        }
        with open(_TRAINING_CANDIDATES, "a", encoding="utf-8") as f:
            f.write(json.dumps(pair, ensure_ascii=False) + "\n")
        return True
    except Exception as exc:
        log.warning("Training candidate quarantine failed: %s", exc)
        return False


def _nim_score(text: str) -> float | None:
    try:
        from docagent import nvidia_nim  # type: ignore
        s, _ = nvidia_nim.score(user_request="document quality", document_text=text)
        return s
    except Exception:
        return None


def _detect_doc_type(text: str) -> str:
    text_lower = text.lower()
    scores = {dt: sum(1 for kw in kws if kw in text_lower) for dt, kws in _DOC_TYPE_RULES.items()}
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "РД"


_STANDARD_FAMILY_KEYWORDS: dict[str, list[str]] = {
    "iec": ["iec", "61511", "61508", "61882", "60529", "60534", "60044", "60092", "60245", "60439"],
    "isa": ["isa", "s84", "isa-84", "isa 84", "instrumentation symbols", "loop diagram"],
    "iso": ["iso 14224", "iso 13702", "iso 14313", "iso 13705", "iso 15463", "iso 13503"],
    "dep": ["dep", "shell", "design engineering practice", "32.80", "30006", "40100"],
}

def _detect_standard_family(query: str) -> str | None:
    """Return doc_type filter string if query clearly targets one standard family."""
    q = query.lower()
    for family, keywords in _STANDARD_FAMILY_KEYWORDS.items():
        if any(kw in q for kw in keywords):
            return family
    return None


def _sqlite_table_columns(conn: sqlite3.Connection, table: str) -> set[str]:
    rows = conn.execute(f"PRAGMA table_info({table})").fetchall()
    return {str(row[1]) for row in rows if row and len(row) > 1}


def _sqlite_pick_column(columns: set[str], candidates: tuple[str, ...]) -> str | None:
    for candidate in candidates:
        if candidate in columns:
            return candidate
    return None


def _qdrant_search(query: str, top_k: int = 10,
                   filter_doc_type: str | None = None) -> list[dict]:
    try:
        from docagent.standards_rag import query as rag_query  # type: ignore
        result = rag_query(query, top_k=top_k, filter_doc_type=filter_doc_type)
        # result is RAGResult dataclass with .clauses list[Clause]
        clauses = getattr(result, "clauses", [])
        return [
            {
                "score": c.score,
                "payload": {
                    "doc_id":       c.standard_id,
                    "title":        c.clause_title or c.standard_id,
                    "doc_type":     c.doc_type,
                    "content":      c.text,
                    "standard_id":  c.standard_id,
                    "clause_ref":   c.clause_ref,
                    "clause_title": c.clause_title,
                    "source_file":  c.source_file,
                },
            }
            for c in clauses
        ]
    except Exception as exc:
        log.debug("Qdrant search unavailable: %s", exc)
        return []


def _sqlite_search(query: str, limit: int = 10, doc_type: str | None = None) -> list[dict]:
    if not _DB_PATH.exists():
        return []
    try:
        conn = sqlite3.connect(str(_DB_PATH))
        conn.row_factory = sqlite3.Row
        columns = _sqlite_table_columns(conn, "documents")
        if not columns:
            conn.close()
            return []
        doc_type_col = _sqlite_pick_column(columns, ("doc_type", "document_type", "type", "doc_kind"))
        if doc_type and not doc_type_col:
            log.warning(
                "SQLite search failed: no compatible document-type column in schema; columns=%s",
                sorted(columns),
            )
            conn.close()
            return []
        title_col = _sqlite_pick_column(
            columns,
            ("title", "file_name", "source_filename", "original_file_name", "canonical_file_name"),
        ) or "id"
        summary_col = _sqlite_pick_column(columns, ("summary", "anonymized_report_md", "master_doc_json", title_col)) or title_col
        updated_col = _sqlite_pick_column(
            columns,
            ("updated_at", "modified_at", "date_modified", "created_at", "date_added"),
        ) or "id"
        created_col = _sqlite_pick_column(
            columns,
            ("created_at", "date_added", "updated_at", "date_modified"),
        ) or updated_col
        q = f"%{query}%"
        filters = [f"({title_col} LIKE ? OR {summary_col} LIKE ?)"]
        params: list[Any] = [q, q]
        if doc_type_col and doc_type:
            filters.append(f"{doc_type_col} = ?")
            params.append(doc_type)
        sql = (
            f"SELECT id, {title_col} AS title, "
            f"{doc_type_col if doc_type_col else 'NULL'} AS doc_type, "
            f"{created_col} AS created_at, "
            f"substr(COALESCE({summary_col}, {title_col}, ''), 1, 200) as excerpt "
            f"FROM documents WHERE {' AND '.join(filters)} "
            f"ORDER BY {updated_col} DESC LIMIT ?"
        )
        params.append(limit)
        rows = conn.execute(sql, params).fetchall()
        results = []
        for r in rows:
            d = dict(r)
            if doc_type and doc_type_col and str(d.get("doc_type", "")).strip() != doc_type:
                continue
            d["score"] = 1.0
            results.append(d)
        conn.close()
        return results
    except Exception as exc:
        log.warning("SQLite search failed: %s", exc)
        return []


# ── Skill implementations ─────────────────────────────────────────────────────

def skill_generate(
    topic: str,
    doc_type: str = "procedure",
    context: str = "",
    language: str = "RU",
    **_: Any,
) -> dict:
    """Generate a new document using the AIMS pipeline (draft → score → revise)."""
    prompt = (
        f"Создай профессиональный технический документ типа «{doc_type}» на тему:\n{topic}\n\n"
        f"{'Контекст: ' + context if context else ''}\n"
        "Структура: цель, область применения, ответственность, порядок, приложения.\n"
        "Стиль: официально-деловой, ГОСТ Р 7.0.97-2016. Язык: Русский."
    )
    draft = _ollama_generate(prompt, temperature=0.3, num_predict=6000)
    nim_score_val = _nim_score(draft)
    final = draft
    if nim_score_val is not None and nim_score_val < 0.8:
        revise_prompt = (
            f"Улучши документ. Текущая оценка качества: {nim_score_val:.2f}/1.0.\n"
            "Добавь недостающие разделы, улучши формулировки.\n\n"
            f"ДОКУМЕНТ:\n{draft[:4000]}"
        )
        final = _ollama_generate(revise_prompt, temperature=0.2, num_predict=6000)
        nim_score_val = _nim_score(final) or nim_score_val

    ts = datetime.now(TZ_DUBAI).strftime("%Y%m%d_%H%M")
    out_path = _OUTBOX / f"generated_{ts}.md"
    out_path.write_text(final, encoding="utf-8")
    doc_id = _db_register(topic, doc_type, final, str(out_path))
    training_candidate_quarantined = False
    if nim_score_val and nim_score_val >= 0.8:
        training_candidate_quarantined = _quarantine_training_candidate(
            "generate",
            topic,
            final,
            nim_score_val,
        )

    return {
        "skill": "doc-generate",
        "topic": topic,
        "doc_type": doc_type,
        "words": len(final.split()),
        "quality_score": nim_score_val,
        "doc_id": doc_id,
        "output_path": str(out_path),
        "training_pair_saved": False,
        "training_candidate_quarantined": training_candidate_quarantined,
        "status": "generated",
    }


def skill_ingest(
    file_path: str,
    doc_type: str | None = None,
    force: bool = False,
    **_: Any,
) -> dict:
    """Ingest a document file into AIMS registry + Qdrant."""
    text, method = _load_text(file_path)
    if not text:
        return {"skill": "doc-ingest", "source_file": file_path, "status": "failed",
                "notes": "Could not extract text"}

    # Duplicate check (by content hash)
    content_hash = hashlib.sha256(text.encode()).hexdigest()
    if not force and _DB_PATH.exists():
        conn = sqlite3.connect(str(_DB_PATH))
        row = conn.execute(
            "SELECT id FROM documents WHERE content_hash=?", (content_hash,)
        ).fetchone()
        conn.close()
        if row:
            return {"skill": "doc-ingest", "source_file": file_path, "status": "duplicate",
                    "doc_id": row[0], "notes": "Already in registry; use force=True to re-ingest"}

    detected_type = doc_type or _detect_doc_type(text)
    title = Path(file_path).stem
    doc_id = _db_register(title, detected_type, text, file_path)

    # Qdrant indexing
    qdrant_ok = False
    try:
        from docagent.standards_rag import ingest_file  # type: ignore
        ingest_file(Path(file_path))
        qdrant_ok = True
    except Exception as exc:
        log.debug("Qdrant ingest skipped: %s", exc)

    return {
        "skill": "doc-ingest",
        "source_file": file_path,
        "doc_id": doc_id,
        "doc_type": detected_type,
        "chars_extracted": len(text),
        "ocr_method": method,
        "duplicate": False,
        "qdrant_indexed": qdrant_ok,
        "status": "registered",
    }


def skill_search(
    query: str,
    mode: str = "hybrid",
    limit: int = 10,
    doc_type: str | None = None,
    **kwargs: Any,
) -> dict:
    """Search AIMS registry: SQLite FTS5 + Qdrant semantic, hybrid ranking."""
    supported_modes = {"full-text", "semantic", "hybrid"}
    if mode not in supported_modes:
        return {
            "skill": "doc-search",
            "status": "failed",
            "query": query,
            "mode": mode,
            "results": [],
            "notes": (
                f"Unsupported search mode {mode!r}; "
                f"supported modes: {sorted(supported_modes)}"
            ),
        }
    if "max_results" in kwargs:
        limit = int(kwargs["max_results"])
    limit = max(1, min(int(limit), 100))

    sqlite_results: list[dict] = []
    qdrant_results: list[dict] = []

    if mode in ("full-text", "hybrid"):
        sqlite_results = _sqlite_search(query, limit=limit, doc_type=doc_type)

    if mode in ("semantic", "hybrid"):
        # Auto-detect standard family from query if doc_type not explicit
        qdrant_filter = doc_type or _detect_standard_family(query)
        qdrant_results = _qdrant_search(query, top_k=limit,
                                        filter_doc_type=qdrant_filter)

    # Merge & deduplicate
    seen: set[str] = set()
    merged: list[dict] = []
    for r in sqlite_results:
        did = str(r.get("id", ""))
        if did not in seen:
            seen.add(did)
            r["source"] = "fts5"
            r["score"] = 1.0
            merged.append(r)
    for r in qdrant_results:
        payload = r.get("payload", {})
        # Key = standard_id + clause_ref to deduplicate clauses
        did = f"{payload.get('standard_id','')}|{payload.get('clause_ref','')}"
        score = r.get("score", 0.0)
        if did not in seen and score >= 0.3:
            seen.add(did)
            merged.append({
                "id":          payload.get("doc_id", did),
                "title":       payload.get("clause_title") or payload.get("title", ""),
                "doc_type":    payload.get("doc_type", ""),
                "standard_id": payload.get("standard_id", ""),
                "clause_ref":  payload.get("clause_ref", ""),
                "score":       score,
                "source":      "qdrant",
                "excerpt":     str(payload.get("content", ""))[:200],
            })
    merged.sort(key=lambda x: x.get("score", 0), reverse=True)

    return {
        "skill": "doc-search",
        "status": "ok",
        "query": query,
        "mode": mode,
        "results": merged[:limit],
        "total_sqlite": len(sqlite_results),
        "total_qdrant": len(qdrant_results),
    }


def _knowledge_ops_audit(limit: int = 20) -> dict[str, Any]:
    """Inspect registry SSoT risks without acquiring a writable DB connection."""
    base = {
        "skill": "doc-knowledge-ops",
        "action": "audit",
        "source_skill": "third_party/ecc/skills/knowledge-ops/SKILL.md",
        "read_only": True,
        "database": str(_DB_PATH),
        "audited_at": _now(),
    }
    if not _DB_PATH.exists():
        return {
            **base,
            "status": "unavailable",
            "notes": "Registry database does not exist",
            "metrics": {},
            "risks": [],
        }

    limit = max(1, min(int(limit), 100))
    conn: sqlite3.Connection | None = None
    try:
        conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
        conn.row_factory = sqlite3.Row
        columns = _sqlite_table_columns(conn, "documents")
        if not columns:
            return {
                **base,
                "status": "unavailable",
                "notes": "documents table is missing or has no columns",
                "metrics": {},
                "risks": [],
            }

        total = int(conn.execute("SELECT COUNT(*) FROM documents").fetchone()[0])
        metrics: dict[str, int] = {"documents": total}
        risks: list[dict[str, Any]] = []

        if "content_hash" in columns:
            missing_hashes = int(
                conn.execute(
                    "SELECT COUNT(*) FROM documents "
                    "WHERE content_hash IS NULL OR trim(content_hash) = ''"
                ).fetchone()[0]
            )
            duplicate_hash_groups = int(
                conn.execute(
                    "SELECT COUNT(*) FROM ("
                    "SELECT content_hash FROM documents "
                    "WHERE content_hash IS NOT NULL AND trim(content_hash) <> '' "
                    "GROUP BY content_hash HAVING COUNT(*) > 1)"
                ).fetchone()[0]
            )
            duplicate_hash_rows = int(
                conn.execute(
                    "SELECT COALESCE(SUM(n - 1), 0) FROM ("
                    "SELECT COUNT(*) AS n FROM documents "
                    "WHERE content_hash IS NOT NULL AND trim(content_hash) <> '' "
                    "GROUP BY content_hash HAVING COUNT(*) > 1)"
                ).fetchone()[0]
            )
            metrics.update(
                {
                    "missing_content_hashes": missing_hashes,
                    "duplicate_content_hash_groups": duplicate_hash_groups,
                    "noncanonical_duplicate_rows": duplicate_hash_rows,
                }
            )
            if duplicate_hash_groups:
                samples = conn.execute(
                    "SELECT content_hash, COUNT(*) AS copies "
                    "FROM documents "
                    "WHERE content_hash IS NOT NULL AND trim(content_hash) <> '' "
                    "GROUP BY content_hash HAVING COUNT(*) > 1 "
                    "ORDER BY copies DESC, content_hash LIMIT ?",
                    (limit,),
                ).fetchall()
                risks.append(
                    {
                        "type": "duplicate_content",
                        "severity": "high",
                        "groups": duplicate_hash_groups,
                        "samples": [dict(row) for row in samples],
                    }
                )
            if missing_hashes:
                risks.append(
                    {
                        "type": "missing_content_hash",
                        "severity": "medium",
                        "rows": missing_hashes,
                    }
                )
        else:
            metrics["missing_content_hash_column"] = 1
            risks.append(
                {
                    "type": "missing_content_hash_column",
                    "severity": "high",
                    "rows": total,
                }
            )

        title_col = _sqlite_pick_column(
            columns,
            ("title", "canonical_file_name", "source_filename", "original_file_name", "file_name"),
        )
        if title_col:
            duplicate_title_groups = int(
                conn.execute(
                    f"SELECT COUNT(*) FROM ("
                    f"SELECT lower(trim({title_col})) AS normalized_title FROM documents "
                    f"WHERE {title_col} IS NOT NULL AND trim({title_col}) <> '' "
                    f"GROUP BY lower(trim({title_col})) HAVING COUNT(*) > 1)"
                ).fetchone()[0]
            )
            metrics["duplicate_normalized_title_groups"] = duplicate_title_groups
            if duplicate_title_groups:
                samples = conn.execute(
                    f"SELECT lower(trim({title_col})) AS normalized_title, COUNT(*) AS copies "
                    f"FROM documents WHERE {title_col} IS NOT NULL AND trim({title_col}) <> '' "
                    f"GROUP BY lower(trim({title_col})) HAVING COUNT(*) > 1 "
                    f"ORDER BY copies DESC, normalized_title LIMIT ?",
                    (limit,),
                ).fetchall()
                risks.append(
                    {
                        "type": "duplicate_normalized_title",
                        "severity": "medium",
                        "groups": duplicate_title_groups,
                        "samples": [dict(row) for row in samples],
                    }
                )

        return {
            **base,
            "status": "ok",
            "notes": "Read-only audit; no registry rows were changed",
            "metrics": metrics,
            "risks": risks,
            "ssot_consistent": not risks,
        }
    except sqlite3.Error as exc:
        return {
            **base,
            "status": "failed",
            "notes": f"Read-only SQLite audit failed: {exc}",
            "metrics": {},
            "risks": [],
        }
    finally:
        if conn is not None:
            conn.close()


def skill_knowledge_ops(
    action: str = "audit",
    query: str = "",
    mode: str = "hybrid",
    limit: int = 10,
    doc_type: str | None = None,
    **_: Any,
) -> dict:
    """
    Local knowledge-operations adapter for OMI.

    This adapts the ECC knowledge-ops search-before-store and canonical-home
    pattern without invoking an external document service. It intentionally
    exposes no write/delete action.
    """
    if action == "audit":
        return _knowledge_ops_audit(limit=limit)
    if action == "search":
        if not query.strip():
            return {
                "skill": "doc-knowledge-ops",
                "action": action,
                "status": "failed",
                "read_only": True,
                "notes": "query is required for action='search'",
            }
        result = skill_search(query=query, mode=mode, limit=limit, doc_type=doc_type)
        result.update(
            {
                "skill": "doc-knowledge-ops",
                "action": action,
                "source_skill": "third_party/ecc/skills/knowledge-ops/SKILL.md",
                "read_only": True,
                "search_before_store": True,
            }
        )
        return result
    return {
        "skill": "doc-knowledge-ops",
        "action": action,
        "status": "forbidden",
        "read_only": True,
        "notes": "Supported actions: audit, search; writes and deletes are not exposed",
    }


def skill_edit(
    source: str,
    operation: str = "rewrite",
    target_lang: str = "RU",
    **_: Any,
) -> dict:
    """Edit a document: rewrite, proofread, translate, shorten, expand, formalize."""
    text, method = _load_text(source)
    if not text:
        return {"skill": "doc-edit", "source": source, "status": "failed",
                "notes": "Could not load document"}

    prompts = {
        "rewrite": (
            "Перепиши следующий документ. Сохрани весь смысл и содержание, улучши стиль, "
            "структуру и читабельность. Приведи к официально-деловому стилю технической документации.\n\n"
            f"ИСХОДНЫЙ ТЕКСТ:\n{text[:5000]}\n\nПЕРЕПИСАННЫЙ ТЕКСТ:"
        ),
        "proofread": (
            "Выполни корректуру текста: исправь грамматику, орфографию, пунктуацию и "
            "стилистические ошибки. Верни ТОЛЬКО исправленный текст и список исправлений.\n\n"
            f"ТЕКСТ:\n{text[:5000]}"
        ),
        "translate": (
            f"Translate the following technical document from {'Russian to English' if target_lang == 'EN' else 'English to Russian'}. "
            "Maintain formal technical style, preserve all structure and numbering.\n\n"
            f"TEXT:\n{text[:5000]}\n\nTRANSLATION:"
        ),
        "shorten": (
            "Сократи текст до ключевых тезисов. Сохрани все важные факты, цифры и требования. "
            "Убери повторения. Длина результата — не более 30% от исходного.\n\n"
            f"ТЕКСТ:\n{text[:5000]}"
        ),
        "expand": (
            "Расширь текст, добавив детали, примеры, ссылки на нормативные документы и "
            "разъяснения для ключевых понятий. Сохрани исходную структуру.\n\n"
            f"ТЕКСТ:\n{text[:4000]}"
        ),
        "formalize": (
            "Приведи текст к официально-деловому стилю технической документации по ГОСТ Р 7.0.97-2016. "
            "Убери разговорные выражения, добавь стандартные формулировки.\n\n"
            f"ТЕКСТ:\n{text[:5000]}\n\nОФИЦИАЛЬНЫЙ ВАРИАНТ:"
        ),
        "structure": (
            "Добавь к документу чёткую структуру: заголовки разделов, нумерацию, подзаголовки. "
            "Сохрани весь исходный контент, только улучши структуру.\n\n"
            f"ТЕКСТ:\n{text[:5000]}"
        ),
    }
    temperature_map = {"rewrite": 0.25, "proofread": 0.1, "translate": 0.2,
                       "shorten": 0.2, "expand": 0.3, "formalize": 0.15, "structure": 0.2}

    prompt = prompts.get(operation, prompts["rewrite"])
    edited = _ollama_generate(prompt, temperature=temperature_map.get(operation, 0.2), num_predict=6000)

    ts = datetime.now(TZ_DUBAI).strftime("%Y%m%d_%H%M")
    out_path = _OUTBOX / f"edited_{operation}_{ts}.md"
    out_path.write_text(edited, encoding="utf-8")
    nim_score_val = _nim_score(edited)
    training_candidate_quarantined = False
    if nim_score_val and nim_score_val >= 0.8:
        training_candidate_quarantined = _quarantine_training_candidate(
            operation,
            text[:500],
            edited,
            nim_score_val,
        )

    orig_words = len(text.split())
    edit_words = len(edited.split())
    return {
        "skill": "doc-edit",
        "source": source,
        "operation": operation,
        "original_words": orig_words,
        "edited_words": edit_words,
        "change_pct": round(abs(edit_words - orig_words) / max(orig_words, 1) * 100, 1),
        "quality_score": nim_score_val,
        "output_path": str(out_path),
        "training_pair_saved": False,
        "training_candidate_quarantined": training_candidate_quarantined,
        "status": "edited",
    }


def skill_analyze(
    source: str,
    analysis_type: str = "full",
    language: str = "RU",
    **_: Any,
) -> dict:
    """Analyze document: structure, ISO 55001 compliance, quality, entities, topics."""
    text, method = _load_text(source)
    if not text:
        return {"skill": "doc-analyze", "source": source, "status": "failed",
                "notes": "Could not load document"}

    result: dict[str, Any] = {
        "skill": "doc-analyze",
        "source": source,
        "analysis_type": analysis_type,
        "analyzed_at": _now(),
    }

    # Structure
    if analysis_type in ("full", "structure"):
        words = text.split()
        sentences = re.split(r"[.!?]+", text)
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        lines = [l for l in text.split("\n") if l.strip()]
        headings = [l for l in lines if re.match(r"^#{1,4}\s+", l) or
                    (len(l) < 80 and l.isupper() and len(l.split()) > 1)]
        required = ["цель", "область применения", "ответственность", "порядок"]
        result["structure"] = {
            "words": len(words),
            "sentences": len([s for s in sentences if s.strip()]),
            "paragraphs": len(paragraphs),
            "headings": len(headings),
            "heading_list": headings[:8],
            "missing_sections": [s for s in required if s not in text.lower()],
        }

    # ISO 55001
    if analysis_type in ("full", "iso"):
        text_lower = text.lower()
        clause_scores: dict[str, Any] = {}
        total = 0.0
        for clause, keywords in _ISO_CLAUSES.items():
            hits = [kw for kw in keywords if kw in text_lower]
            cov = len(hits) / len(keywords)
            total += cov
            clause_scores[clause] = {
                "coverage": round(cov, 2),
                "found": hits,
                "missing": [kw for kw in keywords if kw not in text_lower],
                "status": "covered" if cov >= 0.5 else ("partial" if cov > 0 else "missing"),
            }
        result["iso_55001_score"] = round(total / len(_ISO_CLAUSES), 3)
        result["iso_clause_coverage"] = clause_scores

    # NIM quality score
    if analysis_type in ("full", "quality"):
        result["quality_score_nim"] = _nim_score(text)

    # Entity extraction
    if analysis_type in ("full", "entities"):
        dates = re.findall(r"\d{2}\.\d{2}\.\d{4}|\d{4}-\d{2}-\d{2}", text)
        standards = re.findall(r"ГОСТ\s*Р?\s*[\w.\-]+|ISO\s+[\d:\-]+|СП\s+[\d.\-]+|РД\s+[\d.\-]+",
                               text, re.IGNORECASE)
        units = re.findall(r"\d+(?:[.,]\d+)?\s*(?:кВт|МВт|кПа|МПа|°C|кг|т|м³|л/с|мм|см|м|км|%)", text)
        result["entities"] = {
            "dates": list(dict.fromkeys(dates))[:20],
            "standards": list(dict.fromkeys(standards))[:20],
            "numbers_with_units": list(dict.fromkeys(units))[:20],
        }

    # Topics via LLM
    if analysis_type in ("full", "topics"):
        topics_prompt = (
            "Извлеки из документа в формате JSON:\n"
            '{"main_topic": "...", "key_points": [...], "audience": "...", "requirements": [...]}\n\n'
            f"ДОКУМЕНТ:\n{text[:3000]}"
        )
        raw_topics = _ollama_generate(topics_prompt, temperature=0.15, num_predict=800)
        match = re.search(r"\{.*\}", raw_topics, re.DOTALL)
        if match:
            try:
                result["topics"] = json.loads(match.group())
            except Exception:
                result["topics"] = {"raw": raw_topics[:500]}
        else:
            result["topics"] = {"raw": raw_topics[:500]}

    # Recommendations
    recommendations = []
    if result.get("iso_55001_score", 1.0) < 0.5:
        recommendations.append("ISO 55001 coverage < 50% — add missing clause sections")
    if result.get("quality_score_nim") and result["quality_score_nim"] < 0.7:
        recommendations.append("Quality score < 0.7 — consider doc-edit rewrite")
    structure = result.get("structure", {})
    if structure.get("missing_sections"):
        recommendations.append(f"Missing sections: {structure['missing_sections']}")
    result["recommendations"] = recommendations

    ts = datetime.now(TZ_DUBAI).strftime("%Y%m%d_%H%M")
    report_path = _OUTBOX / f"analysis_{ts}.json"
    report_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    result["report_path"] = str(report_path)
    result["status"] = "analyzed"
    return result


def skill_tables(
    source: str,
    operation: str = "extract",
    **_: Any,
) -> dict:
    """Table operations: extract from PDF/DOCX, create from description, export to CSV."""
    if operation == "create":
        # source = description of the table to create
        prompt = (
            f"Создай таблицу в формате Markdown для: {source}\n"
            "Требования: чёткие заголовки, единицы измерения, примеры значений. Только таблица."
        )
        table_md = _ollama_generate(prompt, temperature=0.2, num_predict=2000)
        ts = datetime.now(TZ_DUBAI).strftime("%Y%m%d_%H%M")
        out_path = _OUTBOX / f"table_{ts}.md"
        out_path.write_text(table_md, encoding="utf-8")
        return {"skill": "doc-tables", "operation": "create", "source": source,
                "output_path": str(out_path), "status": "created"}

    p = Path(source)
    if not p.exists():
        return {"skill": "doc-tables", "operation": operation, "source": source,
                "status": "failed", "notes": "File not found"}

    tables_found = 0
    export_path: str | None = None

    if operation == "extract":
        if p.suffix == ".pdf":
            try:
                import pdfplumber  # type: ignore
                with pdfplumber.open(str(p)) as pdf:
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        tables_found += len(tables)
            except ImportError:
                pass
        elif p.suffix == ".docx":
            try:
                from docx import Document  # type: ignore
                doc = Document(str(p))
                tables_found = len(doc.tables)
            except ImportError:
                pass

    elif operation in ("export", "csv"):
        try:
            import pandas as pd  # type: ignore
            df = pd.read_csv(str(p)) if p.suffix == ".csv" else pd.read_excel(str(p))
            ts = datetime.now(TZ_DUBAI).strftime("%Y%m%d_%H%M")
            out_path = _OUTBOX / f"table_export_{ts}.csv"
            df.to_csv(str(out_path), index=False, encoding="utf-8-sig")
            export_path = str(out_path)
            tables_found = 1
        except ImportError:
            pass

    return {
        "skill": "doc-tables",
        "operation": operation,
        "source": source,
        "tables_found": tables_found,
        "exported_to": export_path,
        "status": "done",
    }


def skill_presentation(
    topic: str,
    slides: int = 10,
    format: str = "md",
    style: str = "technical",
    language: str = "RU",
    **_: Any,
) -> dict:
    """Generate a presentation: Marp MD (default), PPTX, or Reveal.js HTML."""
    prompt = (
        f"Создай структуру профессиональной технической презентации на тему: {topic}\n\n"
        "Формат ответа — строго JSON:\n"
        '{"title": "...", "slides": [{"slide_num": 1, "title": "...", '
        '"bullets": ["...", "..."], "speaker_notes": "..."}]}\n'
        f"Количество слайдов: {slides}. Стиль: {style}. Язык: {'русский' if language == 'RU' else 'английский'}."
    )
    raw = _ollama_generate(prompt, temperature=0.3, num_predict=3000)
    match = re.search(r"\{.*\}", raw, re.DOTALL)
    slides_data: dict = {}
    if match:
        try:
            slides_data = json.loads(match.group())
        except Exception:
            slides_data = {"title": topic, "slides": []}

    ts = datetime.now(TZ_DUBAI).strftime("%Y%m%d_%H%M")
    slides_count = len(slides_data.get("slides", []))

    if format == "md":
        lines = ["---", "marp: true", "theme: default", "paginate: true", "---", "",
                 f"# {slides_data.get('title', topic)}", "", "---"]
        for slide in slides_data.get("slides", [])[1:]:
            lines.append(f"## {slide.get('title', '')}")
            lines.append("")
            for bullet in slide.get("bullets", []):
                lines.append(f"- {bullet}")
            if slide.get("speaker_notes"):
                lines += ["", f"<!-- {slide['speaker_notes']} -->"]
            lines += ["", "---", ""]
        out_path = _OUTBOX / f"slides_{ts}.md"
        out_path.write_text("\n".join(lines), encoding="utf-8")

    elif format == "pptx":
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches, Pt  # type: ignore
            from pptx.dml.color import RGBColor  # type: ignore
            BLUE = RGBColor(0x1E, 0x40, 0x7A)
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            for slide_data in slides_data.get("slides", []):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")
                    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for bullet in slide_data.get("bullets", []):
                        p = tf.add_paragraph()
                        p.text = str(bullet)
                        p.font.size = Pt(18)
                slide.notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")
            out_path = _OUTBOX / f"presentation_{ts}.pptx"
            prs.save(str(out_path))
        except ImportError:
            out_path = _OUTBOX / f"slides_{ts}.md"
            out_path.write_text(raw, encoding="utf-8")
            format = "md"

    else:  # html
        def make_slide(s: dict) -> str:
            bullets = "".join(f"<li>{b}</li>" for b in s.get("bullets", []))
            return f"<section><h2>{s.get('title','')}</h2><ul>{bullets}</ul></section>"
        sections = "\n".join(make_slide(s) for s in slides_data.get("slides", []))
        html = (
            f'<!DOCTYPE html><html><head><meta charset="utf-8"><title>{slides_data.get("title", topic)}</title></head>'
            f'<body><div class="reveal"><div class="slides">{sections}</div></div></body></html>'
        )
        out_path = _OUTBOX / f"slides_{ts}.html"
        out_path.write_text(html, encoding="utf-8")

    return {
        "skill": "doc-presentation",
        "topic": topic,
        "slides_count": slides_count,
        "format": format,
        "output_path": str(out_path),
        "status": "created",
    }


# ── Unified runner ─────────────────────────────────────────────────────────────

_SKILL_MAP = {
    "doc-generate": skill_generate,
    "doc-ingest": skill_ingest,
    "doc-search": skill_search,
    "doc-edit": skill_edit,
    "doc-analyze": skill_analyze,
    "doc-tables": skill_tables,
    "doc-presentation": skill_presentation,
    "doc-knowledge-ops": skill_knowledge_ops,
}


class DocSkillRunner:
    """
    Unified entry point for all AIMS document skills.

    Agents import this class and call:
        runner = DocSkillRunner()
        result = runner.invoke("doc-analyze", source="path/to/doc.pdf")
        result = runner.invoke("doc-search", query="управление рисками")
        result = runner.invoke("doc-generate", topic="Инструкция по ТБ", doc_type="ИНС")

    Each skill returns a dict matching its output contract.
    All errors are caught; status="failed" is returned on failure.
    """

    def __init__(self) -> None:
        self._skills = _SKILL_MAP

    def available_skills(self) -> list[str]:
        return list(self._skills.keys())

    def invoke(self, skill: str, **kwargs: Any) -> dict:
        fn = self._skills.get(skill)
        if fn is None:
            return {"skill": skill, "status": "failed",
                    "notes": f"Unknown skill '{skill}'. Available: {list(self._skills)}"}
        try:
            log.info("DocSkillRunner.invoke: %s kwargs=%s", skill, list(kwargs.keys()))
            return fn(**kwargs)
        except Exception as exc:
            log.exception("DocSkillRunner.invoke error: skill=%s", skill)
            return {"skill": skill, "status": "failed", "notes": str(exc)}

    # Convenience aliases
    def generate(self, topic: str, **kw: Any) -> dict:
        return self.invoke("doc-generate", topic=topic, **kw)

    def ingest(self, file_path: str, **kw: Any) -> dict:
        return self.invoke("doc-ingest", file_path=file_path, **kw)

    def search(self, query: str, **kw: Any) -> dict:
        return self.invoke("doc-search", query=query, **kw)

    def edit(self, source: str, operation: str = "rewrite", **kw: Any) -> dict:
        return self.invoke("doc-edit", source=source, operation=operation, **kw)

    def analyze(self, source: str, analysis_type: str = "full", **kw: Any) -> dict:
        return self.invoke("doc-analyze", source=source, analysis_type=analysis_type, **kw)

    def tables(self, source: str, operation: str = "extract", **kw: Any) -> dict:
        return self.invoke("doc-tables", source=source, operation=operation, **kw)

    def presentation(self, topic: str, **kw: Any) -> dict:
        return self.invoke("doc-presentation", topic=topic, **kw)

    def knowledge_ops(self, action: str = "audit", **kw: Any) -> dict:
        return self.invoke("doc-knowledge-ops", action=action, **kw)
