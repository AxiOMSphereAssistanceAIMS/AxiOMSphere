"""
axi_template_validator.py — валидация, анализ и очистка шаблонов документов.

Место: ops/axi_template_validator.py

Этапы обработки входящего шаблона:
  1. validate_security()   — magic bytes, VBA макросы, размер, целостность
  2. analyse_structure()   — стили, заголовки, колонтитулы, мастер, макеты
  3. build_report()        — текстовый отчёт для пользователя
  4. clean_template()      — очистить текст, оставить 1 символ на стиль
  5. apply_text_changes()  — применить правки из текстового описания (через Gemini)
"""

from __future__ import annotations

import io
import logging
import os
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger("axi_bot")

MAX_TEMPLATE_BYTES = int(os.environ.get("AXI_TEMPLATE_MAX_BYTES", str(50 * 1024 * 1024)))  # 50 MB
# Символ-заглушка оставляемый в каждом стиле после очистки
STYLE_PLACEHOLDER_CHAR = "·"


# ─────────────────────────────────────────────────────────────────────────────
# Результаты валидации и анализа
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class SecurityResult:
    ok: bool
    reason: str = ""
    has_macros: bool = False
    file_size_bytes: int = 0


@dataclass
class StructureResult:
    fmt: str = ""                        # docx / pptx / xlsx
    fonts: list[str] = field(default_factory=list)
    colors: list[str] = field(default_factory=list)
    heading_levels: list[str] = field(default_factory=list)   # ["Heading 1", "Heading 2"]
    named_styles: list[str] = field(default_factory=list)     # все именованные стили
    has_header: bool = False
    has_footer: bool = False
    has_logo: bool = False
    slide_layouts: list[str] = field(default_factory=list)    # PPTX макеты
    sheet_names: list[str] = field(default_factory=list)      # XLSX листы
    warnings: list[str] = field(default_factory=list)
    quality_ok: bool = True              # False если структура слишком бедная


# ─────────────────────────────────────────────────────────────────────────────
# 1. Валидация безопасности
# ─────────────────────────────────────────────────────────────────────────────

def validate_security(path: Path) -> SecurityResult:
    """
    Проверить файл перед признанием шаблоном:
    - существует и не пустой
    - не превышает MAX_TEMPLATE_BYTES
    - magic bytes соответствуют Office формату
    - не содержит VBA макросов
    - zip-архив не повреждён (для docx/pptx/xlsx)
    """
    if not path.is_file():
        return SecurityResult(ok=False, reason="File not found.")

    size = path.stat().st_size
    if size == 0:
        return SecurityResult(ok=False, reason="File is empty.")
    if size > MAX_TEMPLATE_BYTES:
        mb = size // (1024 * 1024)
        return SecurityResult(
            ok=False,
            reason=f"File too large: {mb} MB (max {MAX_TEMPLATE_BYTES // (1024*1024)} MB).",
            file_size_bytes=size,
        )

    ext = path.suffix.lower()
    if ext not in (".docx", ".pptx", ".xlsx"):
        return SecurityResult(
            ok=False,
            reason=f"Unsupported format: {ext}. Supported: .docx, .pptx, .xlsx.",
            file_size_bytes=size,
        )

    # Magic bytes: все три формата — ZIP архив (PK\x03\x04)
    try:
        with open(path, "rb") as f:
            magic = f.read(4)
        if magic[:2] != b"PK":
            return SecurityResult(
                ok=False,
                reason="File is not a valid Office document (invalid magic bytes).",
                file_size_bytes=size,
            )
    except OSError as e:
        return SecurityResult(ok=False, reason=f"Cannot read file: {e}", file_size_bytes=size)

    # Целостность ZIP + проверка VBA
    has_macros = False
    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            # VBA макросы присутствуют если есть vbaProject.bin
            vba_files = [n for n in names if "vbaProject" in n or n.endswith(".bin") and "vba" in n.lower()]
            if vba_files:
                has_macros = True
                return SecurityResult(
                    ok=False,
                    reason=f"Template contains VBA macros ({', '.join(vba_files)}). "
                           "Remove macros before using as template.",
                    has_macros=True,
                    file_size_bytes=size,
                )
            # Проверить что основные части читаемы
            for name in names[:5]:
                zf.read(name)
    except zipfile.BadZipFile:
        return SecurityResult(
            ok=False,
            reason="File is corrupted (cannot open as ZIP/Office document).",
            file_size_bytes=size,
        )
    except Exception as e:
        return SecurityResult(ok=False, reason=f"File integrity check failed: {e}", file_size_bytes=size)

    return SecurityResult(ok=True, has_macros=has_macros, file_size_bytes=size)


# ─────────────────────────────────────────────────────────────────────────────
# 2. Анализ структуры
# ─────────────────────────────────────────────────────────────────────────────

def analyse_structure(path: Path) -> StructureResult:
    """Извлечь структурные атрибуты шаблона для отчёта и последующей генерации."""
    ext = path.suffix.lower()
    if ext == ".docx":
        return _analyse_docx(path)
    elif ext == ".pptx":
        return _analyse_pptx(path)
    elif ext == ".xlsx":
        return _analyse_xlsx(path)
    return StructureResult(warnings=["Unknown format"])


def _analyse_docx(path: Path) -> StructureResult:
    from docx import Document
    r = StructureResult(fmt="docx")
    try:
        doc = Document(str(path))

        # Стили
        style_names = []
        heading_levels = []
        for style in doc.styles:
            try:
                name = style.name
                style_names.append(name)
                if name.startswith("Heading"):
                    heading_levels.append(name)
                # Шрифты
                if style.font and style.font.name and style.font.name not in r.fonts:
                    r.fonts.append(style.font.name)
            except Exception:
                pass
        r.named_styles = style_names[:30]
        r.heading_levels = sorted(set(heading_levels))

        # Колонтитулы
        for section in doc.sections:
            try:
                if any(p.text.strip() for p in section.header.paragraphs):
                    r.has_header = True
                if any(p.text.strip() for p in section.footer.paragraphs):
                    r.has_footer = True
            except Exception:
                pass

        # Логотип (изображения)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                r.has_logo = True
                break

        # Цвета из темы
        try:
            theme_part = doc.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
            )
            xml = theme_part.blob.decode("utf-8", errors="ignore")
            colors = re.findall(r'val="([0-9A-Fa-f]{6})"', xml)
            r.colors = list(dict.fromkeys(colors))[:8]
        except Exception:
            pass

        # Оценка качества структуры
        non_normal = [s for s in r.named_styles if s not in ("Normal", "Default Paragraph Style")]
        if not r.heading_levels:
            r.warnings.append("No heading styles found — document structure may be flat.")
            r.quality_ok = False
        if len(non_normal) < 3:
            r.warnings.append("Very few named styles — consider adding Heading 1-3, Body Text, Caption.")
            r.quality_ok = False

    except Exception as e:
        r.warnings.append(f"DOCX analysis error: {e}")
        r.quality_ok = False
    return r


def _analyse_pptx(path: Path) -> StructureResult:
    from pptx import Presentation
    r = StructureResult(fmt="pptx")
    try:
        prs = Presentation(str(path))

        if not prs.slide_masters:
            r.warnings.append("No slide master found — presentation has no theme.")
            r.quality_ok = False
        else:
            master = prs.slide_masters[0]
            # Макеты слайдов
            r.slide_layouts = [lay.name for lay in master.slide_layouts]

            # Шрифты из темы мастера
            try:
                tf = master.theme.font_scheme
                for attr in ("major_font", "minor_font"):
                    font = getattr(tf, attr, None)
                    if font and hasattr(font, "latin") and font.latin.typeface:
                        name = font.latin.typeface
                        if name and name not in r.fonts:
                            r.fonts.append(name)
            except Exception:
                pass

            # Логотип в мастере
            for shape in master.shapes:
                if shape.shape_type == 13:  # PICTURE
                    r.has_logo = True
                    break

            # Цвета из темы
            try:
                from pptx.util import Pt
                clr_map = master.theme.element
                colors = re.findall(r'val="([0-9A-Fa-f]{6})"', clr_map.xml)
                r.colors = list(dict.fromkeys(colors))[:8]
            except Exception:
                pass

        if len(r.slide_layouts) < 2:
            r.warnings.append("Few slide layouts found — consider adding Title, Content, Section Header.")
            r.quality_ok = False

    except Exception as e:
        r.warnings.append(f"PPTX analysis error: {e}")
        r.quality_ok = False
    return r


def _analyse_xlsx(path: Path) -> StructureResult:
    import openpyxl
    r = StructureResult(fmt="xlsx")
    try:
        wb = openpyxl.load_workbook(str(path), read_only=False, data_only=True)
        r.sheet_names = wb.sheetnames

        # Именованные стили
        try:
            for ns in wb.named_styles:
                r.named_styles.append(ns.name)
                if ns.font and ns.font.name and ns.font.name not in r.fonts:
                    r.fonts.append(ns.font.name)
        except Exception:
            pass

        if not r.named_styles:
            r.warnings.append("No named styles found — consider adding Header, Data, Total styles.")
            r.quality_ok = False

        wb.close()
    except Exception as e:
        r.warnings.append(f"XLSX analysis error: {e}")
        r.quality_ok = False
    return r


# ─────────────────────────────────────────────────────────────────────────────
# 3. Отчёт для пользователя
# ─────────────────────────────────────────────────────────────────────────────

def build_report(path: Path, sec: SecurityResult, struct: StructureResult) -> str:
    """
    Сформировать текстовый отчёт об анализе шаблона.
    Отправляется пользователю перед кнопками выбора.
    """
    lines = [f"Template analysis: {path.name}"]
    lines.append(f"Format: {struct.fmt.upper()}  |  Size: {sec.file_size_bytes // 1024} KB")
    lines.append("")

    # Шрифты
    if struct.fonts:
        lines.append(f"Fonts: {', '.join(struct.fonts[:5])}")
    else:
        lines.append("Fonts: not detected")

    # Цвета
    if struct.colors:
        lines.append(f"Theme colors: {len(struct.colors)} ({', '.join('#' + c for c in struct.colors[:4])})")

    # Структура по формату
    if struct.fmt == "docx":
        if struct.heading_levels:
            lines.append(f"Headings: {', '.join(struct.heading_levels)}")
        else:
            lines.append("Headings: none found")
        lines.append(f"Named styles: {len(struct.named_styles)}")
        hf_parts = []
        if struct.has_header:
            hf_parts.append("header")
        if struct.has_footer:
            hf_parts.append("footer")
        lines.append(f"Header/footer: {', '.join(hf_parts) if hf_parts else 'none'}")
        lines.append(f"Logo/image: {'yes' if struct.has_logo else 'no'}")

    elif struct.fmt == "pptx":
        lines.append(f"Slide layouts: {len(struct.slide_layouts)}")
        if struct.slide_layouts:
            lines.append(f"  {', '.join(struct.slide_layouts[:6])}")
        lines.append(f"Logo in master: {'yes' if struct.has_logo else 'no'}")

    elif struct.fmt == "xlsx":
        lines.append(f"Sheets: {', '.join(struct.sheet_names[:6])}")
        lines.append(f"Named styles: {len(struct.named_styles)}")

    # Предупреждения о структуре
    if struct.warnings:
        lines.append("")
        lines.append("Warnings:")
        for w in struct.warnings:
            lines.append(f"  · {w}")

    # Итог
    lines.append("")
    if struct.quality_ok:
        lines.append("Structure: OK — ready for review.")
    else:
        lines.append("Structure: needs improvement — see warnings above.")

    return "\n".join(lines)


# ─────────────────────────────────────────────────────────────────────────────
# 4. Очистка шаблона (оставить 1 символ на стиль)
# ─────────────────────────────────────────────────────────────────────────────

def clean_template(src_path: Path, dest_path: Path) -> bool:
    """
    Очистить текстовое содержимое документа, оставив по одному
    символу-заглушке (STYLE_PLACEHOLDER_CHAR) в каждом именованном стиле.

    src_path  — исходный файл (оригинал, не изменяется)
    dest_path — куда сохранить очищенный шаблон

    Returns True при успехе.
    """
    ext = src_path.suffix.lower()
    try:
        if ext == ".docx":
            return _clean_docx(src_path, dest_path)
        elif ext == ".pptx":
            return _clean_pptx(src_path, dest_path)
        elif ext == ".xlsx":
            return _clean_xlsx(src_path, dest_path)
        else:
            logger.error("clean_template: unsupported format %s", ext)
            return False
    except Exception as e:
        logger.error("clean_template error src=%s: %s", src_path, e)
        return False


def _clean_docx(src: Path, dest: Path) -> bool:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(str(src))

    # Собрать все стили которые используются в документе
    used_styles: set[str] = set()
    for para in doc.paragraphs:
        if para.style and para.style.name:
            used_styles.add(para.style.name)

    # Пройти по параграфам: первый параграф каждого стиля — оставить placeholder,
    # остальные — удалить текст (но сохранить параграф чтобы стиль не потерялся)
    seen_styles: set[str] = set()
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else "Normal"
        # Очистить все runs в параграфе
        for run in para.runs:
            run.text = ""
        # Если стиль встречаем первый раз — поставить placeholder в первый run
        if style_name not in seen_styles:
            seen_styles.add(style_name)
            if para.runs:
                para.runs[0].text = STYLE_PLACEHOLDER_CHAR
            else:
                # Нет runs — добавить новый
                run = para.add_run(STYLE_PLACEHOLDER_CHAR)

    # Очистить таблицы
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        run.text = ""

    # Колонтитулы — оставить как есть (могут содержать логотип, номера страниц)

    doc.save(str(dest))
    return True


def _clean_pptx(src: Path, dest: Path) -> bool:
    from pptx import Presentation
    from pptx.util import Pt
    import shutil

    # Копируем оригинал — мастер и макеты НЕ трогаем (в них стили)
    shutil.copy2(str(src), str(dest))
    prs = Presentation(str(dest))

    # Удалить все слайды с контентом, оставить только мастер
    # (шаблон PPTX = только мастер + макеты, без слайдов)
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rId = slide_id.get("r:id")
        if rId:
            prs.part.drop_rel(rId)
            slide_id.getparent().remove(slide_id)

    # Добавить по одному слайду на каждый макет с placeholder символом
    seen_layouts: set[str] = set()
    for layout in prs.slide_layouts:
        layout_name = layout.name
        if layout_name in seen_layouts:
            continue
        seen_layouts.add(layout_name)
        try:
            slide = prs.slides.add_slide(layout)
            # Поставить placeholder в первый текстовый фрейм
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        if p.runs:
                            p.runs[0].text = STYLE_PLACEHOLDER_CHAR
                        else:
                            p.add_run().text = STYLE_PLACEHOLDER_CHAR
                    break
        except Exception as e:
            logger.debug("clean_pptx layout=%s err=%s", layout_name, e)

    prs.save(str(dest))
    return True


def _clean_xlsx(src: Path, dest: Path) -> bool:
    import openpyxl
    import shutil

    shutil.copy2(str(src), str(dest))
    wb = openpyxl.load_workbook(str(dest))

    seen_styles: set[str] = set()
    for ws in wb.worksheets:
        first_row_written = False
        for row in ws.iter_rows():
            for cell in row:
                # Запомнить стиль первой непустой ячейки листа
                style_key = f"{ws.title}_{cell.style_id if hasattr(cell, 'style_id') else ''}"
                if not first_row_written and cell.value is not None:
                    first_row_written = True
                    cell.value = STYLE_PLACEHOLDER_CHAR
                else:
                    cell.value = None

    wb.save(str(dest))
    wb.close()
    return True


# ─────────────────────────────────────────────────────────────────────────────
# 5. Применение текстовых правок через LLM
# ─────────────────────────────────────────────────────────────────────────────

async def apply_text_changes(
    path: Path,
    changes_description: str,
    llm_reply_fn,  # callable: async (text: str) -> str
) -> tuple[Path | None, str]:
    """
    Применить текстовые правки к шаблону используя LLM.

    Подход:
    1. Извлечь текущую структуру шаблона
    2. Попросить LLM описать что нужно изменить в JSON
    3. Применить изменения программно
    4. Вернуть путь к изменённому файлу

    Returns: (new_path, status_message)
    new_path = None при ошибке
    """
    ext = path.suffix.lower()
    struct = analyse_structure(path)

    # Строим промпт для LLM
    struct_summary = build_report(path, SecurityResult(ok=True, file_size_bytes=path.stat().st_size), struct)
    prompt = f"""You are helping modify a document template structure.

Current template:
{struct_summary}

Requested changes:
{changes_description}

Respond with a JSON object describing ONLY structural changes (no content):
{{
  "heading_levels_to_add": ["Heading 1", "Heading 2"],
  "styles_to_add": ["Body Text", "Caption"],
  "add_header": true or false,
  "add_footer": true or false,
  "notes": "Brief description of what was changed"
}}

Respond with valid JSON only, no markdown.
"""
    try:
        raw = await llm_reply_fn(prompt)
        # Очистить JSON от markdown если есть
        raw = re.sub(r"```json\s*|\s*```", "", raw).strip()
        import json
        changes = json.loads(raw)
    except Exception as e:
        logger.warning("apply_text_changes llm parse error: %s", e)
        return None, f"Could not parse change instructions from AI: {e}"

    # Применить изменения к DOCX
    if ext == ".docx":
        new_path = path.parent / (path.stem + "_revised" + path.suffix)
        ok, msg = _apply_docx_changes(path, new_path, changes)
        return (new_path if ok else None), msg
    else:
        # Для PPTX/XLSX — пока возвращаем оригинал с описанием изменений
        # (структурные правки сложнее, реализуются в следующей итерации)
        notes = changes.get("notes", "Changes noted but not yet auto-applied for this format.")
        return path, f"Notes from AI: {notes}\n\nFor PPTX/XLSX please upload a revised file."


def _apply_docx_changes(src: Path, dest: Path, changes: dict) -> tuple[bool, str]:
    """Применить структурные изменения к DOCX файлу."""
    from docx import Document
    import shutil

    shutil.copy2(str(src), str(dest))
    doc = Document(str(dest))
    applied: list[str] = []

    # Добавить стили заголовков если нужны
    existing_styles = {s.name for s in doc.styles}

    for heading in changes.get("heading_levels_to_add") or []:
        if heading not in existing_styles:
            try:
                # Добавить параграф с нужным стилем — это создаст стиль если его нет
                para = doc.add_paragraph(STYLE_PLACEHOLDER_CHAR)
                para.style = doc.styles.add_style(heading, 1)  # WD_STYLE_TYPE.PARAGRAPH = 1
                applied.append(f"Added style: {heading}")
            except Exception as e:
                logger.debug("add heading style %s: %s", heading, e)

    for style_name in changes.get("styles_to_add") or []:
        if style_name not in existing_styles:
            try:
                para = doc.add_paragraph(STYLE_PLACEHOLDER_CHAR)
                doc.styles.add_style(style_name, 1)
                para.style = doc.styles[style_name]
                applied.append(f"Added style: {style_name}")
            except Exception as e:
                logger.debug("add style %s: %s", style_name, e)

    doc.save(str(dest))
    notes = changes.get("notes", "")
    msg = f"Applied: {', '.join(applied) if applied else 'no structural changes'}"
    if notes:
        msg += f"\n{notes}"
    return True, msg
