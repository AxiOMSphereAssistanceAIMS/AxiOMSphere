"""Test helpers: build OOXML documents with AI metadata for ai_label_removal tests.

Not a test module (no test_ prefix) — imported by the ai_label_removal tests.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

_CUSTOM_CT = (
    b'<Override PartName="/docProps/custom.xml" '
    b'ContentType="application/vnd.openxmlformats-officedocument.custom-properties+xml"/>'
)


def inject_custom_props(path: Path, props: dict[str, str]) -> None:
    """Add a docProps/custom.xml part with the given name->value properties."""
    entries = []
    for i, (name, value) in enumerate(props.items(), start=2):
        entries.append(
            f'<property fmtid="{{D5CDD505-2E9C-101B-9397-08002B2CF9AE}}" '
            f'pid="{i}" name="{name}"><vt:lpwstr>{value}</vt:lpwstr></property>'
        )
    custom = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties" '
        'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        + "".join(entries)
        + "</Properties>"
    ).encode("utf-8")

    data = Path(path).read_bytes()
    bio = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(data)) as zin, zipfile.ZipFile(bio, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            d = zin.read(item.filename)
            if item.filename == "[Content_Types].xml" and b"custom-properties" not in d:
                d = d.replace(b"</Types>", _CUSTOM_CT + b"</Types>")
            zout.writestr(item, d)
        zout.writestr("docProps/custom.xml", custom)
    Path(path).write_bytes(bio.getvalue())


def build_docx(path: Path, *, ai: bool = True) -> Path:
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Report", level=1)
    doc.add_paragraph("This is the visible body text that must be preserved.")
    doc.add_paragraph("Second paragraph with specific value 42 and details.")
    if ai:
        doc.core_properties.author = "ChatGPT"
        doc.core_properties.keywords = "quarterly; AIGenerated; finance"
        doc.core_properties.last_modified_by = "OpenAI GPT-4"
    else:
        doc.core_properties.author = "Jane Engineer"
    doc.save(str(path))
    if ai:
        inject_custom_props(
            path,
            {
                "AIGenerated": "true",
                "Department": "Finance",
                "GenerationTool": "Claude",
                "GrammarlyDocumentId": "g-1234567890",
            },
        )
    return path


def build_pptx(path: Path, *, ai: bool = True) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(2):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.text = f"Slide {i + 1}: visible content preserved."
    if ai:
        prs.core_properties.author = "ChatGPT"
        prs.core_properties.last_modified_by = "Anthropic Claude"
    prs.save(str(path))
    if ai:
        inject_custom_props(path, {"AIGenerated": "true", "Owner": "Marketing"})
    return path


def build_xlsx(path: Path, *, ai: bool = True) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "Item"
    ws["B1"] = "Amount"
    ws["A2"] = "Widget"
    ws["B2"] = 10
    ws["A3"] = "Gadget"
    ws["B3"] = 20
    ws["B4"] = "=SUM(B2:B3)"
    ws2 = wb.create_sheet("Notes")
    ws2["A1"] = "Keep this sheet"
    if ai:
        wb.properties.creator = "ChatGPT"
        wb.properties.lastModifiedBy = "OpenAI"
    wb.save(str(path))
    if ai:
        inject_custom_props(path, {"AIGenerated": "true", "Region": "EMEA"})
    return path


def sign_package(path: Path) -> Path:
    """Return a copy of an OOXML package with a fake digital-signature part."""
    data = Path(path).read_bytes()
    bio = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(data)) as zin, zipfile.ZipFile(bio, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            zout.writestr(item, zin.read(item.filename))
        zout.writestr("_xmlsignatures/sig1.xml", b"<Signature/>")
    signed = path.with_name(path.stem + "_signed" + path.suffix)
    signed.write_bytes(bio.getvalue())
    return signed


def make_encrypted_office(path: Path) -> Path:
    """Create a fake encrypted Office file (OLE compound magic) with office extension."""
    path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512)
    return path
